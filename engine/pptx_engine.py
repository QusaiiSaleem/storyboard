"""
PPTX Template Engine — LectureBuilder
======================================

Builds Interactive Lecture presentations using a HYBRID approach:
1. Opens the real template PPTX as the base (gets backgrounds, headers, footers, logos)
2. Deletes the template's example slides
3. Adds new slides using the template's own layouts (Layout 0 = title, Layout 1 = content)
4. Places content using add_textbox() / add_picture() — NEVER touches placeholder shapes

Template: "قالب المحاضرة التفاعلية- عربي.pptx"

Design specs from the template:
- Slide dimensions: 12192000 x 6858000 EMU (16:9 widescreen)
- Primary font: Tajawal (ExtraBold, Medium, Regular)
- All text is Arabic RTL
- Color palette: #2D588C (primary blue), #333333 (body), #262626 (subtitle)

CRITICAL RULES:
- NEVER modify placeholder shapes (causes overlapping text)
- ALWAYS use slide.shapes.add_textbox() for text
- ALWAYS use slide.shapes.add_picture() for images

Usage:
    from engine.pptx_engine import LectureBuilder

    builder = LectureBuilder(
        project_code="DSAI",
        unit_number=1,
        unit_name="المهارات الرقمية",
        institution="جامعة نجران - كلية علوم الحاسب ونظم المعلومات"
    )
    builder.add_title_slide(title="المحاضرة الأولى", subtitle="مقدمة في علوم الحاسوب")
    builder.add_objectives_slide(objectives=["تعريف ماهي التقنية", "التعرف على الفوائد"])
    builder.add_content_slide(title="المقدمة", bullets=["نقطة أولى", "نقطة ثانية"])
    builder.save("output/DSAI/U01/DSAI_U01_Interactive_Lecture.pptx")
"""

import os
from datetime import datetime
from typing import Optional

from PIL import Image as PILImage  # For reading image dimensions (aspect ratio)

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

# Shared RTL helpers — critical workarounds for Arabic text in python-pptx.
# These functions handle XML-level operations that python-pptx doesn't
# expose natively (paragraph RTL direction, complex script font assignment).
from engine.rtl_helpers import (
    pptx_set_paragraph_rtl,
    pptx_set_paragraph_ltr,
    pptx_set_run_font_arabic,
)

# Image generation — auto-generates images when agent provides a prompt
# but no pre-existing image file. Falls back gracefully if API key missing.
from engine.image_gen import generate_storyboard_image


# ---------------------------------------------------------------------------
# Design Constants — extracted from the real template
# ---------------------------------------------------------------------------

# Template file path — the engine opens this as the base presentation
TEMPLATE_PATH = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "templates",
    "pptx_assets",
)

# Slide dimensions (EMU) — standard 16:9 widescreen
SLIDE_WIDTH = 12192000
SLIDE_HEIGHT = 6858000

# Color palette — hex values from the template analysis
PRIMARY_BLUE = RGBColor(0x2D, 0x58, 0x8C)     # #2D588C — headings, slide numbers
ACCENT1_BLUE = RGBColor(0x15, 0x60, 0x82)     # #156082 — theme accent1, button fills
BODY_TEXT = RGBColor(0x33, 0x33, 0x33)          # #333333 — body and section titles
SUBTITLE_TEXT = RGBColor(0x26, 0x26, 0x26)      # #262626 — lecture subtitle
LINK_BLUE = RGBColor(0x2E, 0x6C, 0xEC)         # #2E6CEC — summary link text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)              # #FFFFFF — button text, light bg
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)           # Dark background for quiz/card slides
BUTTON_BORDER = RGBColor(0x08, 0x28, 0x36)      # #082836 — button border
NOTES_YELLOW = RGBColor(0xFF, 0xFF, 0x00)       # #FFFF00 — notes callout

# Accent colors for cards and interactive elements
TEAL = RGBColor(0x00, 0x96, 0x88)              # Teal for accent bars
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)      # Green for correct answers
ACCENT_RED = RGBColor(0xF4, 0x43, 0x36)         # Red for wrong answers
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)      # Orange for cards
LIGHT_BLUE_BG = RGBColor(0xE3, 0xF2, 0xFD)     # Light blue background

# Professional design colors — added for consultancy-quality slides
CONTENT_CARD_BG = RGBColor(0xF5, 0xF7, 0xFA)    # Light gray for content cards
CONTENT_CARD_BORDER = RGBColor(0xE0, 0xE5, 0xEC) # Subtle border for content cards
CARD_LIGHT_BG = RGBColor(0xFA, 0xFB, 0xFC)       # Very light card body background
OPTION_ALT_BG = RGBColor(0xF0, 0xF4, 0xF8)       # Alternating option background
DIVIDER_BG = RGBColor(0x2D, 0x58, 0x8C)          # Section divider background
BULLET_MARKER_COLOR = RGBColor(0x2D, 0x58, 0x8C) # Blue bullet circles
SHADOW_COLOR = RGBColor(0xE0, 0xE0, 0xE0)        # Lighter shadow color

# Extended palette — tints and shades for visual depth
PRIMARY_BLUE_LIGHT = RGBColor(0x4A, 0x7A, 0xAE)    # Lighter tint of primary
PRIMARY_BLUE_DARK = RGBColor(0x1E, 0x3D, 0x63)      # Darker shade for depth
TEAL_LIGHT = RGBColor(0x4D, 0xBF, 0xB3)             # Lighter teal
WARM_GRAY = RGBColor(0x6B, 0x6B, 0x6B)              # Secondary text color
ACCENT_DEFINITION = RGBColor(0x00, 0x96, 0x88)       # Teal for definitions
ACCENT_EXAMPLE = RGBColor(0xFF, 0x98, 0x00)          # Orange for examples

# Header bar color — a slightly darker blue for the top banner
HEADER_BAR_BLUE = RGBColor(0x2D, 0x58, 0x8C)

# PNG asset file names (extracted from the template)
ASSET_BANNER_NARROW = "banner_narrow.png"   # Section banner (objectives, content slides)
ASSET_BANNER_WIDE = "banner_wide.png"       # Activity/summary banner (wider)
ASSET_OBJECTIVE_ROW = "objective_row.png"   # Gradient bar for objective rows
ASSET_TARGET_ICON = "target_icon.png"       # Target/circle icon at end of objective rows
ASSET_PLAY_ICON = "play_icon.png"           # Play button triangle icon (title slide)
ASSET_HAND_CURSOR = "hand_cursor.png"       # Hand cursor icon (title slide)

# Additional template assets — decorative elements for visual variety
ASSET_CORNER_TR = "corner_tr.png"          # Decorative corner, top-right
ASSET_CORNER_BL = "corner_bl.png"          # Decorative corner, bottom-left
ASSET_TEXT_BUBBLE = "text_bubble.png"       # Speech/thought bubble shape
ASSET_IMAGE_FRAME = "image_frame.png"      # Frame for image placeholders
ASSET_CONTENT_BG = "content_bg.png"        # Subtle texture background

# Font names — Tajawal is the primary font from the template.
# We set it on cs_font (Complex Script) for Arabic rendering,
# and also on latin_font and ea_font for consistency.
# STORYLINE REQUIREMENT: Developer must install Tajawal fonts
# before importing. Download: https://fonts.google.com/specimen/Tajawal
FONT_EXTRABOLD = "Tajawal ExtraBold"
FONT_MEDIUM = "Tajawal Medium"
FONT_REGULAR = "Tajawal"
FONT_FALLBACK = "Sakkal Majalla"  # Fallback if Tajawal not installed

# Positions (EMU) — extracted from the template's exact coordinates
# These ensure shapes land in the same spots as the original template.

# Lecture title bar — appears on slides 2-8 at the top center
TITLE_BAR_LEFT = 3405034
TITLE_BAR_TOP = 114300
TITLE_BAR_WIDTH = 5181600
TITLE_BAR_HEIGHT = 369332

# Section banner — centered below title bar
BANNER_LEFT = 4790969
BANNER_TOP = 898751
BANNER_WIDTH = 2610062
BANNER_HEIGHT = 695099

# Wider banner — used on activity and summary slides
WIDE_BANNER_LEFT = 3884635
WIDE_BANNER_TOP = 860142
WIDE_BANNER_WIDTH = 4422731
WIDE_BANNER_HEIGHT = 695099

# Wide banner text position (from template)
WIDE_BANNER_TEXT_LEFT = 3818244
WIDE_BANNER_TEXT_TOP = 977750
WIDE_BANNER_TEXT_WIDTH = 4555512
WIDE_BANNER_TEXT_HEIGHT = 400110

# Narrow banner text position (from template)
NARROW_BANNER_TEXT_LEFT = 4947367
NARROW_BANNER_TEXT_TOP = 1035917
NARROW_BANNER_TEXT_WIDTH = 2297266
NARROW_BANNER_TEXT_HEIGHT = 369332

# Content area — main body region for text
CONTENT_LEFT = 900000       # ~2.5cm from left
CONTENT_TOP = 2000000       # ~5.5cm from top
CONTENT_WIDTH = 10300000    # ~28.6cm wide
CONTENT_HEIGHT = 4000000    # ~11.1cm tall

# Text margins (EMU) — the template uses 0.25cm left/right, 0.13cm top/bottom
TEXT_MARGIN_LR = Cm(0.25)
TEXT_MARGIN_TB = Cm(0.13)


class LectureBuilder:
    """
    Builds an Interactive Lecture PPTX from scratch.

    This class creates presentations that match the visual design of
    "قالب المحاضرة التفاعلية- عربي.pptx" without editing the template
    file directly.

    The approach:
    1. Create a blank presentation with correct dimensions
    2. Add slides using blank layouts
    3. Build each slide's visual structure using shapes, textboxes, colors
    4. All text is set to RTL Arabic

    Attributes:
        project_code: Short code like "DSAI" or "NJR01"
        unit_number: Integer unit number (1, 2, 3, ...)
        unit_name: Arabic name of the unit
        institution: Arabic name of the university/institution
        designer: Name of the instructional designer
        prs: The python-pptx Presentation object
        slide_count: Running count of slides (for page numbering)

    Example:
        >>> builder = LectureBuilder("DSAI", 1, "المهارات الرقمية", "جامعة نجران")
        >>> builder.add_title_slide("المحاضرة الأولى", "مقدمة في علوم الحاسوب")
        >>> builder.save("output.pptx")
    """

    def __init__(
        self,
        project_code: str,
        unit_number: int,
        unit_name: str,
        institution: str,
        designer: str = "",
        template_path: str = None,
    ):
        """
        Initialize a new LectureBuilder.

        Opens the template PPTX as the base presentation (gets backgrounds,
        headers, footers, logos from the layouts), then deletes all example
        slides so we start with a clean slate.

        Args:
            project_code: Short project identifier (e.g., "DSAI")
            unit_number: Unit number (e.g., 1)
            unit_name: Arabic name of the unit
            institution: Arabic name of the institution
            designer: Name of the instructional designer (optional)
            template_path: Path to the template PPTX file (optional override)

        Visual output:
            Sets up an empty presentation with these dimensions:
            +------------------------------------------+
            |                                          |
            |          12192000 EMU (33.87cm)           |
            |                                          |
            |  6858000 EMU                             |
            |  (19.05cm)                               |
            |                                          |
            +------------------------------------------+
        """
        # Resolve the template file path
        if template_path and os.path.exists(template_path):
            tpl_path = template_path
        else:
            # Default: look for the template in the standard download location
            tpl_path = os.path.join(
                os.path.expanduser("~"),
                "Downloads",
                "storyboard template",
                "قالب المحاضرة التفاعلية- عربي.pptx",
            )

        # Open the template as the base presentation — this gives us all
        # the layout backgrounds, header bars, footer bars, and logos
        if os.path.exists(tpl_path):
            self.prs = Presentation(tpl_path)
            # Delete ALL existing example slides from the template
            # The rId attribute uses a namespace prefix, so we use the full URI
            _REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
            while len(self.prs.slides) > 0:
                sld_id_elem = self.prs.slides._sldIdLst[0]
                rId = sld_id_elem.get(f'{{{_REL_NS}}}id')
                self.prs.part.drop_rel(rId)
                self.prs.slides._sldIdLst.remove(sld_id_elem)
        else:
            # Fallback: create blank presentation if template not found
            self.prs = Presentation()
            self.prs.slide_width = SLIDE_WIDTH
            self.prs.slide_height = SLIDE_HEIGHT

        # Store the assets directory path for PNG images
        self.assets_dir = TEMPLATE_PATH

        # Store project metadata for reuse across slides
        self.project_code = project_code
        self.unit_number = unit_number
        self.unit_name = unit_name
        self.institution = institution
        self.designer = designer

        # Build the lecture title string that appears on every slide (2+)
        # Format: "المحاضرة [N]: [unit_name]"
        self.lecture_title = ""  # Will be set by add_title_slide

        # Track slide count for automatic page numbering
        self.slide_count = 0

        # Layout variant cycle for content slides (0=A, 1=B, 2=C)
        self._content_layout_cycle = 0

    # -----------------------------------------------------------------------
    # PUBLIC METHODS — Each adds one slide type
    # -----------------------------------------------------------------------

    def add_title_slide(
        self,
        title: str,
        subtitle: str = "",
        start_button_text: str = "ابدأ المحاضرة",
    ):
        """
        Add the opening title slide (matches template slide 1).

        This is the first thing learners see — it shows the institution
        name, lecture title, and a "Start" button.

        Args:
            title: Main lecture title (e.g., "المحاضرة الأولى")
            subtitle: Subtitle text (e.g., "مقدمة في علوم الحاسوب")
            start_button_text: Text for the start button (default: "ابدأ المحاضرة")

        Visual output (ASCII mockup):
            +------------------------------------------+
            |  [=== Blue header decoration bar ===]    |
            |                                          |
            |                                          |
            |              [Institution Name]          |
            |                                          |
            |              [Lecture Title:]             |
            |              [Subtitle]                  |
            |                                          |
            |              [  ابدأ المحاضرة  ]         |
            +------------------------------------------+

        Example:
            >>> builder.add_title_slide(
            ...     title="المحاضرة الأولى:",
            ...     subtitle="المهارات الرقميّة: المشهد التحوليّ"
            ... )
        """
        self.slide_count += 1

        # Store the lecture title for reuse on other slides
        # Combines title + subtitle for the top bar on subsequent slides
        self.lecture_title = f"{title} {subtitle}".strip()

        # Use Layout 0 ("Title Slide") — has background image, logo, etc.
        slide = self._add_slide_with_layout(0)

        # Set hidden TOC title for Storyline sidebar menu
        self._set_slide_title_for_toc(slide, title)

        # --- Institution name ---
        # Positioned in the right-center area (RTL layout puts content on right)
        self._add_arabic_textbox(
            slide,
            left=6096000,       # ~16.93cm from left
            top=3198167,        # ~8.88cm from top
            width=5181600,      # ~14.39cm wide
            height=461665,      # ~1.28cm tall
            text=self.institution,
            font_name=FONT_EXTRABOLD,
            font_size=Pt(24),
            bold=False,
            color=PRIMARY_BLUE,
            alignment=PP_ALIGN.CENTER,
            name="txt_institution",
        )

        # --- Lecture title + subtitle in a single textbox (template uses 1 box) ---
        title_box = slide.shapes.add_textbox(6096000, 4257368, 5181600, 1077218)
        title_box.name = "txt_title"
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.margin_left = TEXT_MARGIN_LR
        tf.margin_right = TEXT_MARGIN_LR
        tf.margin_top = TEXT_MARGIN_TB
        tf.margin_bottom = TEXT_MARGIN_TB

        # Title paragraph
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        run1 = p1.add_run()
        run1.text = title
        self._set_run_font(run1, FONT_EXTRABOLD, Pt(24), False, PRIMARY_BLUE)
        self._set_rtl(p1)

        # Subtitle paragraph (in the same textbox, as 3rd paragraph)
        if subtitle:
            # Empty line between title and subtitle
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = subtitle
            self._set_run_font(run2, FONT_EXTRABOLD, Pt(20), False, SUBTITLE_TEXT)
            self._set_rtl(p2)

        # --- Start button ---
        # Rounded rectangle with accent1 blue fill (#156082) and dark border
        button = self._add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left=7398084,       # ~20.55cm
            top=5599525,        # ~15.55cm
            width=2773680,      # ~7.7cm
            height=665193,      # ~1.85cm
            fill_color=ACCENT1_BLUE,   # accent1 #156082 (not #2D588C)
            border_color=BUTTON_BORDER,
            border_width=Pt(1.5),
            name="btn_start",
        )
        # Add text to the button
        tf_btn = button.text_frame
        tf_btn.word_wrap = True
        tf_btn.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf_btn.margin_left = TEXT_MARGIN_LR
        tf_btn.margin_right = TEXT_MARGIN_LR
        tf_btn.margin_top = TEXT_MARGIN_TB
        tf_btn.margin_bottom = TEXT_MARGIN_TB
        p = tf_btn.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = start_button_text
        self._set_run_font(run, FONT_REGULAR, Pt(20), False, WHITE)
        self._set_rtl(p)

        # --- Play icon (triangle) to the right of the button ---
        play_path = os.path.join(self.assets_dir, ASSET_PLAY_ICON)
        if os.path.exists(play_path):
            pic = slide.shapes.add_picture(
                play_path,
                9476078,    # left
                5599525,    # top
                619211,     # width
                657317,     # height
            )
            pic.name = "icon_play"

        # --- Hand cursor icon below the button ---
        hand_path = os.path.join(self.assets_dir, ASSET_HAND_CURSOR)
        if os.path.exists(hand_path):
            pic = slide.shapes.add_picture(
                hand_path,
                7570916,    # left
                5888428,    # top
                724001,     # width
                752580,     # height
            )
            pic.name = "icon_hand"

        # Add Storyline import instructions as speaker notes
        self._add_notes(slide, self._build_import_instructions(title))

    def add_objectives_slide(self, objectives: list):
        """
        Add a Learning Objectives slide (matches template slide 2).

        Shows numbered objectives in RTL with colored accent bars as
        row backgrounds for each objective.

        Args:
            objectives: List of objective strings in Arabic.
                        Typically 3-5 objectives per lecture.

        Visual output (ASCII mockup):
            +------------------------------------------+
            | [Title Bar: Lecture Name]                 |
            |          [الأهداف التعليمية]              |
            | يتوقع منك في نهاية هذه المحاضرة...       |
            |                                          |
            | [===== Objective Row 1 (colored bg) =====]|
            | [===== Objective Row 2 (colored bg) =====]|
            | [===== Objective Row 3 (colored bg) =====]|
            |                                          |
            | [2]                                       |
            +------------------------------------------+

        Example:
            >>> builder.add_objectives_slide([
            ...     "تعريف ماهي التقنية الناشئة.",
            ...     "التعرف إلى فوائد التقنيات الرقميّة.",
            ...     "اكتشاف عيوب التقنية الرقميّة.",
            ... ])
        """
        self.slide_count += 1
        slide = self._add_content_slide_with_layout()

        # Set hidden TOC title for Storyline sidebar menu
        self._set_slide_title_for_toc(slide, "الأهداف التعليمية")

        # --- Lecture title bar at top ---
        self._add_header_bar(slide, self.lecture_title)

        # --- Section banner (PNG image instead of colored rectangle) ---
        self._add_section_banner(slide, "الأهداف التعليمية")

        # --- Intro text ---
        # "يتوقع منك في نهاية هذه المحاضرة أن تكون قادرًا على:"
        self._add_arabic_textbox(
            slide,
            left=6280654,       # ~17.45cm
            top=1830945,        # ~5.09cm
            width=5361940,      # ~14.89cm
            height=369332,
            text="يتوقع منك في نهاية هذه المحاضرة أن تكون قادرًا على:",
            font_name=FONT_MEDIUM,
            font_size=Pt(18),
            bold=False,
            color=BODY_TEXT,
            alignment=PP_ALIGN.RIGHT,
            name="txt_obj_intro",
        )

        # --- Objective rows ---
        # Each objective gets a gradient PNG background bar + target icon + text
        # Adaptive spacing: fits up to 8 objectives without overflow
        row_top_start = 2315612     # ~6.43cm from top
        row_left = 612770           # ~1.7cm from left
        row_width = 11029824        # ~30.64cm wide
        safe_bottom = 6300000       # Safe zone above page number

        # Adaptive row sizing — shrinks to fit more items
        preferred_row_height = 600002   # ~1.67cm
        row_height, row_spacing = self._calculate_adaptive_spacing(
            item_count=len(objectives),
            available_top=row_top_start,
            available_bottom=safe_bottom,
            min_item_height=preferred_row_height,
        )
        # row_spacing here is the full gap; convert to "step" = height + gap
        row_step = row_height + row_spacing

        # Paths for PNG row assets
        row_img_path = os.path.join(self.assets_dir, ASSET_OBJECTIVE_ROW)
        icon_img_path = os.path.join(self.assets_dir, ASSET_TARGET_ICON)

        for i, objective in enumerate(objectives):
            row_top = row_top_start + (i * row_step)

            # Background gradient bar (image6.png) — the template uses a PNG
            obj_num = i + 1
            if os.path.exists(row_img_path):
                pic = slide.shapes.add_picture(
                    row_img_path,
                    row_left,
                    row_top,
                    row_width,
                    row_height,
                )
                pic.name = f"bg_obj_{obj_num}"
            else:
                # Fallback: colored rectangle if PNG not found
                self._add_shape(
                    slide,
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left=row_left,
                    top=row_top,
                    width=row_width,
                    height=row_height,
                    fill_color=LIGHT_BLUE_BG,
                    name=f"bg_obj_{obj_num}",
                )

            # Target/circle icon at the right end of the row (image13.png)
            icon_left = 10922693   # from template
            icon_width = 703228
            if os.path.exists(icon_img_path):
                pic = slide.shapes.add_picture(
                    icon_img_path,
                    icon_left,
                    row_top,
                    icon_width,
                    row_height,
                )
                pic.name = f"icon_obj_{obj_num}"

            # Objective text — positioned within the row
            text_left = 1462617   # ~4.06cm
            text_width = 9443403  # ~26.23cm
            text_height = 338554  # ~0.94cm
            # Center text vertically within the row
            text_top = row_top + (row_height - text_height) // 2

            self._add_arabic_textbox(
                slide,
                left=text_left,
                top=text_top,
                width=text_width,
                height=text_height,
                text=objective,
                font_name=FONT_REGULAR,
                font_size=Pt(18),
                bold=False,
                color=BODY_TEXT,
                alignment=PP_ALIGN.RIGHT,
                name=f"txt_obj_{obj_num}",
            )

    def add_content_slide(
        self,
        title: str,
        bullets: Optional[list] = None,
        paragraphs: Optional[list] = None,
        image_placeholder: Optional[str] = None,
        image_path: Optional[str] = None,
        image_prompt: Optional[str] = None,
        notes: str = "",
    ):
        """
        Add a content slide with a header and body text (matches template slide 3).

        This is the main workhorse slide type — it presents topic content
        with a title banner and body text. Can optionally include an image
        area on the left side.

        Args:
            title: Section title (e.g., "المقدمة")
            bullets: List of bullet point strings (use this OR paragraphs)
            paragraphs: List of paragraph strings (use this OR bullets)
            image_placeholder: Optional text describing what image to add (gray box)
            image_path: Optional path to a real image file (PNG/JPG). Takes
                        priority over image_placeholder if both are given.
            notes: Speaker notes / Storyline instructions (added to slide notes)

        Visual output (ASCII mockup):
            +------------------------------------------+
            | [Title Bar: Lecture Name]                 |
            |          [Section Title]                  |
            |                                          |
            |  [Image]   |  • Bullet point 1           |
            |  [Area]    |  • Bullet point 2           |
            |            |  • Bullet point 3           |
            |                                          |
            | [3]                                       |
            +------------------------------------------+

        Example:
            >>> builder.add_content_slide(
            ...     title="المقدمة",
            ...     bullets=[
            ...         "أصبحت التقنية الرقمية جزءاً من حياتنا",
            ...         "تؤثر على جميع المجالات",
            ...     ],
            ...     image_path="output/NJR01/U02/images/intro.png",
            ... )
        """
        self.slide_count += 1
        slide = self._add_content_slide_with_layout()

        # Set hidden TOC title for Storyline sidebar menu
        self._set_slide_title_for_toc(slide, title)

        # --- Lecture title bar at top ---
        self._add_header_bar(slide, self.lecture_title)

        # --- Section banner ---
        self._add_section_banner(slide, title)

        # --- Image area (left side) ---
        # Priority: real image > auto-generated > gray placeholder > no image
        # Auto-generate image if prompt provided but no path
        if image_prompt and not image_path:
            image_path = self._generate_image_for_slide(image_prompt, "content")
        has_image = False

        if image_path and os.path.exists(image_path):
            # Real image — add it with aspect ratio preservation
            self._add_image(
                slide, image_path,
                left=Cm(2.5), top=Cm(5.5),
                max_width=Cm(9), max_height=Cm(9),
                name="img_content",
            )
            has_image = True
        elif image_placeholder:
            # Fallback: gray placeholder rectangle with text label
            img_shape = self._add_shape(
                slide,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left=Cm(2.5),
                top=Cm(5.5),
                width=Cm(9),
                height=Cm(9),
                fill_color=RGBColor(0xE0, 0xE0, 0xE0),
                border_color=RGBColor(0xBD, 0xBD, 0xBD),
            )
            tf = img_shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = image_placeholder
            self._set_run_font(run, FONT_REGULAR, Pt(12), False, BODY_TEXT)
            has_image = True

        # --- Content body (with layout variants for visual variety) ---
        variant = self._content_layout_cycle % 3
        self._content_layout_cycle += 1

        content_top = Cm(5)
        content_height = Cm(11.5)

        if has_image:
            # Image mode — always use Variant A (full-width card) since layout is already varied
            variant = 0
            content_left = Cm(13)
            content_width = Cm(18)
        else:
            content_left = Cm(3)
            content_width = Cm(28)

        if variant == 1 and not has_image:
            # --- Variant B: Accent stripe on right, narrower content ---
            self._add_accent_stripe(slide)
            content_width = Cm(26)  # Slightly narrower to make room for stripe

        if variant == 2 and bullets and not has_image:
            # --- Variant C: Numbered points instead of bullet list ---
            self._add_numbered_points(slide, bullets, start_top=content_top + Cm(0.5))
        else:
            # --- Variant A or B: Card with bullets/paragraphs ---
            if bullets:
                card_shape = self._add_shape(
                    slide,
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left=content_left - Cm(0.5),
                    top=content_top - Cm(0.3),
                    width=content_width + Cm(1),
                    height=content_height + Cm(0.6),
                    fill_color=CONTENT_CARD_BG,
                    border_color=CONTENT_CARD_BORDER,
                    border_width=Pt(1),
                    name="bg_content_card",
                    corner_radius=0.04,
                )
                self._add_shadow_to_shape(card_shape, blur_pt=4, opacity_pct=15)

                self._add_bullet_list(
                    slide,
                    left=content_left,
                    top=content_top,
                    width=content_width,
                    height=content_height,
                    items=bullets,
                    font_size=Pt(20),
                    name="txt_body",
                )
            elif paragraphs:
                card_shape = self._add_shape(
                    slide,
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left=content_left - Cm(0.5),
                    top=content_top - Cm(0.3),
                    width=content_width + Cm(1),
                    height=content_height + Cm(0.6),
                    fill_color=CONTENT_CARD_BG,
                    border_color=CONTENT_CARD_BORDER,
                    border_width=Pt(1),
                    name="bg_content_card",
                    corner_radius=0.04,
                )
                self._add_shadow_to_shape(card_shape, blur_pt=4, opacity_pct=15)

                text = "\n\n".join(paragraphs)
                self._add_arabic_textbox(
                    slide,
                    left=content_left,
                    top=content_top,
                    width=content_width,
                    height=content_height,
                    text=text,
                    font_name=FONT_REGULAR,
                    font_size=Pt(18),
                    bold=False,
                    color=BODY_TEXT,
                    alignment=PP_ALIGN.RIGHT,
                    word_wrap=True,
                    auto_size=MSO_AUTO_SIZE.NONE,
                    name="txt_body",
                )

        # --- Speaker notes ---
        if notes:
            self._add_notes(slide, notes)

    def add_content_with_cards(
        self,
        title: str,
        cards: list,
        notes: str = "",
    ):
        """
        Add a content slide with 2-4 card layout (matches template slide 4 pattern).

        Cards are used for concepts like "examples of emerging technologies"
        where each card represents one concept with a title and optional body.

        Args:
            title: Section title (e.g., "أمثلة على التقنيات الناشئة")
            cards: List of dicts, each with:
                   - "title": Card title text
                   - "body": Card body text (optional)
                   - "color": RGBColor for the card (optional, auto-assigned)
                   - "image": Path to thumbnail image file (optional)
            notes: Speaker notes / Storyline instructions

        Visual output (ASCII mockup):
            +------------------------------------------+
            | [Title Bar: Lecture Name]                 |
            |          [Section Title]                  |
            |                                          |
            | +--------+ +--------+ +--------+         |
            | | Card 1 | | Card 2 | | Card 3 |         |
            | |  body  | |  body  | |  body  |         |
            | +--------+ +--------+ +--------+         |
            |                                          |
            | [6]                                       |
            +------------------------------------------+

        Example:
            >>> builder.add_content_with_cards(
            ...     title="أمثلة على التقنيات الناشئة",
            ...     cards=[
            ...         {"title": "الذكاء الاصطناعي", "body": "AI وصف"},
            ...         {"title": "إنترنت الأشياء", "body": "IoT وصف"},
            ...         {"title": "الحوسبة السحابية", "body": "Cloud وصف"},
            ...     ]
            ... )
        """
        self.slide_count += 1
        slide = self._add_content_slide_with_layout()

        # Set hidden TOC title for Storyline sidebar menu
        self._set_slide_title_for_toc(slide, title)

        # --- Lecture title bar ---
        self._add_header_bar(slide, self.lecture_title)

        # --- Section banner ---
        self._add_section_banner(slide, title, wide=True)

        # --- Card layout ---
        # Calculate card dimensions based on count
        card_count = len(cards)
        # Extended palette for 4+ cards — includes TEAL and AMBER for variety
        CARD_DARK2 = RGBColor(0x0E, 0x28, 0x41)  # dark navy from template theme
        AMBER = RGBColor(0xFF, 0x8F, 0x00)        # amber accent
        default_colors = [PRIMARY_BLUE, ACCENT1_BLUE, TEAL, CARD_DARK2, AMBER, PRIMARY_BLUE]

        # Layout area for cards
        cards_area_left = Cm(2.5)
        cards_area_width = Cm(28.5)
        cards_top = Cm(5.5)
        card_height = Cm(9)  # Taller to fit larger text (was Cm(8))

        # Calculate card width with gaps
        gap = Cm(0.8)
        total_gaps = gap * (card_count - 1) if card_count > 1 else 0
        card_width = int((cards_area_width - total_gaps) / card_count)

        for i, card_data in enumerate(cards):
            card_num = i + 1
            card_left = int(cards_area_left + i * (card_width + gap))
            card_color = card_data.get("color", default_colors[i % len(default_colors)])

            # Card background rectangle — light tinted fill instead of pure white
            card_shape = self._add_shape(
                slide,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left=card_left,
                top=cards_top,
                width=card_width,
                height=card_height,
                fill_color=CARD_LIGHT_BG,
                border_color=card_color,
                border_width=Pt(2),
                name=f"card_{card_num}",
            )
            # Real shadow effect on card
            self._add_shadow_to_shape(card_shape)

            # Thicker colored accent bar at top of card (Cm(1.2) for visual impact)
            self._add_shape(
                slide,
                MSO_SHAPE.RECTANGLE,
                left=card_left,
                top=cards_top,
                width=card_width,
                height=Cm(1.2),
                fill_color=card_color,
            )

            # Optional card thumbnail image (below accent bar, above title)
            # Auto-generate from prompt if no image path provided
            card_image = card_data.get("image")
            card_image_prompt = card_data.get("image_prompt")
            if card_image_prompt and not card_image:
                card_image = self._generate_image_for_slide(
                    card_image_prompt, "card",
                    topic_key=f"card_{i+1}" if not card_data.get("title") else None,
                )
            card_has_image = False
            if card_image:
                pic = self._add_image(
                    slide, card_image,
                    left=card_left + Cm(0.5),
                    top=cards_top + Cm(1.5),
                    max_width=card_width - Cm(1),
                    max_height=Cm(3),
                    name=f"img_card_{card_num}",
                )
                if pic is not None:
                    card_has_image = True

            # Vertical offset: shift title/body down when image is present
            title_top = cards_top + Cm(4.8) if card_has_image else cards_top + Cm(1.2)
            body_top = cards_top + Cm(6.5) if card_has_image else cards_top + Cm(3)
            body_height = Cm(2.5) if card_has_image else Cm(5.5)

            # Card title — Pt(20) for QM compliance, vertically centered
            self._add_arabic_textbox(
                slide,
                left=card_left + Cm(0.5),
                top=title_top,
                width=card_width - Cm(1),
                height=Cm(1.5),
                text=card_data.get("title", ""),
                font_name=FONT_EXTRABOLD,
                font_size=Pt(20),
                bold=False,
                color=card_color if isinstance(card_color, RGBColor) else BODY_TEXT,
                alignment=PP_ALIGN.CENTER,
                name=f"txt_card_{card_num}_title",
            )

            # Card body (if provided) — Pt(18) for QM compliance (was Pt(14))
            body = card_data.get("body", "")
            if body:
                self._add_arabic_textbox(
                    slide,
                    left=card_left + Cm(0.5),
                    top=body_top,
                    width=card_width - Cm(1),
                    height=body_height,
                    text=body,
                    font_name=FONT_REGULAR,
                    font_size=Pt(18),
                    bold=False,
                    color=BODY_TEXT,
                    alignment=PP_ALIGN.RIGHT,
                    word_wrap=True,
                    auto_size=MSO_AUTO_SIZE.NONE,
                    name=f"txt_card_{card_num}_body",
                )

        # --- Speaker notes ---
        if notes:
            self._add_notes(slide, notes)

    def add_section_divider(
        self,
        section_title: str,
        section_subtitle: str = "",
        section_number: int = None,
        total_sections: int = None,
        image_path: Optional[str] = None,
        image_prompt: Optional[str] = None,
    ):
        """
        Add a full-color section transition slide with decorative elements.

        Used to mark the boundary between major sections of the lecture
        (e.g., transitioning from "Introduction" to "Main Content").

        Enhanced with:
        - Decorative corner images (top-right, bottom-left) from template assets
        - Thicker accent bar on the right side (RTL primary side)
        - Progress dots showing current section position
        - Optional background illustration

        Args:
            section_title: Main section title
            section_subtitle: Optional subtitle text
            section_number: Current section number (1-based) for progress dots
            total_sections: Total number of sections for progress dots
            image_path: Optional path to a subtle background illustration (PNG)

        Visual output (ASCII mockup):
            +------------------------------------------+
            |                             [corner_tr] █|
            |  ████████████████████████████████████ █  |
            |  ██                              ██  █  |
            |  ██      [Section Title]         ██  █  |
            |  ██      [Section Subtitle]      ██  █  |
            |  ██                              ██     |
            |  ████████████████████████████████████    |
            | [corner_bl]       ● ● ◉ ● ●            |
            +------------------------------------------+

        Example:
            >>> builder.add_section_divider(
            ...     section_title="المحور الثاني",
            ...     section_subtitle="فوائد التقنية الرقمية",
            ...     section_number=2,
            ...     total_sections=5,
            ... )
        """
        self.slide_count += 1
        slide = self._add_content_slide_with_layout()

        # Set hidden TOC title for Storyline sidebar menu
        self._set_slide_title_for_toc(slide, section_title)

        # --- Bold PRIMARY_BLUE background for visual impact ---
        # Full-color rectangle covering most of the slide
        card_margin_h = Cm(2)
        card_margin_v = Cm(2)
        self._add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left=card_margin_h,
            top=card_margin_v,
            width=SLIDE_WIDTH - (card_margin_h * 2),
            height=SLIDE_HEIGHT - (card_margin_v * 2),
            fill_color=DIVIDER_BG,
            name="bg_divider",
        )

        # --- Subtle depth overlay at bottom edge (darker strip for depth) ---
        self._add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            left=card_margin_h,
            top=SLIDE_HEIGHT - Cm(4),
            width=SLIDE_WIDTH - (card_margin_h * 2),
            height=Cm(2),
            fill_color=PRIMARY_BLUE_DARK,
            name="bg_divider_depth",
        )

        # --- Refined accent bar on the right (RTL primary side) ---
        self._add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            left=SLIDE_WIDTH - Cm(2.5),
            top=Cm(2),
            width=Cm(0.3),
            height=SLIDE_HEIGHT - Cm(4),
            fill_color=WHITE,
            name="accent_right_bar",
        )

        # --- Optional background illustration (behind text, above bg) ---
        # Auto-generate image if prompt provided but no path
        if image_prompt and not image_path:
            image_path = self._generate_image_for_slide(image_prompt, "section")
        if image_path:
            self._add_image(
                slide, image_path,
                left=Cm(3), top=Cm(3),
                max_width=Cm(8), max_height=Cm(8),
                name="img_section_bg",
            )

        # --- Decorative corners — code-drawn shapes (no blurry PNGs) ---
        self._add_decorative_corner(slide, "top_right", WHITE, Cm(4))
        self._add_decorative_corner(slide, "bottom_left", WHITE, Cm(4))

        # --- Section title — large white text, moved up for better balance ---
        title_box = self._add_arabic_textbox(
            slide,
            left=Cm(4),
            top=Cm(5),
            width=Cm(26),
            height=Cm(3.5),
            text=section_title,
            font_name=FONT_EXTRABOLD,
            font_size=Pt(40),
            bold=False,
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
            name="txt_section_title",
        )
        title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # --- Thin decorative white line — increased gap below title ---
        self._add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            left=Cm(9),
            top=Cm(8.7),
            width=Cm(15),
            height=Cm(0.08),
            fill_color=WHITE,
            name="divider_line",
        )

        # --- Section subtitle — more room between line and subtitle ---
        if section_subtitle:
            self._add_arabic_textbox(
                slide,
                left=Cm(4),
                top=Cm(9.3),
                width=Cm(26),
                height=Cm(2.5),
                text=section_subtitle,
                font_name=FONT_MEDIUM,
                font_size=Pt(24),
                bold=False,
                color=WHITE,
                alignment=PP_ALIGN.CENTER,
                name="txt_section_subtitle",
            )

        # --- Progress dots showing current section position ---
        if section_number is not None and total_sections is not None:
            dot_size = Cm(0.6)
            dot_gap = Cm(1)
            total_width = total_sections * dot_size + (total_sections - 1) * dot_gap
            dots_left = (SLIDE_WIDTH - total_width) // 2
            dots_top = SLIDE_HEIGHT - Cm(3)

            for i in range(total_sections):
                dot_left = int(dots_left + i * (dot_size + dot_gap))
                is_current = (i + 1) == section_number
                dot_color = WHITE if is_current else RGBColor(0x80, 0x9F, 0xBF)
                dot_shape_size = Cm(0.8) if is_current else dot_size
                dot_offset = (dot_shape_size - dot_size) // 2 if is_current else 0

                self._add_shape(
                    slide,
                    MSO_SHAPE.OVAL,
                    left=dot_left - dot_offset,
                    top=dots_top - dot_offset,
                    width=dot_shape_size,
                    height=dot_shape_size,
                    fill_color=dot_color,
                    name=f"dot_section_{i + 1}",
                )

    def add_quiz_slide(
        self,
        question: str,
        options: list,
        correct_index: int,
        quiz_number: int = 1,
        total_quizzes: int = 5,
        image_path: Optional[str] = None,
        image_prompt: Optional[str] = None,
    ):
        """
        Add an MCQ quiz slide with dark background (matches activity pattern).

        Shows a multiple-choice question with lettered option badges.
        The correct answer is stored in the slide notes.

        Args:
            question: The question text in Arabic
            options: List of answer option strings (2-4 options)
            correct_index: Zero-based index of the correct answer
            quiz_number: Which quiz number this is (for display)
            total_quizzes: Total number of quizzes (for display)
            image_path: Optional illustration next to the question text

        Visual output (ASCII mockup):
            +------------------------------------------+
            | [Title Bar: Lecture Name]                 |
            | [نشاط X.Y (---)]                         |
            |                                          |
            | Question text goes here?                 |
            |                                          |
            |  [أ] Option 1                            |
            |  [ب] Option 2                            |
            |  [ج] Option 3                            |
            |  [د] Option 4                            |
            |                                          |
            | [5]                                       |
            +------------------------------------------+

        Example:
            >>> builder.add_quiz_slide(
            ...     question="أي من العبارات الآتية تعبر بدقة عن التقنية الناشئة؟",
            ...     options=[
            ...         "لأنها تستخدم أجهزة حديثة",
            ...         "لأنها تقدم طرقاً جديدة",
            ...         "لأنها مرتبطة بالإنترنت",
            ...     ],
            ...     correct_index=1,
            ... )
        """
        self.slide_count += 1
        slide = self._add_content_slide_with_layout()

        # Set hidden TOC title for Storyline sidebar menu
        self._set_slide_title_for_toc(slide, f"نشاط: {question[:30]}")

        # --- Lecture title bar ---
        self._add_header_bar(slide, self.lecture_title)

        # --- Activity banner — descriptive label per QM 5.2 ---
        activity_title = "نشاط تفاعلي: اختيار من متعدد"
        self._add_section_banner(slide, activity_title, wide=True)

        # --- Optional question illustration (left side) ---
        # Auto-generate image if prompt provided but no path
        if image_prompt and not image_path:
            image_path = self._generate_image_for_slide(image_prompt, "quiz")
        quiz_has_image = False
        if image_path:
            pic = self._add_image(
                slide, image_path,
                left=Cm(2), top=Cm(4),
                max_width=Cm(7), max_height=Cm(5),
                name="img_quiz",
            )
            if pic is not None:
                quiz_has_image = True

        # --- Question text — Pt(24) bold for emphasis ---
        # Narrower when image is present (shifts right to make room)
        q_left = Cm(10) if quiz_has_image else Cm(2.5)
        q_width = Cm(21) if quiz_has_image else Cm(29)
        self._add_arabic_textbox(
            slide,
            left=q_left,
            top=Cm(5),
            width=q_width,
            height=Cm(2),
            text=question,
            font_name=FONT_EXTRABOLD,
            font_size=Pt(24),
            bold=False,
            color=BODY_TEXT,
            alignment=PP_ALIGN.RIGHT,
            word_wrap=True,
            auto_size=MSO_AUTO_SIZE.NONE,
            name="txt_question",
        )

        # --- Answer options ---
        # Arabic letter badges: أ ب ج د
        arabic_letters = ["أ", "ب", "ج", "د"]
        option_top_start = Cm(7.5)
        option_spacing = Cm(2.2)
        option_height = Cm(1.7)

        for i, option_text in enumerate(options):
            opt_letter = arabic_letters[i] if i < len(arabic_letters) else str(i + 1)
            opt_id = ["a", "b", "c", "d"][i] if i < 4 else str(i + 1)
            option_top = int(option_top_start + i * option_spacing)

            # Option card with integrated badge (no separate accent border)
            option_bg = CONTENT_CARD_BG if i % 2 == 0 else WHITE
            option_bg_shape = self._add_shape(
                slide,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left=Cm(2.5),
                top=option_top - Cm(0.1),
                width=Cm(28),
                height=option_height + Cm(0.2),
                fill_color=option_bg,
                border_color=CONTENT_CARD_BORDER,
                border_width=Pt(0.5),
                name=f"bg_opt_{opt_id}",
                corner_radius=0.06,
            )
            self._add_shadow_to_shape(option_bg_shape, blur_pt=3, opacity_pct=12)

            # Letter badge INSIDE the card (right side for RTL)
            badge_left = Cm(28.5)
            badge_size = Cm(1.5)
            badge = self._add_shape(
                slide,
                MSO_SHAPE.OVAL,
                left=badge_left,
                top=option_top + Cm(0.1),
                width=badge_size,
                height=badge_size,
                fill_color=PRIMARY_BLUE,
                name=f"opt_{opt_id}",
            )
            tf = badge.text_frame
            tf.word_wrap = False
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = opt_letter
            self._set_run_font(run, FONT_EXTRABOLD, Pt(16), False, WHITE)

            # Option text
            self._add_arabic_textbox(
                slide,
                left=Cm(3),
                top=option_top,
                width=Cm(24.5),
                height=option_height,
                text=option_text,
                font_name=FONT_REGULAR,
                font_size=Pt(20),
                bold=False,
                color=BODY_TEXT,
                alignment=PP_ALIGN.RIGHT,
                word_wrap=True,
                auto_size=MSO_AUTO_SIZE.NONE,
                name=f"txt_opt_{opt_id}",
            )

        # --- "Check Answer" button at bottom — larger with more presence ---
        check_btn = self._add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left=Cm(10),
            top=Cm(16),
            width=Cm(10),
            height=Cm(1.6),
            fill_color=ACCENT1_BLUE,
            border_color=BUTTON_BORDER,
            border_width=Pt(1.5),
            name="btn_check",
        )
        tf_btn = check_btn.text_frame
        tf_btn.word_wrap = True
        tf_btn.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_btn = tf_btn.paragraphs[0]
        p_btn.alignment = PP_ALIGN.CENTER
        run_btn = p_btn.add_run()
        run_btn.text = "تحقق من الإجابة"
        self._set_run_font(run_btn, FONT_EXTRABOLD, Pt(22), False, WHITE)
        self._set_rtl(p_btn)

        # Feedback instruction moved to notes to avoid edge overflow

        # --- Structured notes for Storyline import ---
        correct_letter = arabic_letters[correct_index] if correct_index < len(arabic_letters) else str(correct_index + 1)
        correct_opt_id = ["a", "b", "c", "d"][correct_index] if correct_index < 4 else str(correct_index + 1)
        notes_text = (
            f"=== STORYLINE INSTRUCTIONS ===\n"
            f"Slide Type: Quiz - Multiple Choice\n"
            f"Correct Answer: {correct_letter} (opt_{correct_opt_id})\n"
            f"Feedback (Correct): احسنت! الاجابة صحيحة\n"
            f"Feedback (Incorrect): الاجابة غير صحيحة، حاول مرة اخرى\n"
            f"Points: 10\n"
            f"Attempts: 2\n\n"
            f"=== FREEFORM SETUP ===\n"
            f"1. Insert > Convert to Freeform > Pick One\n"
            f"2. Assign opt_a, opt_b, opt_c, opt_d as choices\n"
            f"3. Set opt_{correct_opt_id} as correct answer\n"
            f"4. btn_check triggers submit\n\n"
            f"=== NARRATOR SCRIPT ===\n"
            f"{question}"
        )
        self._add_notes(slide, notes_text)

    def add_drag_drop_slide(
        self,
        question: str,
        items: list,
        correct_order: list,
        quiz_number: int = 1,
    ):
        """
        Add a drag-and-drop interaction slide.

        Shows a classification or ordering activity where learners
        drag items to the correct positions.

        Args:
            question: Instruction text for the activity
            items: List of item strings that learners will drag
            correct_order: List showing the correct classification/order
            quiz_number: Activity number (for display)

        Visual output (ASCII mockup):
            +------------------------------------------+
            | [Title Bar: Lecture Name]                 |
            | [نشاط X.Y]                               |
            |                                          |
            | [Question/Instructions]                  |
            |                                          |
            | +------+ +------+ +------+               |
            | |Item 1| |Item 2| |Item 3|               |
            | +------+ +------+ +------+               |
            |                                          |
            | [9]                                       |
            +------------------------------------------+

        Example:
            >>> builder.add_drag_drop_slide(
            ...     question="صنف العبارات الآتية: فوائد أم سلبيات؟",
            ...     items=["تسريع الوصول", "الإدمان الرقمي", "زيادة الإنتاجية"],
            ...     correct_order=["فائدة", "سلبية", "فائدة"],
            ... )
        """
        self.slide_count += 1
        slide = self._add_content_slide_with_layout()

        # Set hidden TOC title for Storyline sidebar menu
        self._set_slide_title_for_toc(slide, "نشاط: سحب وإفلات")

        # --- Lecture title bar ---
        self._add_header_bar(slide, self.lecture_title)

        # --- Activity banner — descriptive label per QM 5.2 ---
        activity_title = "نشاط تفاعلي: سحب وترتيب"
        self._add_section_banner(slide, activity_title, wide=True)

        # --- Question text — Pt(20) bold for emphasis ---
        self._add_arabic_textbox(
            slide,
            left=Cm(2.5),
            top=Cm(5),
            width=Cm(29),
            height=Cm(2),
            text=question,
            font_name=FONT_EXTRABOLD,
            font_size=Pt(20),
            bold=False,
            color=BODY_TEXT,
            alignment=PP_ALIGN.RIGHT,
            word_wrap=True,
            auto_size=MSO_AUTO_SIZE.NONE,
            name="txt_question",
        )

        # --- Clear instruction text ---
        self._add_arabic_textbox(
            slide,
            left=Cm(2.5),
            top=Cm(7),
            width=Cm(29),
            height=Cm(1.2),
            text="اسحب العناصر التالية إلى الترتيب الصحيح",
            font_name=FONT_MEDIUM,
            font_size=Pt(18),
            bold=False,
            color=ACCENT1_BLUE,
            alignment=PP_ALIGN.RIGHT,
            word_wrap=True,
            auto_size=MSO_AUTO_SIZE.NONE,
            name="txt_instruction",
        )

        # --- Draggable items (left side) ---
        item_count = len(items)
        items_area_left = Cm(2.5)
        items_top = Cm(8.5)
        item_width = Cm(12)
        safe_bottom = 6300000  # Safe zone above page number

        # Adaptive spacing — shrinks items to fit more
        item_height, gap = self._calculate_adaptive_spacing(
            item_count=item_count,
            available_top=items_top,
            available_bottom=safe_bottom,
            min_item_height=Cm(2),
        )
        item_step = item_height + gap

        for i, item_text in enumerate(items):
            # Stack items vertically
            item_top = int(items_top + i * item_step)

            item_shape = self._add_shape(
                slide,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left=items_area_left,
                top=item_top,
                width=item_width,
                height=item_height,
                fill_color=CONTENT_CARD_BG,
                border_color=PRIMARY_BLUE,
                border_width=Pt(1.5),
                name=f"drag_item_{i + 1}",
            )
            # Real shadow effect on drag item
            self._add_shadow_to_shape(item_shape)
            # Item text — Pt(20) for better readability
            tf = item_shape.text_frame
            tf.word_wrap = True
            tf.margin_left = TEXT_MARGIN_LR
            tf.margin_right = TEXT_MARGIN_LR
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = item_text
            self._set_run_font(run, FONT_REGULAR, Pt(20), False, BODY_TEXT)
            self._set_rtl(p)

            # Grip indicator (shows this item is draggable)
            self._add_arabic_textbox(
                slide,
                left=items_area_left + Cm(0.2),
                top=item_top + Cm(0.2),
                width=Cm(1),
                height=Cm(1),
                text="\u2630",  # ☰ trigram/hamburger icon
                font_name=FONT_REGULAR,
                font_size=Pt(14),
                bold=False,
                color=WARM_GRAY,
                alignment=PP_ALIGN.LEFT,
                word_wrap=False,
                auto_size=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
                name=f"icon_grip_{i + 1}",
            )

        # --- Numbered drop positions (right side) ---
        drop_left = Cm(18)
        drop_width = Cm(12)
        for i in range(item_count):
            drop_top = int(items_top + i * item_step)

            # Drop zone rectangle with subtle card styling
            drop_shape = self._add_shape(
                slide,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left=drop_left,
                top=drop_top,
                width=drop_width,
                height=item_height,
                name=f"drop_zone_{i + 1}",
                fill_color=CARD_LIGHT_BG,
                border_color=CONTENT_CARD_BORDER,
                border_width=Pt(1.5),
            )

            # Make drop zone border dashed
            drop_shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH

            # "Drag here" hint text inside drop zone
            drop_tf = drop_shape.text_frame
            drop_tf.word_wrap = True
            drop_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p_hint = drop_tf.paragraphs[0]
            p_hint.alignment = PP_ALIGN.CENTER
            run_hint = p_hint.add_run()
            run_hint.text = "اسحب هنا"
            self._set_run_font(run_hint, FONT_REGULAR, Pt(14), False, WARM_GRAY)
            self._set_rtl(p_hint)

            # Number badge in the drop zone
            badge_size = Cm(1.5)
            badge = self._add_shape(
                slide,
                MSO_SHAPE.OVAL,
                left=drop_left + drop_width - badge_size - Cm(0.3),
                top=drop_top + (item_height - badge_size) // 2,
                width=badge_size,
                height=badge_size,
                fill_color=PRIMARY_BLUE,
            )
            tf_b = badge.text_frame
            tf_b.word_wrap = False
            tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
            p_b = tf_b.paragraphs[0]
            p_b.alignment = PP_ALIGN.CENTER
            run_b = p_b.add_run()
            run_b.text = str(i + 1)
            self._set_run_font(run_b, FONT_EXTRABOLD, Pt(16), False, WHITE)

        # --- Structured notes for Storyline import ---
        mapping_lines = "\n".join(
            f"drag_item_{i+1} → drop_zone_{i+1}: {correct_order[i]}"
            for i in range(len(correct_order))
        )
        notes_text = (
            f"=== STORYLINE INSTRUCTIONS ===\n"
            f"Slide Type: Drag and Drop - Ordering\n"
            f"Correct Order:\n{mapping_lines}\n\n"
            f"=== FREEFORM SETUP ===\n"
            f"1. Insert > Convert to Freeform > Drag and Drop\n"
            f"2. Match drag_item shapes to drop_zone shapes\n"
            f"3. Set correct order as shown above\n\n"
            f"=== NARRATOR SCRIPT ===\n"
            f"{question}"
        )
        self._add_notes(slide, notes_text)

    def add_two_column_slide(
        self,
        title: str,
        left_title: str,
        left_points: list,
        right_title: str,
        right_points: list,
        notes: str = "",
        right_image: Optional[str] = None,
        left_image: Optional[str] = None,
        right_image_prompt: Optional[str] = None,
        left_image_prompt: Optional[str] = None,
    ):
        """
        Add a two-column comparison slide.

        Used for comparing concepts side-by-side (e.g., pros vs cons,
        before vs after, two categories).

        Args:
            title: Section title
            left_title: Title for the left column
            left_points: Bullet points for the left column
            right_title: Title for the right column
            right_points: Bullet points for the right column
            notes: Speaker notes
            right_image: Optional image above the right column header
            left_image: Optional image above the left column header

        Visual output (ASCII mockup):
            +------------------------------------------+
            | [Title Bar: Lecture Name]                 |
            |          [Section Title]                  |
            |                                          |
            |  [Right Title]    |    [Left Title]       |
            |  • Point 1        |    • Point 1          |
            |  • Point 2        |    • Point 2          |
            |  • Point 3        |    • Point 3          |
            |                                          |
            | [7]                                       |
            +------------------------------------------+

            Note: In RTL layout, "right" column appears on the right
            side of the slide (first visually for Arabic readers).

        Example:
            >>> builder.add_two_column_slide(
            ...     title="مقارنة بين الفوائد والسلبيات",
            ...     left_title="السلبيات",
            ...     left_points=["الإدمان الرقمي", "فقدان الخصوصية"],
            ...     right_title="الفوائد",
            ...     right_points=["تسريع الوصول", "زيادة الإنتاجية"],
            ... )
        """
        self.slide_count += 1
        slide = self._add_content_slide_with_layout()

        # Set hidden TOC title for Storyline sidebar menu
        self._set_slide_title_for_toc(slide, title)

        # --- Lecture title bar ---
        self._add_header_bar(slide, self.lecture_title)

        # --- Section banner ---
        self._add_section_banner(slide, title, wide=True)

        # --- Auto-generate column images from prompts ---
        if right_image_prompt and not right_image:
            right_image = self._generate_image_for_slide(right_image_prompt, "two_column")
        if left_image_prompt and not left_image:
            left_image = self._generate_image_for_slide(left_image_prompt, "two_column")

        # --- Column layout ---
        # Right column (primary in RTL) — positioned on the right side
        col_top = Cm(5)
        col_gap = Cm(1)
        col_width = Cm(13.5)

        right_col_left = Cm(17)  # Right side of slide
        left_col_left = Cm(2.5)  # Left side of slide

        # Check if images present — shifts content down
        right_has_img = right_image and os.path.exists(right_image)
        left_has_img = left_image and os.path.exists(left_image)
        any_has_img = right_has_img or left_has_img

        # When images are present, use shorter cards to prevent overflow
        # Without images: col_top(5) + col_height(10) + padding(0.6) = 15.6cm ✓
        # With images: col_top(5) + col_height(8) + img_shift(3.5) + padding(0.6) = 17.1cm ✓
        # Slide height = 19.05cm, so both fit with safe margin
        col_height = Cm(8) if any_has_img else Cm(10)
        img_shift = Cm(3.5)  # How much to shift content down for image

        # --- Right column card ---
        right_card_height = col_height + (img_shift if right_has_img else 0)
        right_card = self._add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left=right_col_left - Cm(0.3),
            top=col_top - Cm(0.3),
            width=col_width + Cm(0.6),
            height=right_card_height + Cm(0.6),
            fill_color=WHITE,
            border_color=CONTENT_CARD_BORDER,
            border_width=Pt(1),
            name="bg_col1_card",
            corner_radius=0.04,
        )
        self._add_shadow_to_shape(right_card, blur_pt=4, opacity_pct=12)

        # Optional right column header image (3cm max for proportional thumbnails)
        right_content_offset = 0
        if right_has_img:
            self._add_image(
                slide, right_image,
                left=right_col_left, top=col_top + Cm(0.3),
                max_width=col_width, max_height=Cm(3),
                name="img_col1",
            )
            right_content_offset = img_shift

        # Right column title — Pt(20) bold
        self._add_arabic_textbox(
            slide,
            left=right_col_left,
            top=col_top + right_content_offset,
            width=col_width,
            height=Cm(1.5),
            text=right_title,
            font_name=FONT_EXTRABOLD,
            font_size=Pt(20),
            bold=False,
            color=PRIMARY_BLUE,
            alignment=PP_ALIGN.CENTER,
            name="txt_col1_title",
        )

        # Accent bar under right title
        self._add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            left=right_col_left + Cm(2),
            top=col_top + Cm(1.5) + right_content_offset,
            width=col_width - Cm(4),
            height=Cm(0.15),
            fill_color=PRIMARY_BLUE,
        )

        # Right column bullets — Pt(18) minimum for QM (was Pt(15))
        self._add_bullet_list(
            slide,
            left=right_col_left,
            top=col_top + Cm(2) + right_content_offset,
            width=col_width,
            height=col_height - Cm(2),
            items=right_points,
            font_size=Pt(18),
            name="txt_col1_body",
        )

        # --- Left column card ---
        left_card_height = col_height + (img_shift if left_has_img else 0)
        left_card = self._add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left=left_col_left - Cm(0.3),
            top=col_top - Cm(0.3),
            width=col_width + Cm(0.6),
            height=left_card_height + Cm(0.6),
            fill_color=WHITE,
            border_color=CONTENT_CARD_BORDER,
            border_width=Pt(1),
            name="bg_col2_card",
            corner_radius=0.04,
        )
        self._add_shadow_to_shape(left_card, blur_pt=4, opacity_pct=12)

        # Optional left column header image (3cm max for proportional thumbnails)
        left_content_offset = 0
        if left_has_img:
            self._add_image(
                slide, left_image,
                left=left_col_left, top=col_top + Cm(0.3),
                max_width=col_width, max_height=Cm(3),
                name="img_col2",
            )
            left_content_offset = img_shift

        # Left column title — Pt(20) bold
        self._add_arabic_textbox(
            slide,
            left=left_col_left,
            top=col_top + left_content_offset,
            width=col_width,
            height=Cm(1.5),
            text=left_title,
            font_name=FONT_EXTRABOLD,
            font_size=Pt(20),
            bold=False,
            color=ACCENT1_BLUE,
            alignment=PP_ALIGN.CENTER,
            name="txt_col2_title",
        )

        # Accent bar under left title
        self._add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            left=left_col_left + Cm(2),
            top=col_top + Cm(1.5) + left_content_offset,
            width=col_width - Cm(4),
            height=Cm(0.15),
            fill_color=ACCENT1_BLUE,
        )

        # Left column bullets — Pt(18) minimum for QM (was Pt(15))
        self._add_bullet_list(
            slide,
            left=left_col_left,
            top=col_top + Cm(2) + left_content_offset,
            width=col_width,
            height=col_height - Cm(2),
            items=left_points,
            font_size=Pt(18),
            name="txt_col2_body",
        )

        if notes:
            self._add_notes(slide, notes)

    def add_summary_slide(self, summary_items: list):
        """
        Add a summary/recap slide (matches template slide 8).

        Shows key takeaways from the lecture in a list format
        with blue link-style text.

        Args:
            summary_items: List of dicts with:
                - "title": Summary item title (optional)
                - "text": Summary text

                OR simply a list of strings for plain summary points.

        Visual output (ASCII mockup):
            +------------------------------------------+
            | [Title Bar: Lecture Name]                 |
            | [ملخص الوحدة الدراسية]                   |
            |                                          |
            | • Summary point 1 text here              |
            | • Summary point 2 text here              |
            | • Summary point 3 text here              |
            |                                          |
            | [12]                                      |
            +------------------------------------------+

        Example:
            >>> builder.add_summary_slide([
            ...     "التقنيات الناشئة هي تقنيات في مراحلها الأولى",
            ...     "للتقنية الرقمية فوائد وسلبيات يجب مراعاتها",
            ... ])
        """
        self.slide_count += 1
        slide = self._add_content_slide_with_layout()

        # Set hidden TOC title for Storyline sidebar menu
        self._set_slide_title_for_toc(slide, "ملخص المحاضرة")

        # --- Lecture title bar ---
        self._add_header_bar(slide, self.lecture_title)

        # --- Section banner ---
        self._add_section_banner(slide, "ملخّص الوحدة الدراسيّة", wide=True)

        # --- Summary content card background for professional look ---
        self._add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left=Cm(2),
            top=Cm(4.2),
            width=Cm(30),
            height=Cm(10.5),
            fill_color=CONTENT_CARD_BG,
            border_color=CONTENT_CARD_BORDER,
            border_width=Pt(1),
            name="bg_summary_card",
        )

        # Build each item as a paragraph with bold label + regular text
        txBox = slide.shapes.add_textbox(Cm(2.5), Cm(4.5), Cm(29), Cm(10))
        txBox.name = "txt_summary"
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = TEXT_MARGIN_LR
        tf.margin_right = TEXT_MARGIN_LR
        tf.margin_top = TEXT_MARGIN_TB
        tf.margin_bottom = TEXT_MARGIN_TB

        for idx, item in enumerate(summary_items):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.alignment = PP_ALIGN.RIGHT
            p.line_spacing = 1.5
            p.space_after = Pt(8)

            if isinstance(item, dict):
                title = item.get("title", "")
                text = item.get("text", item.get("body", ""))
                if title:
                    # Bold label run (PRIMARY_BLUE for emphasis)
                    label_run = p.add_run()
                    label_run.text = f"{title}: "
                    self._set_run_font(label_run, FONT_EXTRABOLD, Pt(20), True, PRIMARY_BLUE)
                    # Regular text run (BODY_TEXT for readability)
                    text_run = p.add_run()
                    text_run.text = text
                    self._set_run_font(text_run, FONT_REGULAR, Pt(20), False, BODY_TEXT)
                else:
                    run = p.add_run()
                    run.text = text
                    self._set_run_font(run, FONT_REGULAR, Pt(20), False, BODY_TEXT)
            else:
                run = p.add_run()
                run.text = str(item)
                self._set_run_font(run, FONT_REGULAR, Pt(20), False, BODY_TEXT)

            self._set_rtl(p)

    def add_closing_slide(self, next_steps: list = None, image_path: Optional[str] = None, image_prompt: Optional[str] = None):
        """
        Add the final closing slide with optional next steps.

        This is the last slide of the lecture — it thanks the learner
        and optionally lists what comes next.

        Args:
            next_steps: Optional list of next step strings
            image_path: Optional decorative illustration above thank-you text

        Visual output (ASCII mockup):
            +------------------------------------------+
            |  ████████████████████████████████████████ |
            |  ██                                  ██  |
            |  ██         شكراً لكم                ██  |
            |  ██                                  ██  |
            |  ██   Next Steps:                    ██  |
            |  ██   • Step 1                       ██  |
            |  ██   • Step 2                       ██  |
            |  ██                                  ██  |
            |  ████████████████████████████████████████ |
            +------------------------------------------+

        Example:
            >>> builder.add_closing_slide([
            ...     "مراجعة المحاضرة التفاعلية",
            ...     "حل النشاط التفاعلي",
            ...     "الاستعداد للاختبار البعدي",
            ... ])
        """
        self.slide_count += 1
        slide = self._add_content_slide_with_layout()

        # Set hidden TOC title for Storyline sidebar menu
        self._set_slide_title_for_toc(slide, "شكراً لكم")

        # --- Full-slide colored background ---
        self._add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            left=0,
            top=0,
            width=SLIDE_WIDTH,
            height=SLIDE_HEIGHT,
            fill_color=PRIMARY_BLUE,
        )

        # --- White content area ---
        margin = Cm(3)
        self._add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left=margin,
            top=margin,
            width=SLIDE_WIDTH - (margin * 2),
            height=SLIDE_HEIGHT - (margin * 2),
            fill_color=WHITE,
            name="bg_card",
        )

        # --- Decorative corner — code-drawn shapes (no blurry PNGs) ---
        self._add_decorative_corner(slide, "top_right", PRIMARY_BLUE_LIGHT, Cm(3))

        # --- Optional decorative illustration above thank-you text ---
        # Auto-generate image if prompt provided but no path
        if image_prompt and not image_path:
            image_path = self._generate_image_for_slide(image_prompt, "closing")
        closing_has_image = False
        if image_path:
            pic = self._add_image(
                slide, image_path,
                left=Cm(9), top=Cm(3.5),
                max_width=Cm(15), max_height=Cm(5),
                name="img_closing",
            )
            if pic is not None:
                closing_has_image = True

        # --- Thank you text (shifts down when image present) ---
        thanks_top = Cm(9) if closing_has_image else Cm(5)
        self._add_arabic_textbox(
            slide,
            left=Cm(5),
            top=thanks_top,
            width=Cm(24),
            height=Cm(2.5),
            text="شكراً لكم",
            font_name=FONT_EXTRABOLD,
            font_size=Pt(36),
            bold=False,
            color=PRIMARY_BLUE,
            alignment=PP_ALIGN.CENTER,
            name="txt_thanks",
        )

        # --- Next steps ---
        if next_steps:
            # Accent line above next steps
            self._add_shape(
                slide,
                MSO_SHAPE.RECTANGLE,
                left=Cm(10),
                top=Cm(7.5),
                width=Cm(14),
                height=Cm(0.1),
                fill_color=PRIMARY_BLUE,
                name="accent_steps_line",
            )

            self._add_arabic_textbox(
                slide,
                left=Cm(5),
                top=Cm(8),
                width=Cm(24),
                height=Cm(1.5),
                text="الخطوات القادمة:",
                font_name=FONT_EXTRABOLD,
                font_size=Pt(20),
                bold=False,
                color=BODY_TEXT,
                alignment=PP_ALIGN.CENTER,
                name="txt_next_steps_label",
            )

            # Numbered next steps (circles instead of bullets)
            for i, step in enumerate(next_steps):
                step_num = i + 1
                step_top = Cm(10) + i * Cm(2)

                # Number circle
                circle = self._add_shape(
                    slide,
                    MSO_SHAPE.OVAL,
                    left=Cm(24),
                    top=step_top,
                    width=Cm(1.5),
                    height=Cm(1.5),
                    fill_color=PRIMARY_BLUE,
                    name=f"num_step_{step_num}",
                )
                tf_c = circle.text_frame
                tf_c.vertical_anchor = MSO_ANCHOR.MIDDLE
                p_c = tf_c.paragraphs[0]
                p_c.alignment = PP_ALIGN.CENTER
                run_c = p_c.add_run()
                run_c.text = str(step_num)
                self._set_run_font(run_c, FONT_EXTRABOLD, Pt(16), False, WHITE)

                # Step text
                self._add_arabic_textbox(
                    slide,
                    left=Cm(6),
                    top=step_top,
                    width=Cm(17),
                    height=Cm(1.5),
                    text=step,
                    font_name=FONT_REGULAR,
                    font_size=Pt(18),
                    bold=False,
                    color=BODY_TEXT,
                    alignment=PP_ALIGN.RIGHT,
                    word_wrap=True,
                    auto_size=MSO_AUTO_SIZE.NONE,
                    name=f"txt_step_{step_num}",
                )

    def add_slider_slide(
        self,
        title: str,
        items: list,
        notes: str = "",
    ):
        """
        Add a slider/scroll interaction slide.

        Used for content that appears one item at a time as the learner
        drags a slider (common pattern in Storyline).

        Args:
            title: Instruction text for the interaction
            items: List of dicts with "number" and "text" for each step
            notes: Speaker notes / Storyline instructions

        Visual output (ASCII mockup):
            +------------------------------------------+
            | [Title Bar: Lecture Name]                 |
            |                                          |
            | [Title/Instructions]                     |
            |                                          |
            | [1] Step text 1                          |
            | [2] Step text 2                          |
            | [3] Step text 3                          |
            |                                          |
            |     ◄════════════►  (slider)             |
            | [8]                                       |
            +------------------------------------------+

        Example:
            >>> builder.add_slider_slide(
            ...     title="يوجد العديد من فوائد التقنية الرقمية، اسحب المؤشر:",
            ...     items=[
            ...         {"number": "1", "text": "تسهيل التواصل"},
            ...         {"number": "2", "text": "تحسين الإنتاجية"},
            ...     ]
            ... )
        """
        self.slide_count += 1
        slide = self._add_content_slide_with_layout()

        # Set hidden TOC title for Storyline sidebar menu
        self._set_slide_title_for_toc(slide, title)

        # --- Lecture title bar ---
        self._add_header_bar(slide, self.lecture_title)

        # --- Title/Instructions ---
        self._add_arabic_textbox(
            slide,
            left=Cm(2.5),
            top=Cm(2.5),
            width=Cm(29),
            height=Cm(2),
            text=title,
            font_name=FONT_MEDIUM,
            font_size=Pt(18),
            bold=False,
            color=BODY_TEXT,
            alignment=PP_ALIGN.RIGHT,
            word_wrap=True,
            auto_size=MSO_AUTO_SIZE.NONE,
            name="txt_instruction",
        )

        # --- Numbered items ---
        item_top_start = Cm(5)
        item_spacing = Cm(2)

        for i, item_data in enumerate(items):
            item_top = int(item_top_start + i * item_spacing)
            slider_num = i + 1

            number = item_data.get("number", str(i + 1)) if isinstance(item_data, dict) else str(i + 1)
            text = item_data.get("text", str(item_data)) if isinstance(item_data, dict) else str(item_data)

            # Number badge
            badge = self._add_shape(
                slide,
                MSO_SHAPE.OVAL,
                left=Cm(28),
                top=item_top,
                width=Cm(1.5),
                height=Cm(1.5),
                fill_color=PRIMARY_BLUE,
                name=f"icon_step_{slider_num}",
            )
            tf = badge.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = number
            self._set_run_font(run, FONT_EXTRABOLD, Pt(16), False, WHITE)

            # Item text — Pt(18) for QM compliance
            self._add_arabic_textbox(
                slide,
                left=Cm(3),
                top=item_top,
                width=Cm(24),
                height=Cm(1.5),
                text=text,
                font_name=FONT_REGULAR,
                font_size=Pt(18),
                bold=False,
                color=BODY_TEXT,
                alignment=PP_ALIGN.RIGHT,
                word_wrap=True,
                auto_size=MSO_AUTO_SIZE.NONE,
                name=f"txt_step_{slider_num}",
            )

        if notes:
            self._add_notes(slide, notes)

    def add_click_reveal_slide(
        self,
        title: str,
        instruction: str,
        reveal_items: list,
        notes: str = "",
    ):
        """
        Add a click-to-reveal interaction slide.

        Used for content where learners click on categories/tabs to
        reveal descriptions (common for concept exploration).

        Args:
            title: Section title
            instruction: Instruction text
            reveal_items: List of dicts with "label" and "description"
            notes: Speaker notes

        Visual output (ASCII mockup):
            +------------------------------------------+
            | [Title Bar: Lecture Name]                 |
            |          [Section Title]                  |
            | [Instructions]                           |
            |                                          |
            | +--------+  +--------+  +--------+       |
            | | Label1 |  | Label2 |  | Label3 |       |
            | +--------+  +--------+  +--------+       |
            | [Description appears on click]           |
            |                                          |
            | [10]                                      |
            +------------------------------------------+

        Example:
            >>> builder.add_click_reveal_slide(
            ...     title="تأثير التقنية على الحياة",
            ...     instruction="انقر على كل جانب لاكتشاف المزيد",
            ...     reveal_items=[
            ...         {"label": "التعليم", "description": "أثرت التقنية على..."},
            ...         {"label": "الصحة", "description": "ساهمت التقنية في..."},
            ...     ]
            ... )
        """
        self.slide_count += 1
        slide = self._add_content_slide_with_layout()

        # Set hidden TOC title for Storyline sidebar menu
        self._set_slide_title_for_toc(slide, title)

        # --- Lecture title bar ---
        self._add_header_bar(slide, self.lecture_title)

        # --- Section banner ---
        self._add_section_banner(slide, title, wide=True)

        # --- Instruction text ---
        self._add_arabic_textbox(
            slide,
            left=Cm(2.5),
            top=Cm(4.5),
            width=Cm(29),
            height=Cm(1.5),
            text=instruction,
            font_name=FONT_MEDIUM,
            font_size=Pt(18),
            bold=False,
            color=BODY_TEXT,
            alignment=PP_ALIGN.RIGHT,
            word_wrap=True,
            auto_size=MSO_AUTO_SIZE.NONE,
            name="txt_instruction",
        )

        # --- Layout depends on item count ---
        tab_count = len(reveal_items)

        if tab_count > 4:
            # VERTICAL LIST layout for 5+ items (avoids tiny unreadable tabs)
            list_top_start = Cm(6.5)
            list_left = Cm(2.5)
            list_width = Cm(28.5)
            safe_bottom = 6300000  # Safe zone above page number

            # Adaptive spacing — shrinks rows to fit more items
            row_height, row_gap = self._calculate_adaptive_spacing(
                item_count=tab_count,
                available_top=list_top_start,
                available_bottom=safe_bottom,
                min_item_height=Cm(1.8),
            )
            row_step = row_height + row_gap

            for i, item in enumerate(reveal_items):
                reveal_num = i + 1
                row_top = int(list_top_start + i * row_step)

                # Alternating row backgrounds for visual distinction
                row_bg = CONTENT_CARD_BG if i % 2 == 0 else WHITE
                self._add_shape(
                    slide,
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left=list_left,
                    top=row_top,
                    width=list_width,
                    height=row_height,
                    fill_color=row_bg,
                    border_color=CONTENT_CARD_BORDER,
                    border_width=Pt(1),
                    name=f"bg_reveal_{reveal_num}",
                )

                # Number badge (left side in RTL = right visual side)
                badge_size = Cm(1.5)
                badge = self._add_shape(
                    slide,
                    MSO_SHAPE.OVAL,
                    left=list_left + list_width - badge_size - Cm(0.3),
                    top=row_top + (row_height - badge_size) // 2,
                    width=badge_size,
                    height=badge_size,
                    fill_color=PRIMARY_BLUE,
                    name=f"btn_reveal_{reveal_num}",
                )
                tf_b = badge.text_frame
                tf_b.word_wrap = False
                tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
                p_b = tf_b.paragraphs[0]
                p_b.alignment = PP_ALIGN.CENTER
                run_b = p_b.add_run()
                run_b.text = str(i + 1)
                self._set_run_font(run_b, FONT_EXTRABOLD, Pt(16), False, WHITE)

                # Label text — larger font
                self._add_arabic_textbox(
                    slide,
                    left=list_left + Cm(0.5),
                    top=row_top,
                    width=list_width - badge_size - Cm(1.5),
                    height=row_height,
                    text=item.get("label", ""),
                    font_name=FONT_EXTRABOLD,
                    font_size=Pt(20),
                    bold=False,
                    color=BODY_TEXT,
                    alignment=PP_ALIGN.RIGHT,
                    word_wrap=True,
                    auto_size=MSO_AUTO_SIZE.NONE,
                    name=f"txt_reveal_{reveal_num}",
                )

            # Put all detail/description text in speaker notes (Storyline handles reveal)
        else:
            # HORIZONTAL TABS layout for 4 or fewer items
            tab_area_left = Cm(2.5)
            tab_area_width = Cm(28.5)
            tab_top = Cm(7)
            tab_height = Cm(2.5)  # Taller tabs (was Cm(2))
            gap = Cm(0.5)

            total_gaps = gap * (tab_count - 1) if tab_count > 1 else 0
            tab_width = int((tab_area_width - total_gaps) / tab_count)

            for i, item in enumerate(reveal_items):
                tab_left = int(tab_area_left + i * (tab_width + gap))

                # First tab = active (darker), rest = inactive (lighter)
                tab_fill = PRIMARY_BLUE if i == 0 else PRIMARY_BLUE_LIGHT
                tab_shape = self._add_shape(
                    slide,
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left=tab_left,
                    top=tab_top,
                    width=tab_width,
                    height=tab_height,
                    fill_color=tab_fill,
                    name=f"btn_reveal_{i + 1}",
                    corner_radius=0.08,
                )
                tf = tab_shape.text_frame
                tf.word_wrap = True
                tf.margin_left = TEXT_MARGIN_LR
                tf.margin_right = TEXT_MARGIN_LR
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = item.get("label", "")
                self._set_run_font(run, FONT_EXTRABOLD, Pt(18), False, WHITE)
                self._set_rtl(p)

            # --- Description area (shown on click in Storyline) ---
            self._add_shape(
                slide,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left=Cm(2.5),
                top=Cm(10.5),
                width=Cm(28.5),
                height=Cm(5),
                fill_color=LIGHT_BLUE_BG,
                border_color=PRIMARY_BLUE,
                name="bg_reveal_desc",
            )

            # Left accent bar on description area
            self._add_shape(
                slide,
                MSO_SHAPE.RECTANGLE,
                left=Cm(30.5),
                top=Cm(10.5),
                width=Cm(0.4),
                height=Cm(5),
                fill_color=PRIMARY_BLUE,
                name="accent_desc_bar",
            )

            # Add first item's description as default visible text
            if reveal_items:
                self._add_arabic_textbox(
                    slide,
                    left=Cm(3),
                    top=Cm(11),
                    width=Cm(27.5),
                    height=Cm(4),
                    text=reveal_items[0].get("description", reveal_items[0].get("detail", "")),
                    font_name=FONT_REGULAR,
                    font_size=Pt(18),
                    bold=False,
                    color=BODY_TEXT,
                    alignment=PP_ALIGN.RIGHT,
                    word_wrap=True,
                    auto_size=MSO_AUTO_SIZE.NONE,
                    name="txt_reveal_desc",
                )

        # --- Structured notes for Storyline import ---
        btn_names = ", ".join(f"btn_reveal_{i+1}" for i in range(tab_count))
        layer_lines = "\n".join(
            f"Layer {i+1}: Show content when btn_reveal_{i+1} clicked — {item.get('label', '')}"
            for i, item in enumerate(reveal_items)
        )
        all_descriptions = "\n".join(
            f"btn_reveal_{i+1} ({item.get('label', '')}): {item.get('description', item.get('detail', ''))}"
            for i, item in enumerate(reveal_items)
        )
        structured_notes = (
            f"=== STORYLINE INSTRUCTIONS ===\n"
            f"Slide Type: Click to Reveal\n"
            f"Items: {btn_names}\n\n"
            f"=== LAYER INSTRUCTIONS ===\n"
            f"{layer_lines}\n\n"
            f"=== CONTENT ===\n"
            f"{all_descriptions}"
        )
        if notes:
            structured_notes = f"{notes}\n\n{structured_notes}"
        self._add_notes(slide, structured_notes)

    def add_dropdown_slide(
        self,
        title: str,
        instruction: str,
        items: list,
        notes: str = "",
    ):
        """
        Add a dropdown matching activity slide.

        Used for activities where learners select the correct category
        from a dropdown for each statement.

        Args:
            title: Activity title (e.g., "نشاط 1.3")
            instruction: Instruction text
            items: List of dicts with:
                   - "text": Statement text
                   - "correct": Correct dropdown selection
            notes: Speaker notes

        Visual output (ASCII mockup):
            +------------------------------------------+
            | [Title Bar: Lecture Name]                 |
            | [Activity Title]                         |
            | [Instructions]                           |
            |                                          |
            | [Statement 1]                 [Dropdown▼]|
            | [Statement 2]                 [Dropdown▼]|
            | [Statement 3]                 [Dropdown▼]|
            |                                          |
            | [9]                                       |
            +------------------------------------------+

        Example:
            >>> builder.add_dropdown_slide(
            ...     title="نشاط 1.3 (تأثير التقنية)",
            ...     instruction="اختر من القائمة المنسدلة الجانب المناسب",
            ...     items=[
            ...         {"text": "توفير بيئة محاكاة رقمية", "correct": "التعليم"},
            ...         {"text": "تحليل تفضيلات المستخدم", "correct": "الترفيه"},
            ...     ]
            ... )
        """
        self.slide_count += 1
        slide = self._add_content_slide_with_layout()

        # Set hidden TOC title for Storyline sidebar menu
        self._set_slide_title_for_toc(slide, title)

        # --- Lecture title bar ---
        self._add_header_bar(slide, self.lecture_title)

        # --- Activity banner ---
        self._add_section_banner(slide, title, wide=True)

        # --- Instruction text ---
        self._add_arabic_textbox(
            slide,
            left=Cm(2.5),
            top=Cm(4.5),
            width=Cm(29),
            height=Cm(2),
            text=instruction,
            font_name=FONT_MEDIUM,
            font_size=Pt(18),
            bold=False,
            color=BODY_TEXT,
            alignment=PP_ALIGN.RIGHT,
            word_wrap=True,
            auto_size=MSO_AUTO_SIZE.NONE,
            name="txt_instruction",
        )

        # --- Statement rows with dropdown indicators ---
        row_top_start = Cm(7)
        safe_bottom = 6300000  # Safe zone above page number

        # Adaptive spacing — shrinks rows to fit more items
        dd_row_height, dd_row_gap = self._calculate_adaptive_spacing(
            item_count=len(items),
            available_top=row_top_start,
            available_bottom=safe_bottom,
            min_item_height=Cm(1.8),
        )
        dd_row_step = dd_row_height + dd_row_gap

        for i, item_data in enumerate(items):
            dd_num = i + 1
            row_top = int(row_top_start + i * dd_row_step)

            # Statement text — Pt(18) for QM compliance
            self._add_arabic_textbox(
                slide,
                left=Cm(7),
                top=row_top,
                width=Cm(22),
                height=dd_row_height,
                text=item_data.get("text", ""),
                font_name=FONT_REGULAR,
                font_size=Pt(18),
                bold=False,
                color=BODY_TEXT,
                alignment=PP_ALIGN.RIGHT,
                word_wrap=True,
                auto_size=MSO_AUTO_SIZE.NONE,
                name=f"txt_statement_{dd_num}",
            )

            # Dropdown indicator (simulated with a shape)
            dropdown = self._add_shape(
                slide,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left=Cm(2.5),
                top=row_top,
                width=Cm(4),
                height=dd_row_height,
                fill_color=WHITE,
                border_color=PRIMARY_BLUE,
                name=f"btn_dropdown_{dd_num}",
            )
            tf = dropdown.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = "▼"
            self._set_run_font(run, FONT_REGULAR, Pt(16), False, PRIMARY_BLUE)

        # --- Notes with correct answers ---
        correct_text = "\n".join(
            f"{item.get('text', '')[:40]}... → {item.get('correct', '')}"
            for item in items
        )
        notes_text = notes if notes else f"الإجابة الصحيحة:\n{correct_text}"
        self._add_notes(slide, notes_text)

    def finalize(self):
        """
        Set up cross-slide references after all slides are added.

        Call this BEFORE save() to:
        - Link btn_start on the title slide to slide 2
        - Any other cross-slide click actions

        Example:
            >>> builder.add_title_slide(...)
            >>> builder.add_objectives_slide(...)
            >>> builder.finalize()  # Sets up cross-slide links
            >>> builder.save("output.pptx")
        """
        if len(self.prs.slides) > 1:
            title_slide = self.prs.slides[0]
            for shape in title_slide.shapes:
                if hasattr(shape, 'name') and shape.name == "btn_start":
                    shape.click_action.target_slide = self.prs.slides[1]
                    break

    def save(self, filepath: str):
        """
        Save the presentation to a file.

        Automatically calls finalize() to set up cross-slide links,
        then creates any necessary directories and writes the .pptx file.

        Args:
            filepath: Output file path (e.g., "output/DSAI/U01/lecture.pptx")

        Example:
            >>> builder.save("output/DSAI/U01/DSAI_U01_Interactive_Lecture.pptx")
        """
        # Set up cross-slide references
        self.finalize()

        # Create output directory if it doesn't exist
        output_dir = os.path.dirname(filepath)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        self.prs.save(filepath)

    # -----------------------------------------------------------------------
    # PRIVATE HELPER METHODS
    # -----------------------------------------------------------------------

    def _calculate_adaptive_spacing(self, item_count, available_top, available_bottom, min_item_height):
        """
        Calculate item height and gap to fit N items between top and bottom bounds.

        This prevents content from overflowing past the page number area.
        If items don't fit at their preferred height, they shrink to fit.

        Args:
            item_count: Number of items to fit
            available_top: Top of available area (EMU)
            available_bottom: Bottom of available area (EMU)
            min_item_height: Preferred/minimum height per item (EMU)

        Returns:
            (item_height, gap) — both in EMU
        """
        if item_count <= 0:
            return min_item_height, 0

        available = available_bottom - available_top
        min_gap = Cm(0.2)  # Minimum gap between items

        # Total space needed at preferred height
        total_needed = item_count * min_item_height + max(item_count - 1, 0) * min_gap
        if total_needed <= available:
            # Everything fits — distribute extra space as gaps
            item_height = min_item_height
            remaining = available - (item_count * item_height)
            gap = remaining // max(item_count - 1, 1)
        else:
            # Too many items — shrink height to fit
            gap = min_gap
            total_gaps = gap * max(item_count - 1, 0)
            item_height = (available - total_gaps) // item_count

        return item_height, gap

    def _add_slide_with_layout(self, layout_index):
        """
        Add a slide using a specific layout from the template.

        Layout 0 = "Title Slide" (has full background image, logo, line)
        Layout 1 = "Title and Content" (has bg, header bar, footer bar, logo)

        CRITICAL: Do NOT modify any placeholder shapes on the returned slide.
        Always use slide.shapes.add_textbox() for text content.

        Args:
            layout_index: Index of the layout to use (0 or 1)

        Returns:
            The new slide object.
        """
        # Use the template's layouts (they contain all background elements)
        try:
            slide_layout = self.prs.slide_layouts[layout_index]
        except IndexError:
            # Fallback to layout 0 if index out of range
            slide_layout = self.prs.slide_layouts[0]
        return self.prs.slides.add_slide(slide_layout)

    def _add_content_slide_with_layout(self):
        """
        Add a content slide using Layout 1 ("Title and Content").

        This layout provides:
        - Background image (nearly white/subtle texture)
        - Header bar (blue gradient rounded bar PNG)
        - Footer bar (blue strip at bottom)
        - University logo (top-left corner)

        Returns:
            The new slide object with all template visuals.
        """
        return self._add_slide_with_layout(1)

    def _add_header_bar(self, slide, title: str, subtitle: str = "", color=None):
        """
        Add the lecture title bar at the top of a slide.

        This bar appears on slides 2+ and shows the lecture name
        in centered ExtraBold text.

        Args:
            slide: The slide object
            title: Text to display in the bar
            subtitle: Optional subtitle text (not commonly used)
            color: Text color override (defaults to BODY_TEXT)

        Visual output:
            +------------------------------------------+
            |       [  Lecture Title Text  ]            |
            +------------------------------------------+
            The bar is 14.39cm wide, centered horizontally on the slide.
        """
        text_color = color if color else BODY_TEXT

        self._add_arabic_textbox(
            slide,
            left=TITLE_BAR_LEFT,
            top=TITLE_BAR_TOP,
            width=TITLE_BAR_WIDTH,
            height=TITLE_BAR_HEIGHT,
            text=title,
            font_name=FONT_EXTRABOLD,
            font_size=Pt(18),
            bold=False,
            color=text_color,
            alignment=PP_ALIGN.CENTER,
            name="header_title",
        )

    def _add_section_banner(self, slide, title: str, wide: bool = False):
        """
        Add a section banner using PNG image with dark text.

        Uses the actual template PNG images (banner_narrow.png or banner_wide.png)
        instead of colored rectangles. Text color is #333333 (dark) on the
        light grey/blue PNG background.

        Args:
            slide: The slide object
            title: Banner title text
            wide: If True, uses the wider banner (for activities/summary)

        Visual output:
            +------------------------------------------+
            |        [======= Title =======]           |
            +------------------------------------------+
        """
        if wide:
            banner_left = WIDE_BANNER_LEFT
            banner_top = WIDE_BANNER_TOP
            banner_width = WIDE_BANNER_WIDTH
            banner_height = WIDE_BANNER_HEIGHT
            text_left = WIDE_BANNER_TEXT_LEFT
            text_top = WIDE_BANNER_TEXT_TOP
            text_width = WIDE_BANNER_TEXT_WIDTH
            text_height = WIDE_BANNER_TEXT_HEIGHT
            asset_name = ASSET_BANNER_WIDE
            font_size = Pt(20)
        else:
            banner_left = BANNER_LEFT
            banner_top = BANNER_TOP
            banner_width = BANNER_WIDTH
            banner_height = BANNER_HEIGHT
            text_left = NARROW_BANNER_TEXT_LEFT
            text_top = NARROW_BANNER_TEXT_TOP
            text_width = NARROW_BANNER_TEXT_WIDTH
            text_height = NARROW_BANNER_TEXT_HEIGHT
            asset_name = ASSET_BANNER_NARROW
            font_size = Pt(18)

        # Banner background — PNG image from the template
        banner_path = os.path.join(self.assets_dir, asset_name)
        if os.path.exists(banner_path):
            pic = slide.shapes.add_picture(
                banner_path,
                banner_left,
                banner_top,
                banner_width,
                banner_height,
            )
            pic.name = "header_banner"
        else:
            # Fallback: colored rectangle if PNG not found
            self._add_shape(
                slide,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left=banner_left,
                top=banner_top,
                width=banner_width,
                height=banner_height,
                fill_color=PRIMARY_BLUE,
                name="header_banner",
            )

        # Banner title text — dark color #333333 on light PNG bg
        self._add_arabic_textbox(
            slide,
            left=text_left,
            top=text_top,
            width=text_width,
            height=text_height,
            text=title,
            font_name=FONT_EXTRABOLD,
            font_size=font_size,
            bold=False,
            color=BODY_TEXT,     # #333333 — dark text on light banner
            alignment=PP_ALIGN.CENTER,
            name="header_banner_text",
        )

    def _add_arabic_textbox(
        self,
        slide,
        left: int,
        top: int,
        width: int,
        height: int,
        text: str,
        font_name: str = FONT_REGULAR,
        font_size=Pt(16),
        bold: bool = False,
        color: RGBColor = BODY_TEXT,
        alignment=PP_ALIGN.RIGHT,
        word_wrap: bool = True,
        auto_size=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
        line_spacing: float = None,
        name: str = None,
    ):
        """
        Add a text box with Arabic RTL text to a slide.

        This is the core text-rendering helper. Every text element in the
        presentation goes through this method, which ensures:
        1. Correct RTL paragraph direction
        2. Arabic language tag (ar-JO)
        3. Proper font assignment (cs_font, latin_font, ea_font)
        4. Consistent margins and sizing

        Args:
            slide: The slide object to add the textbox to
            left: Left position in EMU
            top: Top position in EMU
            width: Width in EMU
            height: Height in EMU
            text: The Arabic text to display
            font_name: Font name (default: Tajawal)
            font_size: Font size in Pt (default: 16pt)
            bold: Whether to make the text bold
            color: Text color as RGBColor
            alignment: Paragraph alignment (RIGHT for RTL body text)
            word_wrap: Whether to wrap text
            auto_size: Auto-size behavior
            line_spacing: Line spacing multiplier (e.g., 1.5)

        Returns:
            The created textbox shape.
        """
        txBox = slide.shapes.add_textbox(left, top, width, height)
        if name:
            txBox.name = name
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        tf.auto_size = auto_size
        tf.margin_left = TEXT_MARGIN_LR
        tf.margin_right = TEXT_MARGIN_LR
        tf.margin_top = TEXT_MARGIN_TB
        tf.margin_bottom = TEXT_MARGIN_TB

        # Use the first (default) paragraph
        p = tf.paragraphs[0]
        p.alignment = alignment

        # Set line spacing — default 1.3 for body text (>= 18pt) for Arabic readability
        if line_spacing:
            p.line_spacing = line_spacing
        elif font_size >= Pt(18):
            p.line_spacing = 1.3

        # Add the text run
        run = p.add_run()
        run.text = text

        # Apply font settings
        self._set_run_font(run, font_name, font_size, bold, color)

        # Set RTL direction — critical for Arabic text
        self._set_rtl(p)

        return txBox

    def _set_run_font(self, run, font_name: str, font_size, bold: bool, color: RGBColor):
        """
        Apply font settings to a text run.

        Sets size, bold, and color using the python-pptx API, then delegates
        to the shared rtl_helpers.pptx_set_run_font_arabic() for the font
        name assignment. This ensures the font is set on all three slots
        (cs, latin, ea) via XML for reliable Arabic rendering.

        Args:
            run: The text run to style
            font_name: Font family name (e.g., "Tajawal ExtraBold")
            font_size: Font size (Pt value)
            bold: Whether to bold the text
            color: Text color as RGBColor
        """
        font = run.font
        font.size = font_size
        font.bold = bold
        font.color.rgb = color

        # Delegate font name + language to the shared RTL helper.
        # This sets cs, latin, ea fonts and the ar-JO language tag via XML.
        pptx_set_run_font_arabic(run, font_name)

    def _set_rtl(self, paragraph):
        """
        Set paragraph direction to RTL for Arabic text.

        Delegates to the shared rtl_helpers.pptx_set_paragraph_rtl().

        Args:
            paragraph: The paragraph object to set RTL on
        """
        pptx_set_paragraph_rtl(paragraph)

    def _validate_bounds(self, left, top, width, height, context=""):
        """
        Warn if a shape would extend beyond slide boundaries.

        Prints a console warning when shapes overflow — helps catch
        layout bugs during development without crashing production.

        Args:
            left: Left position in EMU
            top: Top position in EMU
            width: Width in EMU
            height: Height in EMU
            context: Description of the shape (e.g., "bg_col1_card")
        """
        right_edge = left + width
        bottom_edge = top + height
        if right_edge > SLIDE_WIDTH:
            overflow_cm = (right_edge - SLIDE_WIDTH) / 360000
            print(f"⚠ OVERFLOW: {context} extends {overflow_cm:.1f}cm beyond right edge")
        if bottom_edge > SLIDE_HEIGHT:
            overflow_cm = (bottom_edge - SLIDE_HEIGHT) / 360000
            print(f"⚠ OVERFLOW: {context} extends {overflow_cm:.1f}cm beyond bottom edge")

    def _add_shape(
        self,
        slide,
        shape_type,
        left: int,
        top: int,
        width: int,
        height: int,
        fill_color: RGBColor = None,
        border_color: RGBColor = None,
        border_width=None,
        name: str = None,
        corner_radius: float = None,
    ):
        """
        Add a shape to a slide with optional fill and border.

        Used for rectangles, rounded rectangles, ovals, etc. that make up
        the visual structure of slides (banners, cards, buttons, etc.)

        Args:
            slide: The slide object
            shape_type: MSO_SHAPE enum value (e.g., MSO_SHAPE.RECTANGLE)
            left: Left position in EMU
            top: Top position in EMU
            width: Width in EMU
            height: Height in EMU
            fill_color: Optional solid fill color
            border_color: Optional border color
            border_width: Optional border width (Pt value)
            name: Optional shape name (for Storyline identification)
            corner_radius: Optional corner radius for rounded rectangles (0.0 to 1.0)

        Returns:
            The created shape object.
        """
        # Check for overflow before creating the shape
        self._validate_bounds(left, top, width, height, name or "unnamed_shape")

        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        if name:
            shape.name = name

        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color

        if border_color:
            shape.line.color.rgb = border_color
            if border_width:
                shape.line.width = border_width
        else:
            # No border — set to no line
            shape.line.fill.background()

        # Set custom corner radius for rounded rectangles
        if corner_radius is not None and shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
            # adjustments[0] controls corner radius (0.0 to 1.0)
            shape.adjustments[0] = corner_radius

        return shape

    def _add_shadow_to_shape(self, shape, blur_pt=6, dist_pt=3, direction=2700000, opacity_pct=25):
        """
        Add a real OOXML outer shadow effect to a shape.

        Uses effectLst > outerShdw for professional drop shadows that
        Storyline 360 respects (standard OOXML effects).

        Args:
            shape: The shape to add shadow to
            blur_pt: Shadow blur radius in points (default: 6)
            dist_pt: Shadow distance in points (default: 3)
            direction: Shadow direction in 60000ths of degree (default: 2700000 = bottom-right)
            opacity_pct: Shadow opacity 0-100 (default: 25)
        """
        from lxml import etree

        nsmap = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        }

        # Build shadow XML
        blur_emu = blur_pt * 12700  # Points to EMU
        dist_emu = dist_pt * 12700
        alpha_val = (100 - opacity_pct) * 1000  # Convert to OOXML alpha (0=opaque, 100000=transparent)

        shadow_xml = (
            f'<a:outerShdw xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
            f' blurRad="{blur_emu}" dist="{dist_emu}" dir="{direction}" rotWithShape="0">'
            f'<a:srgbClr val="000000"><a:alpha val="{alpha_val}"/></a:srgbClr>'
            f'</a:outerShdw>'
        )
        shadow_elem = etree.fromstring(shadow_xml)

        # Find or create effectLst on the shape's spPr
        spPr = shape._element.spPr
        effectLst = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst')
        if effectLst is None:
            effectLst = etree.SubElement(spPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst')

        effectLst.append(shadow_elem)

    def _add_decorative_corner(self, slide, position="top_right",
                               color=None, size=None):
        """
        Draw a decorative corner accent using shapes instead of low-res PNGs.

        Creates 2 thin lines (horizontal + vertical) meeting at the corner,
        plus a small arc for elegance. Scales perfectly at any resolution.

        Args:
            slide: The slide object
            position: "top_right" or "bottom_left"
            color: Line color (default: WHITE)
            size: Overall size of the corner decoration (default: Cm(4))
        """
        line_color = color if color else WHITE
        corner_size = size if size else Cm(4)
        line_thickness = Pt(2)

        if position == "top_right":
            # Horizontal line extending left from top-right corner area
            self._add_shape(
                slide, MSO_SHAPE.RECTANGLE,
                left=SLIDE_WIDTH - Cm(3) - corner_size,
                top=Cm(2.5),
                width=corner_size,
                height=line_thickness,
                fill_color=line_color,
                name="deco_corner_tr_h",
            )
            # Vertical line extending down from top-right corner area
            self._add_shape(
                slide, MSO_SHAPE.RECTANGLE,
                left=SLIDE_WIDTH - Cm(3),
                top=Cm(2.5),
                width=line_thickness,
                height=corner_size,
                fill_color=line_color,
                name="deco_corner_tr_v",
            )
            # Small circle at the corner junction for elegance
            dot_size = Cm(0.4)
            self._add_shape(
                slide, MSO_SHAPE.OVAL,
                left=SLIDE_WIDTH - Cm(3) - dot_size // 2,
                top=Cm(2.5) - dot_size // 2,
                width=dot_size,
                height=dot_size,
                fill_color=line_color,
                name="deco_corner_tr_dot",
            )
        elif position == "bottom_left":
            # Horizontal line extending right from bottom-left corner area
            self._add_shape(
                slide, MSO_SHAPE.RECTANGLE,
                left=Cm(3),
                top=SLIDE_HEIGHT - Cm(2.5),
                width=corner_size,
                height=line_thickness,
                fill_color=line_color,
                name="deco_corner_bl_h",
            )
            # Vertical line extending up from bottom-left corner area
            self._add_shape(
                slide, MSO_SHAPE.RECTANGLE,
                left=Cm(3),
                top=SLIDE_HEIGHT - Cm(2.5) - corner_size,
                width=line_thickness,
                height=corner_size,
                fill_color=line_color,
                name="deco_corner_bl_v",
            )
            # Small circle at the corner junction
            dot_size = Cm(0.4)
            self._add_shape(
                slide, MSO_SHAPE.OVAL,
                left=Cm(3) - dot_size // 2,
                top=SLIDE_HEIGHT - Cm(2.5) - dot_size // 2,
                width=dot_size,
                height=dot_size,
                fill_color=line_color,
                name="deco_corner_bl_dot",
            )

    # ------------------------------------------------------------------
    # IMAGE GENERATION HELPER — Auto-generate images from prompts
    # ------------------------------------------------------------------

    def _generate_image_for_slide(self, image_prompt, image_type, topic_key=None):
        """
        Generate an image using the project's visual direction.

        Called by slide methods when an image_prompt is provided but no
        image_path. Uses the project's config.json visual direction
        (prefix, suffix, negative rules) to build the final prompt.

        Args:
            image_prompt: Description of the image to generate
            image_type: One of "content", "card", "section", "two_column",
                       "closing", "quiz" — used for aspect ratio lookup
            topic_key: Optional cache key (e.g. "design_thinking").
                      If an image already exists for this key, returns
                      the cached path instead of regenerating.

        Returns:
            Absolute path to the generated image file, or None if
            generation failed or no project_code is set.
        """
        if not image_prompt or not self.project_code:
            return None
        try:
            result = generate_storyboard_image(
                prompt=image_prompt,
                project_code=self.project_code,
                unit_number=self.unit_number or 1,
                image_type=image_type,
                topic_key=topic_key,
            )
            if result["success"]:
                return result["path"]
        except Exception:
            # Graceful fallback — image generation is optional
            pass
        return None

    # ------------------------------------------------------------------
    # IMAGE HELPERS — Smart image placement with aspect ratio preservation
    # ------------------------------------------------------------------

    def _get_image_dimensions(self, image_path, max_width, max_height):
        """
        Calculate display dimensions that fit inside a bounding box
        while preserving the image's natural aspect ratio.

        Uses min(scale_w, scale_h) — so the image NEVER stretches or
        overflows the bounding box. A portrait image stays tall,
        a landscape image stays wide.

        Args:
            image_path: Path to the image file (PNG, JPG, etc.)
            max_width: Maximum allowed width in EMU
            max_height: Maximum allowed height in EMU

        Returns:
            Tuple of (display_width, display_height) in EMU,
            or None if the file can't be read.
        """
        try:
            with PILImage.open(image_path) as img:
                img_w, img_h = img.size  # pixels

            # Calculate scale factors for width and height
            scale_w = max_width / img_w
            scale_h = max_height / img_h

            # Use the SMALLER scale — this ensures the image fits
            # entirely within the box without overflowing either dimension
            scale = min(scale_w, scale_h)

            display_w = int(img_w * scale)
            display_h = int(img_h * scale)
            return (display_w, display_h)
        except Exception:
            return None

    def _add_image(self, slide, image_path, left, top, max_width, max_height,
                   name=None, center_in_area=True):
        """
        Smart image inserter with aspect ratio preservation and centering.

        This is the main method agents use to add images to slides.
        It handles:
        - Aspect ratio preservation (never stretches)
        - Centering within bounding box (so images look balanced)
        - Missing file guard (returns None instead of crashing)
        - Shape naming for Storyline selection pane

        Args:
            slide: The slide object to add the image to
            image_path: Path to the image file (PNG, JPG, etc.)
            left: Left edge of bounding box (EMU)
            top: Top edge of bounding box (EMU)
            max_width: Maximum width of bounding box (EMU)
            max_height: Maximum height of bounding box (EMU)
            name: Shape name for Storyline (e.g., "img_content")
            center_in_area: If True, center the image within the box

        Returns:
            The picture shape object, or None if image file is missing.
        """
        # Guard: don't crash if the file doesn't exist
        if not image_path or not os.path.exists(image_path):
            return None

        # Calculate display size that fits within the bounding box
        dims = self._get_image_dimensions(image_path, max_width, max_height)
        if dims is None:
            return None

        display_w, display_h = dims

        # Center the image within the bounding box
        if center_in_area:
            # Offset to center horizontally and vertically
            offset_left = (max_width - display_w) // 2
            offset_top = (max_height - display_h) // 2
            img_left = left + offset_left
            img_top = top + offset_top
        else:
            img_left = left
            img_top = top

        # Check for overflow before adding
        self._validate_bounds(img_left, img_top, display_w, display_h,
                              name or "unnamed_image")

        # Add the picture to the slide
        pic = slide.shapes.add_picture(
            image_path,
            left=img_left,
            top=img_top,
            width=display_w,
            height=display_h,
        )

        # Name the shape for Storyline's Selection Pane
        if name:
            pic.name = name

        return pic

    def _add_accent_stripe(self, slide, color=None):
        """
        Add a vertical accent stripe on the right side of the slide.

        Used in content slide Variant B for visual variety.
        The stripe provides a colorful accent that breaks the monotony
        of full-width content cards.

        Args:
            slide: The slide object
            color: Stripe color (default: PRIMARY_BLUE_LIGHT)
        """
        stripe_color = color if color else PRIMARY_BLUE_LIGHT
        self._add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            left=Cm(31),
            top=Cm(4),
            width=Cm(1.2),
            height=Cm(13),
            fill_color=stripe_color,
            name="bg_accent_stripe",
        )

    def _add_numbered_points(self, slide, items, start_top=Cm(5.5), left=Cm(3), width=Cm(28)):
        """
        Add content as numbered points with circle badges instead of bullets.

        Used in content slide Variant C for visual variety.
        Each point gets a numbered circle badge (RTL: badge on right side).

        Args:
            slide: The slide object
            items: List of point strings
            start_top: Starting Y position
            left: Starting X position
            width: Width of content area
        """
        point_spacing = Cm(2.5)

        for i, item_text in enumerate(items):
            point_top = int(start_top + i * point_spacing)
            point_num = i + 1

            # Number badge (circle on the right for RTL)
            badge_size = Cm(1.8)
            badge = self._add_shape(
                slide,
                MSO_SHAPE.OVAL,
                left=left + width - badge_size,
                top=point_top,
                width=badge_size,
                height=badge_size,
                fill_color=PRIMARY_BLUE,
                name=f"num_point_{point_num}",
            )
            tf = badge.text_frame
            tf.word_wrap = False
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(point_num)
            self._set_run_font(run, FONT_EXTRABOLD, Pt(18), False, WHITE)

            # Point text (to the left of badge for RTL)
            self._add_arabic_textbox(
                slide,
                left=left,
                top=point_top,
                width=width - badge_size - Cm(0.5),
                height=badge_size,
                text=item_text,
                font_name=FONT_REGULAR,
                font_size=Pt(20),
                bold=False,
                color=BODY_TEXT,
                alignment=PP_ALIGN.RIGHT,
                word_wrap=True,
                auto_size=MSO_AUTO_SIZE.NONE,
                name=f"txt_point_{point_num}",
            )

    def _add_footer(self, slide):
        """
        Add a footer to the slide (placeholder for future use).

        The template doesn't have a visible footer on content slides,
        but this method is here for completeness and can be extended
        if needed for specific project requirements.

        Args:
            slide: The slide object
        """
        pass

    def _add_notes(self, slide, notes_text: str):
        """
        Add speaker notes to a slide.

        Speaker notes are used to store Storyline instructions,
        correct answers, image links, and other metadata that
        shouldn't be visible on the slide itself.

        Args:
            slide: The slide object
            notes_text: The notes content
        """
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = notes_text

    def _set_slide_title_for_toc(self, slide, title_text: str):
        """
        Set a hidden title for Storyline TOC (Table of Contents).

        Storyline 360 reads the slide title placeholder to populate its
        sidebar menu. Without a title, slides appear as "Untitled Slide".

        This adds an off-screen textbox named "title" that Storyline
        reads but learners never see.

        Args:
            slide: The slide object
            title_text: Arabic title text for the TOC entry
        """
        # Place far off-screen so it's invisible in the presentation
        txBox = slide.shapes.add_textbox(
            left=-Cm(20),
            top=-Cm(20),
            width=Cm(10),
            height=Cm(2),
        )
        txBox.name = "title"
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        self._set_run_font(run, FONT_REGULAR, Pt(12), False, WHITE)
        self._set_rtl(p)

    def _build_import_instructions(self, lecture_title: str):
        """
        Build Storyline 360 import instructions for the developer.

        Returns a formatted string to be placed in the title slide's
        speaker notes, giving the Storyline developer everything they
        need to import and configure the presentation.
        """
        now = datetime.now().strftime("%Y-%m-%d")
        return (
            "=== STORYLINE 360 IMPORT INSTRUCTIONS ===\n\n"
            f"Project: {self.project_code}\n"
            f"Unit: {self.unit_number}\n"
            f"Lecture: {lecture_title}\n"
            f"Generated: {now}\n"
            f"Designer: {self.designer}\n\n"
            "--- REQUIRED FONTS ---\n"
            "Install BEFORE importing:\n"
            "  1. Tajawal ExtraBold\n"
            "  2. Tajawal Medium\n"
            "  3. Tajawal (Regular)\n"
            "Download: https://fonts.google.com/specimen/Tajawal\n\n"
            "--- IMPORT STEPS ---\n"
            "1. File > Import > PowerPoint\n"
            "2. Select THIS .pptx file\n"
            "3. Import ALL slides\n"
            "4. Story Size: 1280 x 720 pixels (must match)\n\n"
            "--- POST-IMPORT QA ---\n"
            "[ ] Verify TOC shows Arabic slide titles (not 'Untitled')\n"
            "[ ] Verify fonts render correctly (Tajawal family)\n"
            "[ ] Check RTL text direction on all slides\n"
            "[ ] Test all interactive elements (buttons, quiz, drag-drop)\n"
            "[ ] Verify speaker notes contain Storyline instructions per slide\n\n"
            "--- SHAPE NAMING CONVENTION ---\n"
            "txt_*  = Text elements\n"
            "btn_*  = Clickable buttons (add triggers)\n"
            "icon_* = Decorative icons\n"
            "bg_*   = Background shapes\n"
            "opt_*  = Quiz option shapes (add triggers)\n"
            "num_*  = Numbered elements\n"
            "title  = Hidden TOC title (do not move)\n"
        )

    def _add_bullet_list(
        self,
        slide,
        left: int,
        top: int,
        width: int,
        height: int,
        items: list,
        font_size=Pt(16),
        color: RGBColor = None,
        name: str = None,
    ):
        """
        Add a bullet list as a text box with multiple paragraphs.

        Each item becomes a separate paragraph. RTL direction is set
        on every paragraph for consistent Arabic rendering.

        Args:
            slide: The slide object
            left: Left position in EMU
            top: Top position in EMU
            width: Width in EMU
            height: Height in EMU
            items: List of bullet point strings
            font_size: Font size for bullet text
            color: Text color (default: BODY_TEXT)
            name: Optional shape name

        Returns:
            The created textbox shape.
        """
        text_color = color if color else BODY_TEXT

        txBox = slide.shapes.add_textbox(left, top, width, height)
        if name:
            txBox.name = name
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = TEXT_MARGIN_LR
        tf.margin_right = TEXT_MARGIN_LR
        tf.margin_top = TEXT_MARGIN_TB
        tf.margin_bottom = TEXT_MARGIN_TB

        for i, item_text in enumerate(items):
            # Use the existing first paragraph for the first item,
            # add new paragraphs for subsequent items
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.alignment = PP_ALIGN.RIGHT

            # Colored filled circle marker (BULLET_MARKER_COLOR) + body text
            bullet_run = p.add_run()
            bullet_run.text = "\u25CF "  # Filled circle character
            self._set_run_font(bullet_run, FONT_REGULAR, Pt(16), False, BULLET_MARKER_COLOR)

            text_run = p.add_run()
            text_run.text = item_text
            self._set_run_font(text_run, FONT_REGULAR, font_size, False, text_color)
            self._set_rtl(p)

            # Comfortable spacing for Arabic bullet items
            p.space_before = Pt(10)
            p.space_after = Pt(10)
            p.line_spacing = 1.4

        return txBox
