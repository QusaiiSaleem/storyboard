#!/usr/bin/env python3
"""
NJR01 U02 Interactive Lecture Generator
Generates a 28-slide interactive lecture from the Arabic template.
Uses python-pptx directly - no external dependencies on skill scripts.

Usage: python3 generate_lecture.py
"""
import copy, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── Paths ──
TEMPLATE = "/Users/qusaiabushanap/dev/storyboard/templates/قالب المحاضرة التفاعلية- عربي.pptx"
OUTPUT = "/Users/qusaiabushanap/dev/storyboard/output/NJR01/U02/NJR01_U02_Interactive_Lecture.pptx"
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)

# ── Step 1: Analyze template ──
print("Loading template...")
prs = Presentation(TEMPLATE)
slide_count = len(prs.slides)
print(f"Template has {slide_count} slides")
print(f"Slide width: {prs.slide_width}, height: {prs.slide_height}")

# Print template slide inventory
for i, slide in enumerate(prs.slides):
    layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
    shapes_with_text = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text[:50].strip()
            if txt:
                shapes_with_text.append(txt)
    print(f"  Slide {i}: layout='{layout_name}', text shapes: {len(shapes_with_text)}")
    for t in shapes_with_text[:3]:
        print(f"    -> {t}")

print("\n" + "="*60)
print("Template analysis complete. Now building presentation...")
print("="*60 + "\n")


# ── Step 2: Identify slide layouts ──
# We need to understand which template slides to duplicate for our 28-slide deck.
# We'll map each of our target slides to a template slide index.

# First, let's catalog all available layouts
layouts_available = {}
for i, layout in enumerate(prs.slide_layouts):
    layouts_available[i] = layout.name
    print(f"  Layout {i}: '{layout.name}'")

print(f"\nTotal layouts: {len(layouts_available)}")


# ── Step 3: Build the 28-slide deck by duplicating template slides ──
# Strategy: We'll use the template slides as-is and duplicate them.
# Based on typical Arabic lecture templates:
#   Slide 0 = Title/Cover slide
#   Slide 1 = Content slide with title + body (objectives, content, etc.)
#   Other slides = Various layouts

# We'll duplicate slides from the template to create our 28 slides.
# The key insight: we duplicate the XML of each slide.

def duplicate_slide(prs, template_slide_index):
    """Duplicate a slide from the presentation by copying its XML."""
    template_slide = prs.slides[template_slide_index]
    slide_layout = template_slide.slide_layout

    # Add a new slide with the same layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Copy all shapes from template to new slide
    # First, remove default placeholder shapes from new slide
    # (they come from the layout)

    # We need to copy the XML content
    # Get the template slide's XML
    import lxml.etree as etree

    # Clear the new slide's spTree (shape tree)
    new_sp_tree = new_slide.shapes._spTree
    # Remove all existing shapes except the first two (nvGrpSpPr and grpSpPr)
    children_to_remove = []
    for child in new_sp_tree:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag not in ('nvGrpSpPr', 'grpSpPr'):
            children_to_remove.append(child)
    for child in children_to_remove:
        new_sp_tree.remove(child)

    # Copy shapes from template
    template_sp_tree = template_slide.shapes._spTree
    for child in template_sp_tree:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag not in ('nvGrpSpPr', 'grpSpPr'):
            new_sp_tree.append(copy.deepcopy(child))

    # Copy slide background if present
    template_cSld = template_slide._element
    new_cSld = new_slide._element

    # Copy background
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    template_bg = template_cSld.find('.//p:bg', nsmap)
    if template_bg is not None:
        new_bg = new_cSld.find('.//p:bg', nsmap)
        if new_bg is not None:
            new_bg.getparent().replace(new_bg, copy.deepcopy(template_bg))
        else:
            # Insert bg before spTree
            cSld = new_cSld.find('.//p:cSld', nsmap)
            if cSld is None:
                cSld = new_cSld
            sp_tree = cSld.find('.//p:spTree', nsmap)
            if sp_tree is not None:
                cSld.insert(list(cSld).index(sp_tree), copy.deepcopy(template_bg))

    return new_slide


def set_text_in_shape(shape, paragraphs_data):
    """Set text in a shape with proper RTL Arabic formatting.

    paragraphs_data: list of dicts with keys:
        - text: str
        - bold: bool (optional)
        - font_size: int in Pt (optional)
        - color: str hex like "FFFFFF" (optional)
        - alignment: PP_ALIGN value (optional)
        - bullet: bool (optional)
        - level: int (optional, for bullets)
    """
    if not shape.has_text_frame:
        return

    tf = shape.text_frame
    tf.word_wrap = True

    # Clear existing paragraphs
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)

    for idx, pdata in enumerate(paragraphs_data):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Set paragraph properties
        p.alignment = pdata.get('alignment', PP_ALIGN.RIGHT)  # RTL default

        # Set bullet level
        if pdata.get('bullet'):
            p.level = pdata.get('level', 0)

        # Add run with text
        run = p.add_run()
        run.text = pdata.get('text', '')

        # Font properties
        font = run.font
        if 'font_size' in pdata:
            font.size = Pt(pdata['font_size'])
        if 'bold' in pdata:
            font.bold = pdata['bold']
        if 'italic' in pdata:
            font.italic = pdata['italic']
        if 'color' in pdata:
            font.color.rgb = RGBColor.from_string(pdata['color'])
        if 'font_name' in pdata:
            font.name = pdata['font_name']

        # RTL for Arabic
        from lxml import etree
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        pPr = p._p.find('{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        if pPr is None:
            pPr = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
            # Move pPr to be first child
            p._p.insert(0, pPr)
        pPr.set('rtl', '1')

        # Set space before/after if specified
        if 'space_before' in pdata:
            spc_bef = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcBef')
            spc_pts = etree.SubElement(spc_bef, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcPts')
            spc_pts.set('val', str(int(pdata['space_before'] * 100)))

        if 'space_after' in pdata:
            spc_aft = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcAft')
            spc_pts = etree.SubElement(spc_aft, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcPts')
            spc_pts.set('val', str(int(pdata['space_after'] * 100)))


def find_shape_by_placeholder_type(slide, ph_type):
    """Find a shape by placeholder type (e.g., 'TITLE', 'BODY')."""
    for shape in slide.placeholders:
        if ph_type == 'TITLE' and shape.placeholder_format.idx == 0:
            return shape
        if ph_type == 'BODY' and shape.placeholder_format.idx == 1:
            return shape
        if ph_type == 'SUBTITLE' and shape.placeholder_format.idx == 1:
            return shape
    return None


def find_title_shape(slide):
    """Find the title shape on a slide."""
    # Try placeholder first
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            return shape
    # Fallback: find largest text shape near top
    for shape in slide.shapes:
        if shape.has_text_frame and shape.top < prs.slide_height // 3:
            return shape
    return None


def find_body_shape(slide):
    """Find the main body/content shape on a slide."""
    # Try placeholder first
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            return shape
    # Fallback: find largest text shape
    best = None
    best_area = 0
    for shape in slide.shapes:
        if shape.has_text_frame and shape.top > prs.slide_height // 4:
            area = shape.width * shape.height
            if area > best_area:
                best_area = area
                best = shape
    return best


def get_all_text_shapes(slide):
    """Get all shapes with text frames, sorted by position (top then left)."""
    text_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_shapes.append(shape)
    text_shapes.sort(key=lambda s: (s.top, s.left))
    return text_shapes


def add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = notes_text


# ── Step 4: Build the presentation ──
# We'll work with the existing template and modify its slides.
# First, determine how many slides the template has and what to keep/duplicate.

# Strategy: Keep template slide 0 (title), duplicate slide 1 (content) for all other slides
# Then modify text in each slide.

# Let's figure out the best approach based on the template structure
print("Building 28-slide presentation from template...")

# Determine which template slide to use as the "content" template
# Usually slide 0 = title, slide 1+ = content layouts
title_template_idx = 0
content_template_idx = min(1, slide_count - 1)  # Use slide 1 if available

print(f"Using template slide {title_template_idx} for title")
print(f"Using template slide {content_template_idx} for content slides")

# We need 28 slides total. Template has {slide_count} slides.
# We'll keep slide 0 (title) and duplicate content slides for the rest.

# First, duplicate enough content slides
slides_needed = 28
existing_slides = slide_count

if existing_slides < slides_needed:
    # Need to duplicate slides
    for i in range(slides_needed - existing_slides):
        print(f"  Duplicating slide {content_template_idx} -> new slide {existing_slides + i}")
        duplicate_slide(prs, content_template_idx)

print(f"Total slides now: {len(prs.slides)}")

# ── Step 5: Define all 28 slides content ──
# Each entry: (slide_index, title_text, body_paragraphs, speaker_notes)

SLIDES_CONTENT = [
    # ═══════════════════════════════════════════════════════════
    # SLIDE 0: Title Slide
    # ═══════════════════════════════════════════════════════════
    {
        'index': 0,
        'title': 'الذهنية الرقمية وممارسات الابتكار التقني',
        'subtitle': 'الوحدة الثانية\nجامعة نجران - كلية علوم الحاسب ونظم المعلومات',
        'notes': 'شريحة العنوان - المحاضرة التفاعلية للوحدة الثانية. مدة المحاضرة المقدرة: 45-60 دقيقة.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 1: Learning Objectives
    # ═══════════════════════════════════════════════════════════
    {
        'index': 1,
        'title': 'الاهداف التعليمية',
        'body': [
            {'text': 'بنهاية هذه المحاضرة، سيكون المتعلم قادرا على:', 'bold': True, 'font_size': 16, 'space_after': 6},
            {'text': 'تعداد عناصر الابتكار الثلاثة وتقنيات العصف الذهني الاساسية', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'توضيح مفهوم ريادة الاعمال وعلاقتها بالابتكار التقني', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'شرح المراحل الخمس للتفكير التصميمي ودور كل مرحلة', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'تطبيق طريقة SCAMPER لتوليد افكار ابداعية جديدة', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'استخدام مراحل التفكير التصميمي لبناء نموذج اولي', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'تحليل العلاقة بين التكنولوجيا وملاءمة السوق', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'تقييم فرص العمل الريادية بناء على معايير محددة', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'تصميم حل ابتكاري لمشكلة رقمية واقعية', 'bullet': True, 'level': 0, 'font_size': 14},
        ],
        'notes': 'اعرض الاهداف التعليمية الثمانية. هذه الاهداف مبنية على تصنيف بلوم وتتدرج من التذكر الى الابداع.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 2: Content Overview / Agenda
    # ═══════════════════════════════════════════════════════════
    {
        'index': 2,
        'title': 'محاور المحاضرة',
        'body': [
            {'text': 'المحور الاول: اساسيات الابتكار وتوليد الافكار', 'bullet': True, 'level': 0, 'font_size': 15, 'bold': True},
            {'text': 'المحور الثاني: ريادة الاعمال وعقلية الابتكار التقني', 'bullet': True, 'level': 0, 'font_size': 15, 'bold': True},
            {'text': 'المحور الثالث: التفكير التصميمي - المراحل الخمس', 'bullet': True, 'level': 0, 'font_size': 15, 'bold': True},
            {'text': 'المحور الرابع: المرونة وملاءمة السوق', 'bullet': True, 'level': 0, 'font_size': 15, 'bold': True},
            {'text': 'المحور الخامس: بناء العقلية الريادية والتغلب على التحديات', 'bullet': True, 'level': 0, 'font_size': 15, 'bold': True},
        ],
        'notes': 'استعرض محاور المحاضرة الخمسة. وضح للمتعلمين ان المحاضرة تتضمن اسئلة تفاعلية وانشطة بين الاقسام.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 3: What is Innovation?
    # ═══════════════════════════════════════════════════════════
    {
        'index': 3,
        'title': 'ما هو الابتكار؟',
        'body': [
            {'text': 'الابتكار هو ايجاد حلول جديدة للمشاكل الواقعية وتحويلها الى قيمة حقيقية', 'bold': True, 'font_size': 16, 'space_after': 8},
            {'text': 'عناصر الابتكار الثلاثة:', 'bold': True, 'font_size': 15, 'space_after': 4},
            {'text': 'حلول جديدة: تقديم افكار لم تكن موجودة من قبل لمعالجة مشكلات قائمة', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'تحسين مستمر: تطوير العمليات والمنتجات والخدمات الحالية', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'تقنيات ابداعية: استخدام اساليب وادوات جديدة في التنفيذ', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': '', 'font_size': 8},
            {'text': 'مثال: تطبيق "مرسول" في السعودية ابتكر حلا لتوصيل الطلبات من اي مكان وليس فقط المطاعم', 'font_size': 13, 'italic': True, 'color': '2E86C1'},
        ],
        'notes': 'اشرح مفهوم الابتكار مع التركيز على العناصر الثلاثة. استخدم مثال تطبيق مرسول كمثال محلي سعودي قريب من المتعلمين. اسأل المتعلمين: ما هو اخر ابتكار لاحظتموه في حياتكم اليومية؟'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 4: Importance of Idea Generation
    # ═══════════════════════════════════════════════════════════
    {
        'index': 4,
        'title': 'اهمية توليد الافكار',
        'body': [
            {'text': 'توليد الافكار هو الخطوة الاولى في رحلة الابتكار', 'bold': True, 'font_size': 16, 'space_after': 6},
            {'text': 'يتطلب التفكير خارج الصندوق لايجاد حلول غير تقليدية', 'font_size': 14, 'space_after': 6},
            {'text': 'تقنيات العصف الذهني الاساسية:', 'bold': True, 'font_size': 15, 'space_after': 4},
            {'text': 'التواصل الحر: اطلاق الافكار بحرية دون اصدار احكام او نقد', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'الخرائط الذهنية: ربط الافكار والمفاهيم بصريا في شكل شجري', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'طريقة SCAMPER: اداة منهجية لتوليد افكار جديدة (ستُشرح بالتفصيل)', 'bullet': True, 'level': 0, 'font_size': 14},
        ],
        'notes': '[تعليمات تفاعلية]\nالنوع: انقر للكشف\nالوصف: تظهر التقنيات الثلاث واحدة تلو الاخرى عند النقر. كل تقنية تظهر مع رسم توضيحي بسيط يمثلها.\n\nاشرح كل تقنية مع مثال سريع. التواصل الحر مثل جلسات العصف الذهني في الشركات الكبرى.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 5: SCAMPER Method (Interactive)
    # ═══════════════════════════════════════════════════════════
    {
        'index': 5,
        'title': 'طريقة SCAMPER لتوليد الافكار',
        'body': [
            {'text': '[تعليمات تفاعلية] النوع: انقر للكشف - اضغط على كل حرف لكشف معناه ومثاله', 'font_size': 11, 'color': 'E74C3C', 'bold': True, 'space_after': 4},
            {'text': 'S - Substitute (استبدال): استبدل مكونا بآخر', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'مثال: استبدال المفتاح التقليدي بالبصمة في الهواتف الذكية', 'bullet': True, 'level': 1, 'font_size': 12, 'italic': True, 'color': '7F8C8D'},
            {'text': 'C - Combine (دمج): ادمج فكرتين او اكثر معا', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'مثال: دمج الهاتف والكاميرا والحاسوب في جهاز واحد (الهاتف الذكي)', 'bullet': True, 'level': 1, 'font_size': 12, 'italic': True, 'color': '7F8C8D'},
            {'text': 'A - Adapt (تكيف): كيف يمكن تكييف فكرة موجودة؟', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'M - Modify (تعديل): عدل الشكل او الحجم او اللون', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'P - Put to other use (استخدام آخر): استخدمه في سياق جديد', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'E - Eliminate (حذف): ماذا لو حذفت جزءا منه؟', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'R - Reverse (عكس): اعكس الترتيب او العملية', 'bullet': True, 'level': 0, 'font_size': 13},
        ],
        'notes': '[تعليمات تفاعلية]\nالنوع: انقر للكشف (Tabs / Accordion)\nالوصف: كل حرف من SCAMPER يظهر كتبويب. عند النقر على كل حرف:\n1. يظهر اسم العنصر بالعربي والانجليزي\n2. يظهر التعريف\n3. يظهر مثال تطبيقي من الواقع\n\nالتفاعل: 7 تبويبات يضغط عليها المتعلم واحدا تلو الآخر.\n\nمثال تطبيقي شامل: خذ "تطبيق توصيل طعام" وطبق كل عنصر من SCAMPER عليه:\nS: استبدل السائق بطائرة درون\nC: ادمج التوصيل مع خدمة البقالة\nA: كيّف التطبيق ليعمل في المناطق الريفية\nM: عدل الحد الادنى للطلب\nP: استخدم التطبيق لتوصيل الادوية\nE: احذف رسوم التوصيل\nR: اعكس العملية - العميل يذهب للمطعم والتطبيق يحجز له'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 6: Knowledge Check 1
    # ═══════════════════════════════════════════════════════════
    {
        'index': 6,
        'title': 'اختبر معلوماتك',
        'body': [
            {'text': '[تعليمات تفاعلية] النوع: اختيار من متعدد', 'font_size': 11, 'color': 'E74C3C', 'bold': True, 'space_after': 6},
            {'text': 'السؤال: اي من التالي يمثل عناصر الابتكار الثلاثة بشكل صحيح؟', 'bold': True, 'font_size': 16, 'space_after': 8},
            {'text': 'أ) التخطيط والتنفيذ والتقييم', 'font_size': 15, 'space_after': 4},
            {'text': 'ب) حلول جديدة وتحسين مستمر وتقنيات ابداعية', 'font_size': 15, 'space_after': 4, 'color': '27AE60', 'bold': True},
            {'text': 'ج) البحث والتطوير والتسويق', 'font_size': 15, 'space_after': 4},
            {'text': 'د) الفكرة والتمويل والتنفيذ', 'font_size': 15, 'space_after': 4},
            {'text': '', 'font_size': 6},
            {'text': 'الاجابة الصحيحة: ب', 'font_size': 13, 'bold': True, 'color': '27AE60'},
            {'text': 'التغذية الراجعة: الابتكار يرتكز على ايجاد حلول جديدة وتحسين ما هو موجود واستخدام تقنيات ابداعية', 'font_size': 12, 'italic': True, 'color': '2980B9'},
        ],
        'notes': '[تعليمات تفاعلية]\nالنوع: اختيار من متعدد\nالوصف: يختار المتعلم اجابة واحدة من اربعة خيارات.\nالاجابة الصحيحة: ب) حلول جديدة وتحسين مستمر وتقنيات ابداعية\nالتغذية الراجعة للاجابة الصحيحة: احسنت! الابتكار يقوم على ثلاثة عناصر اساسية.\nالتغذية الراجعة للاجابة الخاطئة: راجع شريحة "ما هو الابتكار" - العناصر الثلاثة هي حلول جديدة وتحسين وتقنيات ابداعية.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 7: What is Entrepreneurship?
    # ═══════════════════════════════════════════════════════════
    {
        'index': 7,
        'title': 'ما هي ريادة الاعمال؟',
        'body': [
            {'text': 'ريادة الاعمال هي تحويل الافكار المبتكرة الى مشاريع ناجحة ومستدامة', 'bold': True, 'font_size': 16, 'space_after': 6},
            {'text': 'المتطلبات العقلية لرائد الاعمال:', 'bold': True, 'font_size': 15, 'space_after': 4},
            {'text': 'حل المشكلات: القدرة على تحديد المشكلات وايجاد حلول عملية', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'الابداع: التفكير بطرق غير تقليدية لتقديم قيمة مضافة', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'المرونة: التكيف مع التغيرات والتحديات غير المتوقعة', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': '', 'font_size': 6},
            {'text': 'رواد الاعمال يحددون الفرص ويتحملون المخاطر المحسوبة لتطوير منتجات او خدمات جديدة تلبي احتياجات السوق', 'font_size': 13, 'italic': True},
            {'text': 'مثال: رائد الاعمال السعودي عبدالله السبيعي مؤسس "هنقرستيشن" حوّل فكرة بسيطة لتوصيل الطعام الى شركة بمليارات', 'font_size': 13, 'italic': True, 'color': '2E86C1'},
        ],
        'notes': 'اشرح مفهوم ريادة الاعمال مع التاكيد على الفرق بينها وبين الابتكار. ريادة الاعمال تتطلب تحويل الابتكار الى مشروع تجاري. استخدم مثال هنقرستيشن كنموذج سعودي ناجح.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 8: Technology Innovation Mindset
    # ═══════════════════════════════════════════════════════════
    {
        'index': 8,
        'title': 'عقلية الابتكار التقني',
        'body': [
            {'text': 'القدرة على تطوير حلول ابداعية لمشاكل واقعية باستخدام التكنولوجيا', 'bold': True, 'font_size': 15, 'space_after': 6},
            {'text': 'لا تقتصر على شركات التكنولوجيا بل تنطبق على اي مجال يشهد تحولا رقميا', 'font_size': 14, 'space_after': 6},
            {'text': 'المتطلبات الرئيسية الثلاثة:', 'bold': True, 'font_size': 15, 'color': '2E86C1', 'space_after': 4},
            {'text': '1. التفكير التصميمي: فهم احتياجات المستخدم قبل تطوير الحلول', 'bullet': True, 'level': 0, 'font_size': 14, 'bold': True},
            {'text': '2. المرونة والقدرة على التكيف: الاستجابة للتغيرات التكنولوجية', 'bullet': True, 'level': 0, 'font_size': 14, 'bold': True},
            {'text': '3. التكنولوجيا وملاءمة السوق: ضمان تلبية المنتجات لاحتياجات السوق', 'bullet': True, 'level': 0, 'font_size': 14, 'bold': True},
        ],
        'notes': '[تعليمات تفاعلية]\nالنوع: نقاط ساخنة (Hotspots)\nالوصف: يظهر مخطط دائري يحتوي على المتطلبات الثلاثة. عند النقر على كل متطلب يظهر شرح تفصيلي مع مثال.\n\nالمتطلب 1 (التفكير التصميمي) سيُشرح بالتفصيل في الشرائح القادمة.\nالمتطلب 2 (المرونة) يُشرح في شريحة 18.\nالمتطلب 3 (ملاءمة السوق) يُشرح في شريحة 20.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 9: Knowledge Check 2
    # ═══════════════════════════════════════════════════════════
    {
        'index': 9,
        'title': 'اختبر معلوماتك',
        'body': [
            {'text': '[تعليمات تفاعلية] النوع: سحب وافلات (مطابقة)', 'font_size': 11, 'color': 'E74C3C', 'bold': True, 'space_after': 6},
            {'text': 'السؤال: طابق كل متطلب من متطلبات عقلية الابتكار التقني مع وصفه الصحيح', 'bold': True, 'font_size': 15, 'space_after': 6},
            {'text': 'المتطلبات:', 'bold': True, 'font_size': 14, 'space_after': 2},
            {'text': 'التفكير التصميمي ← فهم احتياجات المستخدم قبل تطوير الحلول', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'المرونة والتكيف ← الاستجابة للتغيرات التكنولوجية المتسارعة', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'ملاءمة السوق ← ضمان تلبية المنتجات الرقمية لاحتياجات السوق الفعلية', 'bullet': True, 'level': 0, 'font_size': 13},
        ],
        'notes': '[تعليمات تفاعلية]\nالنوع: سحب وافلات (Drag & Drop)\nالوصف: على الجانب الايمن تظهر المتطلبات الثلاثة، وعلى الجانب الايسر تظهر الاوصاف. يسحب المتعلم كل متطلب ويضعه بجانب وصفه الصحيح.\n\nالاجابات:\n- التفكير التصميمي = فهم احتياجات المستخدم\n- المرونة والتكيف = الاستجابة للتغيرات\n- ملاءمة السوق = تلبية الاحتياجات الفعلية\n\nالتغذية الراجعة: عند الاجابة الصحيحة يظهر علامة صح خضراء مع صوت نجاح.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 10: Section Transition - Design Thinking
    # ═══════════════════════════════════════════════════════════
    {
        'index': 10,
        'title': 'المحور الثالث',
        'body': [
            {'text': 'التفكير التصميمي', 'bold': True, 'font_size': 28, 'alignment': PP_ALIGN.CENTER, 'space_after': 8},
            {'text': 'Design Thinking', 'font_size': 20, 'alignment': PP_ALIGN.CENTER, 'color': '2E86C1', 'italic': True, 'space_after': 12},
            {'text': 'نهج ابتكاري متمحور حول الانسان لحل المشكلات', 'font_size': 16, 'alignment': PP_ALIGN.CENTER},
            {'text': 'تستخدمه شركات عالمية مثل Google و Apple و IDEO', 'font_size': 14, 'alignment': PP_ALIGN.CENTER, 'color': '7F8C8D'},
        ],
        'notes': 'شريحة انتقالية للمحور الثالث. اذكر ان التفكير التصميمي هو المتطلب الاول لعقلية الابتكار التقني وسنتعمق فيه خلال الشرائح القادمة.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 11: Design Thinking Overview
    # ═══════════════════════════════════════════════════════════
    {
        'index': 11,
        'title': 'ما هو التفكير التصميمي؟',
        'body': [
            {'text': 'نهج ابتكاري متمحور حول الانسان يركز على فهم احتياجات المستخدم', 'bold': True, 'font_size': 15, 'space_after': 6},
            {'text': 'الخصائص الرئيسية:', 'bold': True, 'font_size': 14, 'space_after': 4},
            {'text': 'يضع المستخدم في مركز عملية التصميم والتطوير', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'يستخدم النهج التكراري: بناء > اختبار > تحسين > تكرار', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'يشجع التعاون بين تخصصات مختلفة للوصول لحلول شاملة', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'يقبل الفشل كجزء طبيعي من عملية التعلم والتحسين', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': '', 'font_size': 6},
            {'text': 'مثال: شركة IDEO صممت عربة تسوق جديدة في 5 ايام فقط باستخدام التفكير التصميمي', 'font_size': 13, 'italic': True, 'color': '2E86C1'},
        ],
        'notes': 'اشرح التفكير التصميمي كمنهجية عملية وليس مجرد نظرية. اذكر ان ستانفورد هي من طورت هذا المنهج وان شركات مثل جوجل وآبل تطبقه يوميا.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 12: 5 Stages of Design Thinking (Interactive)
    # ═══════════════════════════════════════════════════════════
    {
        'index': 12,
        'title': 'المراحل الخمس للتفكير التصميمي',
        'body': [
            {'text': '[تعليمات تفاعلية] النوع: كشف تسلسلي - تظهر كل مرحلة عند النقر', 'font_size': 11, 'color': 'E74C3C', 'bold': True, 'space_after': 4},
            {'text': '[وصف بصري: مخطط خطي افقي يظهر 5 دوائر متصلة بخط]', 'font_size': 11, 'color': '7F8C8D', 'italic': True, 'space_after': 4},
            {'text': '1. التعاطف (Empathize): فهم مشاعر واحتياجات المستخدمين', 'bullet': True, 'level': 0, 'font_size': 14, 'bold': True},
            {'text': '2. التحديد (Define): صياغة المشكلة بوضوح ودقة', 'bullet': True, 'level': 0, 'font_size': 14, 'bold': True},
            {'text': '3. التفكير (Ideate): توليد اكبر عدد من الحلول الممكنة', 'bullet': True, 'level': 0, 'font_size': 14, 'bold': True},
            {'text': '4. النموذج الاولي (Prototype): بناء نموذج بسيط للحل', 'bullet': True, 'level': 0, 'font_size': 14, 'bold': True},
            {'text': '5. الاختبار (Test): تجربة الحل مع المستخدمين وجمع الملاحظات', 'bullet': True, 'level': 0, 'font_size': 14, 'bold': True},
        ],
        'notes': '[تعليمات تفاعلية]\nالنوع: كشف تسلسلي (Sequential Reveal)\nالوصف: يظهر مخطط خطي في اعلى الشريحة يحتوي على 5 دوائر متصلة. عند النقر على كل دائرة:\n1. تتلون الدائرة\n2. يظهر اسم المرحلة بالعربي والانجليزي\n3. يظهر وصف مختصر اسفل المخطط\n\nالمراحل تظهر من اليسار لليمين (Empathize > Define > Ideate > Prototype > Test)\nمع سهم يربط بين كل مرحلة والتالية.\n\nبعد كشف جميع المراحل يظهر سهم عودة من Test الى Empathize ليوضح الطبيعة التكرارية.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 13: Empathy & Needs Identification
    # ═══════════════════════════════════════════════════════════
    {
        'index': 13,
        'title': 'التعاطف وتحديد الاحتياجات',
        'body': [
            {'text': 'المرحلة الاولى والاهم: فهم المستخدم الحقيقي وليس ما نظن انه يحتاجه', 'bold': True, 'font_size': 15, 'space_after': 6},
            {'text': 'اساليب التعاطف وجمع المعلومات:', 'bold': True, 'font_size': 14, 'space_after': 4},
            {'text': 'الملاحظة المباشرة: مراقبة المستخدمين اثناء تفاعلهم مع المنتج', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'المقابلات الشخصية: طرح اسئلة مفتوحة لفهم التجربة بعمق', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'اختبارات المستخدمين: تجربة النماذج الاولية مع مستخدمين حقيقيين', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'استطلاعات الرأي: جمع بيانات كمية من عدد كبير من المستخدمين', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': '', 'font_size': 6},
            {'text': 'مثال: عندما صممت Airbnb تجربتها، عاش المؤسسون مع المضيفين لفهم تجربتهم الحقيقية', 'font_size': 12, 'italic': True, 'color': '2E86C1'},
        ],
        'notes': 'ركز على ان التعاطف ليس مجرد سؤال المستخدم بل فهم سياقه الكامل. استخدم مثال Airbnb لتوضيح كيف ان المؤسسين عاشوا التجربة بانفسهم.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 14: Define, Ideate, Prototype
    # ═══════════════════════════════════════════════════════════
    {
        'index': 14,
        'title': 'التحديد والتفكير والنمذجة',
        'body': [
            {'text': 'المرحلة 2 - التحديد:', 'bold': True, 'font_size': 15, 'color': '2E86C1', 'space_after': 2},
            {'text': 'صياغة المشكلة بجملة واضحة تركز على المستخدم', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'مثال: "طلاب الجامعة يحتاجون طريقة سهلة لتنظيم مشاريعهم الجماعية"', 'bullet': True, 'level': 1, 'font_size': 12, 'italic': True, 'color': '7F8C8D'},
            {'text': 'المرحلة 3 - التفكير (التصور):', 'bold': True, 'font_size': 15, 'color': '2E86C1', 'space_after': 2},
            {'text': 'توليد اكبر عدد ممكن من الافكار دون تقييم (الكمية قبل الجودة)', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'استخدام تقنيات مثل SCAMPER والعصف الذهني', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'المرحلة 4 - النموذج الاولي:', 'bold': True, 'font_size': 15, 'color': '2E86C1', 'space_after': 2},
            {'text': 'بناء نموذج بسيط وسريع لاختبار الفكرة', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'انواع النماذج: ورقية (رسومات) / رقمية (تفاعلية) / مادية (مجسمات)', 'bullet': True, 'level': 0, 'font_size': 13},
        ],
        'notes': 'اشرح المراحل الثلاث بشكل مختصر. ركز على ان مرحلة التفكير تشجع الكمية وليس الجودة في البداية. في مرحلة النمذجة، وضح ان النموذج لا يحتاج ان يكون مثاليا بل سريعا وبسيطا.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 15: Testing & Iteration
    # ═══════════════════════════════════════════════════════════
    {
        'index': 15,
        'title': 'الاختبار والتكرار',
        'body': [
            {'text': 'المرحلة 5 - الاختبار:', 'bold': True, 'font_size': 16, 'color': '2E86C1', 'space_after': 4},
            {'text': 'تقييم الحل مع مستخدمين حقيقيين وجمع ملاحظاتهم', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'تحليل نتائج الاختبار وتحديد نقاط التحسين', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'تكرار العملية: العودة لاي مرحلة سابقة عند الحاجة', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': '', 'font_size': 6},
            {'text': 'الطبيعة التكرارية للتفكير التصميمي:', 'bold': True, 'font_size': 15, 'space_after': 4},
            {'text': '[وصف بصري: مخطط دائري يوضح: بناء > اختبار > تعلم > تحسين > اعادة البناء]', 'font_size': 12, 'italic': True, 'color': '7F8C8D'},
            {'text': 'العملية ليست خطية - يمكن العودة من اي مرحلة لاي مرحلة سابقة حسب نتائج الاختبار', 'font_size': 13},
        ],
        'notes': 'اكد ان التفكير التصميمي عملية تكرارية وليست خطية. نتائج الاختبار قد تعيدنا لمرحلة التعاطف لفهم افضل. اعط مثالا: لو اختبرنا تطبيقا ووجدنا ان المستخدمين لا يفهمون واجهته، نعود لمرحلة التعاطف لفهم كيف يفكرون.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 16: Knowledge Check 3 (Drag & Drop)
    # ═══════════════════════════════════════════════════════════
    {
        'index': 16,
        'title': 'اختبر معلوماتك',
        'body': [
            {'text': '[تعليمات تفاعلية] النوع: سحب وافلات - ترتيب', 'font_size': 11, 'color': 'E74C3C', 'bold': True, 'space_after': 6},
            {'text': 'السؤال: رتب مراحل التفكير التصميمي الخمس بالترتيب الصحيح', 'bold': True, 'font_size': 16, 'space_after': 6},
            {'text': 'العناصر المبعثرة للترتيب:', 'bold': True, 'font_size': 14, 'space_after': 4},
            {'text': 'الاختبار | التعاطف | النموذج الاولي | التحديد | التفكير', 'font_size': 15, 'space_after': 6},
            {'text': 'الترتيب الصحيح:', 'bold': True, 'font_size': 14, 'color': '27AE60', 'space_after': 4},
            {'text': '1. التعاطف > 2. التحديد > 3. التفكير > 4. النموذج الاولي > 5. الاختبار', 'font_size': 14, 'color': '27AE60'},
        ],
        'notes': '[تعليمات تفاعلية]\nالنوع: سحب وافلات - ترتيب (Sorting)\nالوصف: تظهر 5 بطاقات مبعثرة تحتوي على اسماء مراحل التفكير التصميمي. يسحب المتعلم كل بطاقة ويضعها في الترتيب الصحيح من 1 الى 5.\n\nالترتيب الصحيح:\n1. التعاطف (Empathize)\n2. التحديد (Define)\n3. التفكير (Ideate)\n4. النموذج الاولي (Prototype)\n5. الاختبار (Test)\n\nالتغذية الراجعة: بعد الترتيب الصحيح يظهر المخطط الخطي الكامل مع شرح مختصر لكل مرحلة.\nعند الخطأ: يتحرك العنصر للموقع الصحيح مع توضيح السبب.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 17: Flexibility vs Adaptability (Interactive)
    # ═══════════════════════════════════════════════════════════
    {
        'index': 17,
        'title': 'المرونة مقابل القدرة على التكيف',
        'body': [
            {'text': '[تعليمات تفاعلية] النوع: مقارنة بصرية - انقر للكشف', 'font_size': 11, 'color': 'E74C3C', 'bold': True, 'space_after': 4},
            {'text': 'المرونة (Agility):', 'bold': True, 'font_size': 16, 'color': '2E86C1', 'space_after': 2},
            {'text': 'القدرة على الاستجابة للتحديات او الفرص الفورية بكفاءة وسرعة', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'مثال: تعديل تطبيق خلال ايام استجابة لشكاوى المستخدمين', 'bullet': True, 'level': 1, 'font_size': 12, 'italic': True, 'color': '7F8C8D'},
            {'text': '', 'font_size': 4},
            {'text': 'القدرة على التكيف (Adaptability):', 'bold': True, 'font_size': 16, 'color': '27AE60', 'space_after': 2},
            {'text': 'القدرة على التطور والازدهار استجابة للتحولات الجذرية طويلة الامد', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'مثال: تحول نتفليكس من تأجير DVD بالبريد الى البث الرقمي', 'bullet': True, 'level': 1, 'font_size': 12, 'italic': True, 'color': '7F8C8D'},
        ],
        'notes': '[تعليمات تفاعلية]\nالنوع: مقارنة بصرية (Visual Comparison) - انقر للكشف\nالوصف: الشاشة مقسومة نصفين. الجانب الايمن "المرونة" والايسر "التكيف". عند النقر على كل جانب:\n1. يظهر التعريف\n2. يظهر مثال واقعي\n3. يظهر رمز بصري (ساعة للمرونة = استجابة سريعة، شجرة للتكيف = نمو طويل)\n\nالرسالة الرئيسية: كلاهما ضروري - المرونة للتحديات القصيرة والتكيف للتحولات الكبرى.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 18: Identifying Business Opportunities
    # ═══════════════════════════════════════════════════════════
    {
        'index': 18,
        'title': 'تحديد فرص العمل',
        'body': [
            {'text': 'ليست كل فكرة جيدة فرصة عمل قابلة للاستمرار', 'bold': True, 'font_size': 15, 'space_after': 6},
            {'text': 'معايير تقييم فرصة العمل:', 'bold': True, 'font_size': 15, 'space_after': 4},
            {'text': 'طلب السوق: هل هناك عدد كاف من العملاء المحتملين؟', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'حل مشكلة حقيقية: هل الفكرة تعالج مشكلة يعاني منها الناس فعلا؟', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'احتياجات العملاء: هل تم التحقق من ان العملاء يريدون هذا الحل؟', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'دراسة الجدوى: هل المشروع ممكن تقنيا وماليا؟', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'قابلية التوسع: هل يمكن للمشروع ان ينمو ويتوسع؟', 'bullet': True, 'level': 0, 'font_size': 14},
        ],
        'notes': 'وضح الفرق بين "فكرة جيدة" و"فرصة عمل". كثير من الافكار الرائعة تفشل لانها لا تلبي حاجة حقيقية في السوق. اعط مثالا: Google Glass كانت فكرة مبتكرة لكنها فشلت لان السوق لم يكن جاهزا.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 19: Technology & Market Fit
    # ═══════════════════════════════════════════════════════════
    {
        'index': 19,
        'title': 'التكنولوجيا وملاءمة السوق',
        'body': [
            {'text': 'ملاءمة المنتج للسوق (Product-Market Fit) هي اللحظة التي يلبي فيها منتجك حاجة حقيقية', 'bold': True, 'font_size': 14, 'space_after': 6},
            {'text': 'كيف تحقق ملاءمة السوق؟', 'bold': True, 'font_size': 15, 'space_after': 4},
            {'text': 'فهم احتياجات السوق من خلال الابحاث والمقابلات', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'الاستماع لملاحظات العملاء وتحليلها باستمرار', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'التأكد من ان التكنولوجيا تتوافق مع احتياجات المستخدم', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'استخدام النهج التكراري: بناء > قياس > تعلم > تحسين', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': '', 'font_size': 6},
            {'text': 'مثال: Slack بدأ كأداة داخلية لشركة العاب ثم اكتشفوا ان ملاءمته الحقيقية في سوق التواصل المؤسسي', 'font_size': 12, 'italic': True, 'color': '2E86C1'},
        ],
        'notes': 'اشرح مفهوم Product-Market Fit. وضح ان كثيرا من الشركات الناجحة غيرت اتجاهها (pivot) حتى وجدت الملاءمة. مثال Slack: بدأ كأداة لتطوير لعبة Glitch ثم تحول لمنصة تواصل.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 20: Prototyping & A/B Testing
    # ═══════════════════════════════════════════════════════════
    {
        'index': 20,
        'title': 'النمذجة والتجريب واختبار A/B',
        'body': [
            {'text': 'انواع النماذج الاولية:', 'bold': True, 'font_size': 15, 'space_after': 4},
            {'text': 'نماذج ورقية: رسومات تخطيطية سريعة على الورق (اسرع وارخص)', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'نماذج رقمية: نماذج تفاعلية باستخدام ادوات مثل Figma', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'نماذج مادية: مجسمات فعلية للمنتج (للمنتجات المادية)', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': '', 'font_size': 4},
            {'text': 'اختبار A/B:', 'bold': True, 'font_size': 15, 'space_after': 4},
            {'text': 'مقارنة نسختين مختلفتين من المنتج لمعرفة ايهما افضل', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'مثال: اختبار لونين مختلفين لزر "اشتر الان" - الاحمر حقق مبيعات اعلى بـ 21%', 'bullet': True, 'level': 1, 'font_size': 12, 'italic': True, 'color': '7F8C8D'},
            {'text': 'التكرار المستمر يؤدي الى منتجات افضل تلبي احتياجات المستخدمين', 'font_size': 13, 'space_before': 4},
        ],
        'notes': 'اشرح الفرق بين انواع النماذج الثلاثة. وضح ان اختبار A/B يستخدمه الجميع من امازون الى نتفليكس. المبدأ بسيط: اعرض نسختين لمجموعتين مختلفتين وقس النتائج.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 21: Knowledge Check 4
    # ═══════════════════════════════════════════════════════════
    {
        'index': 21,
        'title': 'اختبر معلوماتك',
        'body': [
            {'text': '[تعليمات تفاعلية] النوع: اختيار من متعدد', 'font_size': 11, 'color': 'E74C3C', 'bold': True, 'space_after': 6},
            {'text': 'السؤال: شركة ناشئة طورت تطبيقا لتنظيم الملفات لكن المستخدمين لم يستخدموه. ما السبب الاكثر ترجيحا؟', 'bold': True, 'font_size': 14, 'space_after': 6},
            {'text': 'أ) التكنولوجيا المستخدمة قديمة', 'font_size': 14, 'space_after': 3},
            {'text': 'ب) لم يتحقق ملاءمة المنتج للسوق - لم يحل مشكلة حقيقية', 'font_size': 14, 'space_after': 3, 'color': '27AE60', 'bold': True},
            {'text': 'ج) سعر التطبيق مرتفع جدا', 'font_size': 14, 'space_after': 3},
            {'text': 'د) عدم وجود اعلانات كافية', 'font_size': 14, 'space_after': 6},
            {'text': 'الاجابة: ب - ملاءمة المنتج للسوق هي العامل الاهم لنجاح اي منتج تقني', 'font_size': 12, 'bold': True, 'color': '27AE60'},
        ],
        'notes': '[تعليمات تفاعلية]\nالنوع: اختيار من متعدد (سيناريو)\nالوصف: سيناريو واقعي يتبعه سؤال. يختار المتعلم اجابة واحدة.\nالاجابة الصحيحة: ب\nالتغذية الراجعة: المنتج الناجح يبدأ من فهم حاجة السوق الحقيقية وليس من بناء التكنولوجيا اولا.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 22: Building Entrepreneurial Mindset
    # ═══════════════════════════════════════════════════════════
    {
        'index': 22,
        'title': 'بناء عقلية ريادية',
        'body': [
            {'text': 'كيف تبني عقلية ريادية؟', 'bold': True, 'font_size': 16, 'space_after': 6},
            {'text': 'رؤية الفشل كفرصة للتعلم وليس نهاية الطريق', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'مثال: توماس اديسون فشل 1000 مرة قبل اختراع المصباح الكهربائي', 'bullet': True, 'level': 1, 'font_size': 12, 'italic': True, 'color': '7F8C8D'},
            {'text': 'الانفتاح على التغذية الراجعة والتكيف بناء عليها', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'البقاء فضوليا والسعي للتحسين المستمر', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'بناء شبكة علاقات مع رواد اعمال وخبراء', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'حضور الفعاليات والانضمام الى حاضنات ومسرعات الاعمال', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': '', 'font_size': 4},
            {'text': 'رواد الاعمال الناجحون يتعلمون من كل تجربة ويبنون على ملاحظات الآخرين', 'font_size': 13, 'italic': True},
        ],
        'notes': 'ركز على ان العقلية الريادية مهارة مكتسبة وليست موهبة فطرية. يمكن لاي شخص ان يطورها من خلال الممارسة والتعلم المستمر. اذكر امثلة لحاضنات اعمال سعودية مثل "منشآت" و"بادر".'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 23: Tools & Resources
    # ═══════════════════════════════════════════════════════════
    {
        'index': 23,
        'title': 'ادوات ومصادر لرائد الاعمال',
        'body': [
            {'text': 'Lean Startup (الشركات الناشئة المرنة):', 'bold': True, 'font_size': 15, 'color': '2E86C1', 'space_after': 2},
            {'text': 'منهجية للتطوير السريع تقوم على حلقة: ابنِ > قِس > تعلّم', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'الهدف: اختبار فرضياتك بأقل تكلفة واسرع وقت', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': '', 'font_size': 4},
            {'text': 'Business Model Canvas (نموذج العمل التجاري):', 'bold': True, 'font_size': 15, 'color': '2E86C1', 'space_after': 2},
            {'text': 'صفحة واحدة تلخص مشروعك في 9 عناصر اساسية', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'العناصر: القيمة المقدمة، شرائح العملاء، القنوات، العلاقات، مصادر الايراد، الموارد، الانشطة، الشركاء، التكاليف', 'bullet': True, 'level': 0, 'font_size': 12},
            {'text': '', 'font_size': 4},
            {'text': 'مصادر للتعلم المستمر:', 'bold': True, 'font_size': 15, 'color': '2E86C1', 'space_after': 2},
            {'text': 'Stanford Online (d.school) لدورات التفكير التصميمي المجانية', 'bullet': True, 'level': 0, 'font_size': 13},
        ],
        'notes': '[تعليمات تفاعلية]\nالنوع: تبويبات (Tabs)\nالوصف: 3 تبويبات:\nتبويب 1: Lean Startup - يظهر مخطط حلقة Build-Measure-Learn\nتبويب 2: Business Model Canvas - يظهر النموذج التسعيني\nتبويب 3: مصادر التعلم - روابط ومنصات\n\nعند النقر على كل تبويب يظهر المحتوى الخاص به مع رسم توضيحي.\n\nملاحظة: Lean Startup ابتكرها Eric Ries.\nBusiness Model Canvas ابتكرها Alexander Osterwalder.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 24: Overcoming Innovation Challenges (Expanded)
    # ═══════════════════════════════════════════════════════════
    {
        'index': 24,
        'title': 'التغلب على تحديات الابتكار',
        'body': [
            {'text': 'التحديات الرئيسية:', 'bold': True, 'font_size': 15, 'color': 'E74C3C', 'space_after': 4},
            {'text': 'مقاومة التغيير: الخوف من المجهول والتمسك بالطرق التقليدية', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'الخوف من الفشل: تجنب المخاطرة خوفا من الخسارة', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'نقص الموارد: محدودية التمويل والوقت والكفاءات', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'غياب ثقافة الابتكار: بيئة عمل لا تشجع التجريب', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': '', 'font_size': 4},
            {'text': 'الحلول العملية:', 'bold': True, 'font_size': 15, 'color': '27AE60', 'space_after': 4},
            {'text': 'تشجيع التجربة: خلق بيئة آمنة لتجربة اشياء جديدة دون خوف', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'تعزيز ثقافة الانفتاح: تشجيع التواصل المفتوح ومشاركة الافكار', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'البدء صغيرا: اختبار الافكار بنماذج بسيطة قبل الاستثمار الكبير', 'bullet': True, 'level': 0, 'font_size': 14},
        ],
        'notes': 'هذا القسم تم اثراؤه عن المحتوى الاصلي الذي كان يحتوي على تحديين فقط. اضفنا تحديات اضافية (نقص الموارد وغياب الثقافة) وحلولا عملية اضافية (البدء صغيرا). اسأل المتعلمين: ما اكبر تحد يواجهك عند محاولة الابتكار؟'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 25: Knowledge Check 5 (Classification)
    # ═══════════════════════════════════════════════════════════
    {
        'index': 25,
        'title': 'اختبر معلوماتك',
        'body': [
            {'text': '[تعليمات تفاعلية] النوع: تصنيف - سحب وافلات', 'font_size': 11, 'color': 'E74C3C', 'bold': True, 'space_after': 6},
            {'text': 'السؤال: صنف السلوكيات التالية: هل تدعم او تعيق العقلية الريادية؟', 'bold': True, 'font_size': 15, 'space_after': 6},
            {'text': 'سلوكيات داعمة:', 'bold': True, 'font_size': 14, 'color': '27AE60', 'space_after': 2},
            {'text': 'رؤية الفشل كفرصة للتعلم', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'طلب التغذية الراجعة باستمرار', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'التعلم من تجارب الآخرين', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': '', 'font_size': 4},
            {'text': 'سلوكيات معيقة:', 'bold': True, 'font_size': 14, 'color': 'E74C3C', 'space_after': 2},
            {'text': 'الخوف من النقد وتجنب المخاطرة', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'التمسك بالخطة الاصلية رغم فشلها', 'bullet': True, 'level': 0, 'font_size': 13},
            {'text': 'العمل بمعزل عن الفريق والعملاء', 'bullet': True, 'level': 0, 'font_size': 13},
        ],
        'notes': '[تعليمات تفاعلية]\nالنوع: تصنيف (Classification)\nالوصف: تظهر 8 بطاقات سلوكية مبعثرة في وسط الشاشة. يوجد عمودان:\n- العمود الايمن (اخضر): سلوكيات تدعم العقلية الريادية\n- العمود الايسر (احمر): سلوكيات تعيق العقلية الريادية\n\nيسحب المتعلم كل بطاقة ويضعها في العمود المناسب.\n\nالسلوكيات الداعمة: رؤية الفشل كتعلم، طلب التغذية الراجعة، التعلم من الآخرين، التحسين المستمر\nالسلوكيات المعيقة: الخوف من النقد، التمسك بالخطة الفاشلة، العزلة، انتظار الكمال\n\nالتغذية الراجعة: عند وضع البطاقة في المكان الصحيح تتلون بالاخضر، وعند الخطأ تتلون بالاحمر وتعود لمكانها.'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 26: Summary / Key Takeaways
    # ═══════════════════════════════════════════════════════════
    {
        'index': 26,
        'title': 'ملخص الوحدة',
        'body': [
            {'text': 'النقاط الرئيسية التي تعلمناها:', 'bold': True, 'font_size': 16, 'space_after': 6},
            {'text': 'الابتكار = حلول جديدة + تحسين مستمر + تقنيات ابداعية', 'bullet': True, 'level': 0, 'font_size': 14, 'bold': True},
            {'text': 'ريادة الاعمال تحول الابتكار الى مشاريع ناجحة ومستدامة', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'عقلية الابتكار التقني تتطلب: تفكير تصميمي + مرونة + ملاءمة سوق', 'bullet': True, 'level': 0, 'font_size': 14, 'bold': True},
            {'text': 'التفكير التصميمي: تعاطف > تحديد > تفكير > نمذجة > اختبار', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'SCAMPER اداة فعالة لتوليد الافكار بشكل منهجي', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'ملاءمة المنتج للسوق اهم من التكنولوجيا المتقدمة', 'bullet': True, 'level': 0, 'font_size': 14},
            {'text': 'الفشل جزء طبيعي من رحلة الابتكار وفرصة للتعلم', 'bullet': True, 'level': 0, 'font_size': 14},
        ],
        'notes': 'شريحة الملخص. راجع النقاط الرئيسية مع المتعلمين واسألهم: ما هو اهم مفهوم تعلمتموه اليوم؟ كيف ستطبقونه في مشاريعكم؟'
    },

    # ═══════════════════════════════════════════════════════════
    # SLIDE 27: Next Steps
    # ═══════════════════════════════════════════════════════════
    {
        'index': 27,
        'title': 'الخطوات القادمة',
        'body': [
            {'text': 'ما بعد المحاضرة:', 'bold': True, 'font_size': 16, 'space_after': 6},
            {'text': 'اكمال الانشطة التفاعلية المطلوبة في المنصة', 'bullet': True, 'level': 0, 'font_size': 15},
            {'text': 'المشاركة في منتدى النقاش حول العقلية الرقمية والابتكار', 'bullet': True, 'level': 0, 'font_size': 15},
            {'text': 'تسليم الواجب: تطبيق التفكير التصميمي على مشكلة رقمية واقعية', 'bullet': True, 'level': 0, 'font_size': 15},
            {'text': 'الاستعداد للاختبار البعدي', 'bullet': True, 'level': 0, 'font_size': 15},
            {'text': '', 'font_size': 8},
            {'text': 'المرجع: "عصر رقمي جديد" - يوسف العساف، 2020', 'font_size': 13, 'italic': True, 'color': '7F8C8D'},
        ],
        'notes': 'الشريحة الختامية. ذكّر المتعلمين بالمهام المطلوبة وموعد تسليم الواجب. شجعهم على مراجعة المحاضرة وتطبيق المفاهيم عمليا.'
    },
]

# ── Step 6: Apply content to slides ──
print(f"Applying content to {len(SLIDES_CONTENT)} slides...")

for slide_data in SLIDES_CONTENT:
    idx = slide_data['index']
    if idx >= len(prs.slides):
        print(f"  WARNING: Slide {idx} does not exist (only {len(prs.slides)} slides). Skipping.")
        continue

    slide = prs.slides[idx]
    print(f"  Processing slide {idx}: {slide_data['title'][:40]}...")

    # Find shapes
    text_shapes = get_all_text_shapes(slide)

    if not text_shapes:
        print(f"    WARNING: No text shapes found on slide {idx}")
        continue

    # Strategy: first shape is title (or top-most), rest is body
    title_shape = text_shapes[0]
    body_shapes = text_shapes[1:] if len(text_shapes) > 1 else []

    # Set title
    title_text = slide_data.get('title', '')
    if title_text and title_shape:
        set_text_in_shape(title_shape, [
            {'text': title_text, 'bold': True, 'font_size': 24, 'alignment': PP_ALIGN.RIGHT}
        ])

    # Set subtitle (for title slide)
    if 'subtitle' in slide_data and len(text_shapes) > 1:
        subtitle_shape = text_shapes[1]
        set_text_in_shape(subtitle_shape, [
            {'text': line, 'font_size': 16, 'alignment': PP_ALIGN.RIGHT}
            for line in slide_data['subtitle'].split('\n')
        ])

    # Set body content
    if 'body' in slide_data:
        # Find the best body shape (largest area, not the title)
        if body_shapes:
            body_shape = max(body_shapes, key=lambda s: s.width * s.height)
        elif len(text_shapes) > 1:
            body_shape = text_shapes[1]
        else:
            body_shape = text_shapes[0]  # Use title shape as fallback

        set_text_in_shape(body_shape, slide_data['body'])

    # Add speaker notes
    if 'notes' in slide_data:
        try:
            add_speaker_notes(slide, slide_data['notes'])
        except Exception as e:
            print(f"    Note: Could not add speaker notes: {e}")

# ── Step 7: Save ──
print(f"\nSaving presentation to: {OUTPUT}")
prs.save(OUTPUT)
print(f"SUCCESS! Presentation saved with {len(prs.slides)} slides.")
print(f"File: {OUTPUT}")
