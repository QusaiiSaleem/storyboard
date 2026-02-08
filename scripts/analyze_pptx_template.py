"""
Deep analysis of the PPTX template file.
Extracts: slide dimensions, layouts, shapes, positions, fonts, colors, backgrounds.
"""
import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

TEMPLATE_PATH = "/Users/qusaiabushanap/dev/storyboard/templates/قالب المحاضرة التفاعلية- عربي.pptx"

def emu_to_cm(emu):
    if emu is None:
        return None
    return round(emu / 914400 * 2.54, 2)

def emu_to_inches(emu):
    if emu is None:
        return None
    return round(emu / 914400, 2)

def emu_to_pt(emu):
    if emu is None:
        return None
    return round(emu / 12700, 1)

def get_color_info(color_format):
    """Extract color info from a ColorFormat object."""
    try:
        if color_format is None:
            return None
        info = {}
        try:
            if color_format.rgb is not None:
                info['rgb'] = str(color_format.rgb)
        except:
            pass
        try:
            if color_format.theme_color is not None:
                info['theme_color'] = str(color_format.theme_color)
        except:
            pass
        try:
            info['brightness'] = color_format.brightness
        except:
            pass
        return info if info else None
    except:
        return None

def get_fill_info(fill):
    """Extract fill information."""
    try:
        if fill is None:
            return None
        info = {'type': str(fill.type)}
        try:
            if fill.type is not None:
                info['type_name'] = fill.type.name if hasattr(fill.type, 'name') else str(fill.type)
        except:
            pass
        try:
            fore = fill.fore_color
            info['fore_color'] = get_color_info(fore)
        except:
            pass
        try:
            back = fill.back_color
            info['back_color'] = get_color_info(back)
        except:
            pass
        return info
    except:
        return None

def analyze_shape(shape):
    """Analyze a single shape in detail."""
    shape_info = {
        'shape_id': shape.shape_id,
        'name': shape.name,
        'shape_type': str(shape.shape_type) if shape.shape_type else None,
        'left_emu': shape.left,
        'top_emu': shape.top,
        'width_emu': shape.width,
        'height_emu': shape.height,
        'left_cm': emu_to_cm(shape.left),
        'top_cm': emu_to_cm(shape.top),
        'width_cm': emu_to_cm(shape.width),
        'height_cm': emu_to_cm(shape.height),
        'rotation': shape.rotation,
    }

    # Check for text frame
    if shape.has_text_frame:
        tf = shape.text_frame
        shape_info['text_frame'] = {
            'text': tf.text,
            'word_wrap': tf.word_wrap,
            'paragraphs': [],
        }
        try:
            shape_info['text_frame']['auto_size'] = str(tf.auto_size)
        except:
            pass
        try:
            shape_info['text_frame']['margin_left_cm'] = emu_to_cm(tf.margin_left)
            shape_info['text_frame']['margin_right_cm'] = emu_to_cm(tf.margin_right)
            shape_info['text_frame']['margin_top_cm'] = emu_to_cm(tf.margin_top)
            shape_info['text_frame']['margin_bottom_cm'] = emu_to_cm(tf.margin_bottom)
        except:
            pass

        for para in tf.paragraphs:
            para_info = {
                'text': para.text,
                'alignment': str(para.alignment) if para.alignment else None,
                'level': para.level,
                'runs': [],
            }

            # Check paragraph XML for RTL
            pPr = para._p.find(qn('a:pPr'))
            if pPr is not None:
                para_info['rtl'] = pPr.get('rtl')
                para_info['algn'] = pPr.get('algn')

                # Space before/after
                spcBef = pPr.find(qn('a:spcBef'))
                if spcBef is not None:
                    spcPts = spcBef.find(qn('a:spcPts'))
                    if spcPts is not None:
                        para_info['space_before_pt'] = int(spcPts.get('val')) / 100

                spcAft = pPr.find(qn('a:spcAft'))
                if spcAft is not None:
                    spcPts = spcAft.find(qn('a:spcPts'))
                    if spcPts is not None:
                        para_info['space_after_pt'] = int(spcPts.get('val')) / 100

                # Line spacing
                lnSpc = pPr.find(qn('a:lnSpc'))
                if lnSpc is not None:
                    spcPct = lnSpc.find(qn('a:spcPct'))
                    if spcPct is not None:
                        para_info['line_spacing_pct'] = int(spcPct.get('val')) / 1000

                # Default run properties
                defRPr = pPr.find(qn('a:defRPr'))
                if defRPr is not None:
                    para_info['default_font_size_pt'] = round(int(defRPr.get('sz', '0')) / 100, 1) if defRPr.get('sz') else None

            for run in para.runs:
                run_info = {
                    'text': run.text,
                    'font_name': run.font.name,
                    'font_size_pt': emu_to_pt(run.font.size) if run.font.size else None,
                    'bold': run.font.bold,
                    'italic': run.font.italic,
                    'underline': run.font.underline,
                }
                try:
                    if run.font.color and run.font.color.rgb:
                        run_info['font_color_rgb'] = str(run.font.color.rgb)
                except:
                    pass
                try:
                    if run.font.color and run.font.color.theme_color:
                        run_info['font_color_theme'] = str(run.font.color.theme_color)
                except:
                    pass

                # Check run XML for complex script font
                rPr = run._r.find(qn('a:rPr'))
                if rPr is not None:
                    run_info['lang'] = rPr.get('lang')
                    run_info['altLang'] = rPr.get('altLang')
                    run_info['dirty'] = rPr.get('dirty')

                    # CS font
                    cs = rPr.find(qn('a:cs'))
                    if cs is not None:
                        run_info['cs_font'] = cs.get('typeface')

                    # Latin font
                    latin = rPr.find(qn('a:latin'))
                    if latin is not None:
                        run_info['latin_font'] = latin.get('typeface')

                    # Arabic font
                    ea = rPr.find(qn('a:ea'))
                    if ea is not None:
                        run_info['ea_font'] = ea.get('typeface')

                para_info['runs'].append(run_info)

            shape_info['text_frame']['paragraphs'].append(para_info)

    # Check for table
    if shape.has_table:
        tbl = shape.table
        shape_info['table'] = {
            'row_count': len(tbl.rows),
            'col_count': len(tbl.columns),
            'rows': [],
        }
        # Column widths
        shape_info['table']['column_widths_cm'] = [emu_to_cm(col.width) for col in tbl.columns]

        for row_idx, row in enumerate(tbl.rows):
            row_info = {
                'row_index': row_idx,
                'height_cm': emu_to_cm(row.height),
                'cells': [],
            }
            for col_idx, cell in enumerate(row.cells):
                cell_info = {
                    'row': row_idx,
                    'col': col_idx,
                    'text': cell.text.strip(),
                    'paragraphs': [],
                }
                # Cell fill
                try:
                    tcPr = cell._tc.find(qn('a:tcPr'))
                    if tcPr is not None:
                        solidFill = tcPr.find(qn('a:solidFill'))
                        if solidFill is not None:
                            srgb = solidFill.find(qn('a:srgbClr'))
                            if srgb is not None:
                                cell_info['fill_color'] = srgb.get('val')
                            schemeClr = solidFill.find(qn('a:schemeClr'))
                            if schemeClr is not None:
                                cell_info['fill_scheme'] = schemeClr.get('val')
                except:
                    pass

                for para in cell.text_frame.paragraphs:
                    p_info = {
                        'text': para.text,
                        'alignment': str(para.alignment) if para.alignment else None,
                        'runs': [],
                    }
                    pPr = para._p.find(qn('a:pPr'))
                    if pPr is not None:
                        p_info['rtl'] = pPr.get('rtl')
                    for run in para.runs:
                        r_info = {
                            'text': run.text,
                            'font_name': run.font.name,
                            'font_size_pt': emu_to_pt(run.font.size) if run.font.size else None,
                            'bold': run.font.bold,
                        }
                        try:
                            if run.font.color and run.font.color.rgb:
                                r_info['font_color'] = str(run.font.color.rgb)
                        except:
                            pass
                        p_info['runs'].append(r_info)
                    cell_info['paragraphs'].append(p_info)

                row_info['cells'].append(cell_info)
            shape_info['table']['rows'].append(row_info)

    # Check for image
    if shape.shape_type is not None and 'PICTURE' in str(shape.shape_type):
        shape_info['is_image'] = True
        try:
            shape_info['image_content_type'] = shape.image.content_type
        except:
            pass

    # Fill info
    try:
        shape_info['fill'] = get_fill_info(shape.fill)
    except:
        pass

    # Line/border info
    try:
        if shape.line and shape.line.fill:
            shape_info['line'] = {
                'width_pt': emu_to_pt(shape.line.width) if shape.line.width else None,
            }
            try:
                if shape.line.color and shape.line.color.rgb:
                    shape_info['line']['color'] = str(shape.line.color.rgb)
            except:
                pass
    except:
        pass

    return shape_info

def analyze_slide_layout(layout):
    """Analyze a slide layout."""
    layout_info = {
        'name': layout.name,
        'shapes': [],
    }
    for shape in layout.placeholders:
        ph_info = {
            'idx': shape.placeholder_format.idx,
            'type': str(shape.placeholder_format.type),
            'name': shape.name,
            'left_cm': emu_to_cm(shape.left),
            'top_cm': emu_to_cm(shape.top),
            'width_cm': emu_to_cm(shape.width),
            'height_cm': emu_to_cm(shape.height),
        }
        if shape.has_text_frame:
            ph_info['text'] = shape.text_frame.text
        layout_info['shapes'].append(ph_info)
    return layout_info

def analyze_slide_master(master):
    """Analyze slide master."""
    master_info = {
        'layouts': [],
    }
    for layout in master.slide_layouts:
        master_info['layouts'].append({
            'name': layout.name,
            'placeholder_count': len(layout.placeholders),
        })

    # Master shapes
    master_info['shapes'] = []
    for shape in master.shapes:
        master_info['shapes'].append(analyze_shape(shape))

    return master_info

def main():
    prs = Presentation(TEMPLATE_PATH)
    analysis = {
        'filename': os.path.basename(TEMPLATE_PATH),
        'slide_width_emu': prs.slide_width,
        'slide_height_emu': prs.slide_height,
        'slide_width_cm': emu_to_cm(prs.slide_width),
        'slide_height_cm': emu_to_cm(prs.slide_height),
        'slide_count': len(prs.slides),
        'slides': [],
        'slide_masters': [],
        'slide_layouts': [],
    }

    print(f"Slide dimensions: {analysis['slide_width_cm']}cm x {analysis['slide_height_cm']}cm")
    print(f"Slide count: {analysis['slide_count']}")

    # Analyze slide masters
    for i, master in enumerate(prs.slide_masters):
        print(f"\nSlide Master {i}:")
        master_info = analyze_slide_master(master)
        analysis['slide_masters'].append(master_info)
        for layout in master_info['layouts']:
            print(f"  Layout: {layout['name']} ({layout['placeholder_count']} placeholders)")

    # Analyze all slide layouts
    for i, layout in enumerate(prs.slide_layouts):
        layout_info = analyze_slide_layout(layout)
        analysis['slide_layouts'].append(layout_info)

    # Analyze each slide
    for slide_idx, slide in enumerate(prs.slides):
        slide_info = {
            'slide_index': slide_idx,
            'slide_number': slide_idx + 1,
            'layout_name': slide.slide_layout.name if slide.slide_layout else None,
            'shapes': [],
            'notes': None,
        }

        # Background
        try:
            bg = slide.background
            if bg and bg.fill:
                slide_info['background'] = get_fill_info(bg.fill)
        except:
            pass

        print(f"\nSlide {slide_idx + 1} (Layout: {slide_info['layout_name']}):")

        for shape in slide.shapes:
            shape_info = analyze_shape(shape)
            slide_info['shapes'].append(shape_info)
            print(f"  Shape: {shape.name} | Type: {shape.shape_type} | "
                  f"Pos: ({emu_to_cm(shape.left)}, {emu_to_cm(shape.top)})cm | "
                  f"Size: {emu_to_cm(shape.width)}x{emu_to_cm(shape.height)}cm")
            if shape.has_text_frame:
                text = shape.text_frame.text[:50]
                if text:
                    print(f"    Text: {text}...")

        # Notes
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes:
                    slide_info['notes'] = notes
                    print(f"  Notes: {notes[:80]}...")
        except:
            pass

        analysis['slides'].append(slide_info)

    # Save full analysis
    output_path = "/Users/qusaiabushanap/dev/storyboard/docs/pptx_analysis_raw.json"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(analysis, f, ensure_ascii=False, indent=2, default=str)

    print(f"\n\nFull analysis saved to: {output_path}")
    return analysis

if __name__ == "__main__":
    main()
