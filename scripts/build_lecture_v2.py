#!/usr/bin/env python3
"""
NJR01_U02 Interactive Lecture - Built from scratch with high design quality
28 slides - الذهنية الرقمية وممارسات الابتكار التقني
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ─── DESIGN TOKENS ───────────────────────────────────────────
NAVY      = RGBColor(0x0A, 0x16, 0x28)   # Deep navy - headers
NAVY_MED  = RGBColor(0x12, 0x1E, 0x33)   # Medium navy
TEAL      = RGBColor(0x00, 0xB8, 0x94)   # Teal accent
ORANGE    = RGBColor(0xFF, 0x6B, 0x35)   # Orange accent
INDIGO    = RGBColor(0x63, 0x66, 0xF1)   # Indigo accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF5, 0xF6, 0xFA)   # Light background
DARK_TEXT = RGBColor(0x2D, 0x34, 0x36)   # Body text
MID_GRAY  = RGBColor(0x8B, 0x95, 0xA8)   # Muted text
LIGHT_LINE= RGBColor(0xE0, 0xE4, 0xEB)   # Subtle lines
GOLD      = RGBColor(0xF0, 0xC0, 0x40)   # Gold for stars/highlights

# Section colors
SEC_A = RGBColor(0x0A, 0x16, 0x28)  # Opening - Navy
SEC_B = RGBColor(0x00, 0x7B, 0x83)  # Foundations - Dark Teal
SEC_C = RGBColor(0x2D, 0x5F, 0xA0)  # Innovation Mindset - Blue
SEC_D = RGBColor(0x6C, 0x3A, 0xB2)  # Design Thinking - Purple
SEC_E = RGBColor(0xC0, 0x5A, 0x20)  # Market Fit - Burnt Orange
SEC_F = RGBColor(0x1A, 0x73, 0x5C)  # Entrepreneurial - Forest Green
SEC_G = RGBColor(0x0A, 0x16, 0x28)  # Closing - Navy

# Quiz colors
QUIZ_BG   = RGBColor(0x1B, 0x2A, 0x4A)
CORRECT   = RGBColor(0x00, 0xB8, 0x94)
WRONG     = RGBColor(0xE7, 0x4C, 0x3C)

# Slide dimensions (16:9 widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Font
FONT_AR = "Sakkal Majalla"
FONT_EN = "Calibri"

# ─── HELPER FUNCTIONS ────────────────────────────────────────

def set_rtl(paragraph):
    """Set paragraph to RTL"""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set(qn('w:bidi'), '1')
    # For PowerPoint RTL
    pPr.attrib['{http://schemas.openxmlformats.org/drawingml/2006/main}rtl'] = '1'

def set_rtl_simple(paragraph):
    """Simple RTL setting for pptx"""
    pPr = paragraph._pPr
    if pPr is None:
        pPr = paragraph._p.get_or_add_pPr()
    pPr.set('rtl', '1')

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                color=DARK_TEXT, alignment=PP_ALIGN.RIGHT, font_name=FONT_AR,
                line_spacing=1.5, rtl=True):
    """Add a styled textbox"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.line_spacing = line_spacing

    if rtl:
        set_rtl_simple(p)

    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=16,
                   color=DARK_TEXT, alignment=PP_ALIGN.RIGHT, line_spacing=1.5,
                   bullet=False, bold_first=False):
    """Add textbox with multiple lines"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if bullet:
            p.text = f"  {line}"
        else:
            p.text = line

        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_AR
        p.alignment = alignment
        p.line_spacing = line_spacing
        p.font.bold = (bold_first and i == 0)
        set_rtl_simple(p)

    return txBox

def add_rect(slide, left, top, width, height, fill_color, border=False, border_color=None):
    """Add a rectangle shape"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not border:
        shape.line.fill.background()
    elif border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_header_bar(slide, section_color=NAVY, title_text="", subtitle_text=""):
    """Add the top header bar with title"""
    # Full-width header bar
    bar = add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.3), section_color)

    # Accent line below header
    accent = add_rect(slide, Inches(0), Inches(1.3), SLIDE_W, Inches(0.06), TEAL)

    # Title text in header
    if title_text:
        add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.8),
                   title_text, font_size=32, bold=True, color=WHITE,
                   alignment=PP_ALIGN.RIGHT)

    # Subtitle
    if subtitle_text:
        add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
                   subtitle_text, font_size=14, color=RGBColor(0xA0,0xB0,0xC0),
                   alignment=PP_ALIGN.RIGHT)

def add_slide_number(slide, num, total=28):
    """Add slide number in bottom left"""
    add_textbox(slide, Inches(0.3), Inches(6.9), Inches(1.5), Inches(0.4),
               f"{num} / {total}", font_size=11, color=MID_GRAY,
               alignment=PP_ALIGN.LEFT, rtl=False)

def add_footer_line(slide):
    """Add a subtle footer line"""
    add_rect(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.02), LIGHT_LINE)

def add_institution_footer(slide):
    """Add institution name in footer"""
    add_textbox(slide, Inches(8), Inches(6.9), Inches(5), Inches(0.4),
               "جامعة نجران - كلية علوم الحاسب ونظم المعلومات",
               font_size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

def add_card(slide, left, top, width, height, title, body_lines,
             card_color=WHITE, title_color=NAVY, accent_color=TEAL, border=True):
    """Add a content card with title and body"""
    # Card background
    card = add_rounded_rect(slide, left, top, width, height, card_color,
                           border_color=LIGHT_LINE if border else None)

    # Accent stripe on right side of card
    add_rect(slide, left + width - Inches(0.08), top + Inches(0.15),
             Inches(0.08), height - Inches(0.3), accent_color)

    # Card title
    add_textbox(slide, left + Inches(0.2), top + Inches(0.15),
               width - Inches(0.5), Inches(0.5),
               title, font_size=20, bold=True, color=title_color,
               alignment=PP_ALIGN.RIGHT)

    # Card body
    if body_lines:
        add_multi_text(slide, left + Inches(0.2), top + Inches(0.65),
                      width - Inches(0.5), height - Inches(0.8),
                      body_lines, font_size=15, color=DARK_TEXT,
                      alignment=PP_ALIGN.RIGHT, line_spacing=1.6)

def add_notes(slide, text):
    """Add speaker notes"""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text

def add_number_badge(slide, left, top, number, color=TEAL):
    """Add a circular number badge"""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()

    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

def add_interaction_box(slide, left, top, width, height, interaction_type, instruction):
    """Add a Storyline interaction instruction box"""
    box = add_rounded_rect(slide, left, top, width, height,
                          RGBColor(0xFE, 0xF3, 0xE2),
                          border_color=ORANGE)

    # Icon/label
    add_textbox(slide, left + Inches(0.15), top + Inches(0.08),
               width - Inches(0.3), Inches(0.35),
               f"تعليمات Storyline: {interaction_type}",
               font_size=11, bold=True, color=ORANGE,
               alignment=PP_ALIGN.RIGHT)

    # Instruction text
    add_textbox(slide, left + Inches(0.15), top + Inches(0.4),
               width - Inches(0.3), height - Inches(0.5),
               instruction, font_size=10, color=RGBColor(0x8B, 0x60, 0x20),
               alignment=PP_ALIGN.RIGHT)

# ─── SLIDE BUILDERS ──────────────────────────────────────────

def build_title_slide(prs):
    """Slide 0: Title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Full dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Decorative geometric shapes
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, TEAL)
    add_rect(slide, Inches(0), Inches(5.8), SLIDE_W, Inches(0.04), TEAL)

    # Decorative corner accent
    add_rect(slide, Inches(11.5), Inches(0), Inches(1.8), Inches(0.04), ORANGE)
    add_rect(slide, Inches(13.15), Inches(0), Inches(0.04), Inches(1.2), ORANGE)

    # Unit label
    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.5),
               "الوحدة الثانية", font_size=20, color=TEAL,
               alignment=PP_ALIGN.RIGHT)

    # Main title
    add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.5),
               "الذهنية الرقمية وممارسات الابتكار التقني",
               font_size=44, bold=True, color=WHITE,
               alignment=PP_ALIGN.RIGHT)

    # Subtitle
    add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
               "Digital Mindset & Technology Innovation Practices",
               font_size=18, color=MID_GRAY, alignment=PP_ALIGN.RIGHT,
               font_name=FONT_EN)

    # Divider line
    add_rect(slide, Inches(8), Inches(4.6), Inches(4.3), Inches(0.03), TEAL)

    # Institution info
    add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.4),
               "جامعة نجران - كلية علوم الحاسب ونظم المعلومات",
               font_size=16, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

    # Start button area
    btn = add_rounded_rect(slide, Inches(5), Inches(6.0), Inches(3.3), Inches(0.7), TEAL)
    add_textbox(slide, Inches(5), Inches(6.05), Inches(3.3), Inches(0.6),
               "ابدأ المحاضرة", font_size=20, bold=True, color=WHITE,
               alignment=PP_ALIGN.CENTER)

    add_notes(slide, "شريحة العنوان - الذهنية الرقمية وممارسات الابتكار التقني\nStoryline: اضافة زر 'ابدأ المحاضرة' للانتقال للشريحة التالية")

def build_objectives_slide(prs):
    """Slide 1: Learning Objectives"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SEC_A, "الاهداف التعليمية", "ما ستتعلمه في هذه الوحدة")

    objectives = [
        ("تذكر", "ان يعدد المتعلم عناصر الابتكار الثلاثة وتقنيات العصف الذهني الاساسية", TEAL),
        ("فهم", "ان يوضح المتعلم مفهوم ريادة الاعمال وعلاقتها بالابتكار التقني", RGBColor(0x00,0x96,0x88)),
        ("فهم", "ان يشرح المتعلم المراحل الخمس للتفكير التصميمي ودور كل مرحلة", RGBColor(0x00,0x96,0x88)),
        ("تطبيق", "ان يطبق المتعلم طريقة SCAMPER لتوليد افكار ابداعية جديدة", RGBColor(0xF0,0xC0,0x40)),
        ("تطبيق", "ان يستخدم المتعلم مراحل التفكير التصميمي لبناء نموذج اولي", RGBColor(0xF0,0xC0,0x40)),
        ("تحليل", "ان يحلل المتعلم العلاقة بين التكنولوجيا وملاءمة السوق", ORANGE),
        ("تقييم", "ان يقيم المتعلم فرص العمل الريادية بناء على معايير محددة", RGBColor(0xE7,0x4C,0x3C)),
        ("ابداع", "ان يصمم المتعلم حلا ابتكاريا لمشكلة رقمية واقعية", INDIGO),
    ]

    y_start = Inches(1.6)
    for i, (level, obj, color) in enumerate(objectives):
        y = y_start + Inches(i * 0.62)

        # Level badge
        badge = add_rounded_rect(slide, Inches(10.5), y, Inches(1.3), Inches(0.45), color)
        add_textbox(slide, Inches(10.5), y + Inches(0.02), Inches(1.3), Inches(0.4),
                   level, font_size=13, bold=True, color=WHITE,
                   alignment=PP_ALIGN.CENTER)

        # Objective text
        add_textbox(slide, Inches(0.8), y + Inches(0.02), Inches(9.5), Inches(0.45),
                   obj, font_size=14, color=DARK_TEXT,
                   alignment=PP_ALIGN.RIGHT)

    add_footer_line(slide)
    add_slide_number(slide, 2)
    add_institution_footer(slide)
    add_notes(slide, "Storyline: عرض الاهداف بشكل متتابع مع تأثير ظهور تدريجي من اليمين")

def build_agenda_slide(prs):
    """Slide 2: Agenda"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SEC_A, "محاور المحاضرة")

    topics = [
        ("01", "ما هو الابتكار واهمية توليد الافكار"),
        ("02", "ريادة الاعمال وعقلية الابتكار التقني"),
        ("03", "التفكير التصميمي ومراحله الخمس"),
        ("04", "المرونة والتكيف وملاءمة السوق"),
        ("05", "بناء عقلية ريادية والتغلب على التحديات"),
    ]

    for i, (num, topic) in enumerate(topics):
        y = Inches(1.7) + Inches(i * 1.0)

        # Number circle
        add_number_badge(slide, Inches(11.2), y + Inches(0.05), num, TEAL)

        # Topic text
        add_textbox(slide, Inches(1), y + Inches(0.08), Inches(10), Inches(0.5),
                   topic, font_size=22, color=DARK_TEXT,
                   alignment=PP_ALIGN.RIGHT)

        # Subtle line
        if i < len(topics) - 1:
            add_rect(slide, Inches(1), y + Inches(0.7), Inches(11), Inches(0.01), LIGHT_LINE)

    add_footer_line(slide)
    add_slide_number(slide, 3)
    add_institution_footer(slide)

def build_innovation_slide(prs):
    """Slide 3: What is Innovation?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SEC_B, "ما هو الابتكار؟", "المحور الاول: الاسس")

    # Main definition
    defbox = add_rounded_rect(slide, Inches(2), Inches(1.6), Inches(9.5), Inches(1.0),
                              RGBColor(0xE8,0xF8,0xF5), border_color=TEAL)
    add_textbox(slide, Inches(2.3), Inches(1.7), Inches(9), Inches(0.8),
               "الابتكار هو ايجاد حلول جديدة للمشاكل الواقعية باستخدام تقنيات جديدة او اساليب ابداعية",
               font_size=20, bold=True, color=RGBColor(0x00,0x7B,0x83),
               alignment=PP_ALIGN.CENTER)

    # Three pillars
    cards_data = [
        ("استخدام تقنيات جديدة", "يتضمن الابتكار استخدام تقنيات\nجديدة او اساليب ابداعية\nلم تكن مستخدمة من قبل", INDIGO),
        ("تحسين الموجود", "الابتكار يقوم بتحسين العمليات\nاو المنتجات او الخدمات\nالحالية وتطويرها", ORANGE),
        ("حلول جديدة", "ايجاد حلول مبتكرة وجديدة\nللمشاكل الواقعية التي\nتواجه الافراد والمؤسسات", TEAL),
    ]

    for i, (title, body, color) in enumerate(cards_data):
        x = Inches(1.2) + Inches(i * 4.0)
        y = Inches(3.0)

        card = add_rounded_rect(slide, x, y, Inches(3.6), Inches(3.2), WHITE,
                               border_color=LIGHT_LINE)

        # Color accent bar at top of card
        add_rect(slide, x + Inches(0.8), y, Inches(2.0), Inches(0.06), color)

        # Number
        add_number_badge(slide, x + Inches(1.5), y + Inches(0.3), i+1, color)

        # Card title
        add_textbox(slide, x + Inches(0.2), y + Inches(1.0), Inches(3.2), Inches(0.5),
                   title, font_size=18, bold=True, color=color,
                   alignment=PP_ALIGN.CENTER)

        # Card body
        add_textbox(slide, x + Inches(0.2), y + Inches(1.5), Inches(3.2), Inches(1.5),
                   body, font_size=14, color=DARK_TEXT,
                   alignment=PP_ALIGN.CENTER, line_spacing=1.6)

    add_footer_line(slide)
    add_slide_number(slide, 4)
    add_institution_footer(slide)
    add_notes(slide, "Storyline: عرض البطاقات الثلاث بشكل متتابع عند النقر\nتعليمات: انقر على كل بطاقة لعرض التفاصيل")

def build_idea_generation_slide(prs):
    """Slide 4: Importance of Idea Generation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SEC_B, "اهمية توليد الافكار")

    # Key concept
    add_textbox(slide, Inches(1), Inches(1.6), Inches(11.3), Inches(0.6),
               "اساس الابتكار هو توليد افكار متنوعة تتطلب التفكير خارج الصندوق",
               font_size=20, bold=True, color=NAVY, alignment=PP_ALIGN.RIGHT)

    # Brainstorming techniques title
    add_textbox(slide, Inches(1), Inches(2.4), Inches(11.3), Inches(0.5),
               "تقنيات العصف الذهني:", font_size=18, bold=True, color=TEAL,
               alignment=PP_ALIGN.RIGHT)

    techniques = [
        ("التواصل الحر", "دع الافكار تتدفق دون اصدار احكام - اطلق العنان للتفكير الابداعي بدون قيود او نقد"),
        ("الخرائط الذهنية", "اربط الافكار والمفاهيم بصريا - استخدم الرسومات والتفرعات لتنظيم الافكار وربطها"),
        ("طريقة SCAMPER", "اداة فعالة لتوليد افكار جديدة بشكل منهجي من خلال سبع تقنيات منظمة"),
    ]

    for i, (title, desc) in enumerate(techniques):
        y = Inches(3.0) + Inches(i * 1.3)

        # Card
        card = add_rounded_rect(slide, Inches(1.5), y, Inches(10.5), Inches(1.1),
                               WHITE, border_color=LIGHT_LINE)

        # Accent
        add_rect(slide, Inches(11.92), y + Inches(0.15), Inches(0.08), Inches(0.8), TEAL)

        # Number
        add_number_badge(slide, Inches(11.1), y + Inches(0.25), i+1, TEAL)

        # Title
        add_textbox(slide, Inches(1.8), y + Inches(0.05), Inches(9), Inches(0.45),
                   title, font_size=18, bold=True, color=NAVY,
                   alignment=PP_ALIGN.RIGHT)

        # Description
        add_textbox(slide, Inches(1.8), y + Inches(0.5), Inches(9), Inches(0.5),
                   desc, font_size=13, color=MID_GRAY,
                   alignment=PP_ALIGN.RIGHT)

    add_footer_line(slide)
    add_slide_number(slide, 5)
    add_institution_footer(slide)

def build_scamper_slide(prs):
    """Slide 5: SCAMPER Method - Interactive"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SEC_B, "طريقة SCAMPER لتوليد الافكار", "تفاعلي: انقر على كل حرف لعرض التفاصيل")

    scamper = [
        ("S", "استبدال\nSubstitute", "ما الذي يمكن استبداله\nفي المنتج او الخدمة؟", RGBColor(0xE7,0x4C,0x3C)),
        ("C", "دمج\nCombine", "ما الذي يمكن دمجه\nمع عنصر آخر؟", RGBColor(0xF3,0x9C,0x12)),
        ("A", "تكيف\nAdapt", "ما الذي يمكن تكييفه\nمن مجال آخر؟", RGBColor(0x27,0xAE,0x60)),
        ("M", "تعديل\nModify", "ما الذي يمكن تعديله\nاو تكبيره او تصغيره؟", RGBColor(0x29,0x80,0xB9)),
        ("P", "استخدام آخر\nPut to use", "هل يمكن استخدامه\nفي مكان آخر؟", RGBColor(0x8E,0x44,0xAD)),
        ("E", "حذف\nEliminate", "ما الذي يمكن حذفه\nاو تبسيطه؟", RGBColor(0xD3,0x54,0x00)),
        ("R", "عكس\nReverse", "ماذا لو عكسنا\nالترتيب او الاتجاه؟", RGBColor(0x16,0xA0,0x85)),
    ]

    for i, (letter, ar_name, desc, color) in enumerate(scamper):
        x = Inches(0.6) + Inches(i * 1.8)
        y = Inches(1.7)

        # Letter circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.35), y,
                                        Inches(1.0), Inches(1.0))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = letter
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Arabic name
        add_textbox(slide, x, y + Inches(1.15), Inches(1.7), Inches(0.8),
                   ar_name, font_size=12, bold=True, color=color,
                   alignment=PP_ALIGN.CENTER, line_spacing=1.2)

        # Description
        add_textbox(slide, x, y + Inches(2.0), Inches(1.7), Inches(1.0),
                   desc, font_size=11, color=DARK_TEXT,
                   alignment=PP_ALIGN.CENTER, line_spacing=1.3)

    # Example box
    example_box = add_rounded_rect(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.2),
                                   RGBColor(0xFE,0xF9,0xE7), border_color=RGBColor(0xF0,0xC0,0x40))
    add_textbox(slide, Inches(1.8), Inches(5.05), Inches(9.8), Inches(0.35),
               "مثال تطبيقي: تطبيق SCAMPER على الهاتف الذكي",
               font_size=14, bold=True, color=RGBColor(0xB7,0x95,0x0B),
               alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(1.8), Inches(5.4), Inches(9.8), Inches(0.7),
               "استبدال: الازرار بشاشة لمس | دمج: هاتف + كاميرا + GPS | تكيف: تطبيقات من الحاسوب | حذف: لوحة المفاتيح الفعلية",
               font_size=12, color=RGBColor(0x7D,0x6B,0x0D),
               alignment=PP_ALIGN.RIGHT)

    # Interaction instruction
    add_interaction_box(slide, Inches(0.5), Inches(6.3), Inches(5), Inches(0.9),
                       "Click to Reveal",
                       "اخفاء جميع الحروف في البداية. عند النقر على كل حرف يظهر الاسم العربي والوصف مع تأثير انزلاق")

    add_slide_number(slide, 6)
    add_notes(slide, "Storyline: تفاعل Click-to-Reveal\n- اخفاء الاسماء والاوصاف في البداية\n- عند النقر على كل دائرة حرف، يظهر الاسم العربي والوصف\n- اضافة صوت نقر خفيف عند كل كشف\n- المثال التطبيقي يظهر بعد كشف جميع الحروف")

def build_quiz1_slide(prs):
    """Slide 6: Knowledge Check 1"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = QUIZ_BG

    # Header
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), RGBColor(0x0D,0x1B,0x36))
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
               "اختبر معلوماتك (1/5)", font_size=28, bold=True, color=TEAL,
               alignment=PP_ALIGN.RIGHT)
    add_rect(slide, Inches(0), Inches(1.0), SLIDE_W, Inches(0.04), TEAL)

    # Question
    add_textbox(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(0.8),
               "ما هو الحرف الذي يمثل 'الحذف' في طريقة SCAMPER؟",
               font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.RIGHT)

    # Options
    options = [
        ("أ", "S - Substitute", False),
        ("ب", "E - Eliminate", True),
        ("ج", "M - Modify", False),
        ("د", "R - Reverse", False),
    ]

    for i, (letter, text, correct) in enumerate(options):
        y = Inches(2.8) + Inches(i * 0.95)
        color = RGBColor(0x1A, 0x2E, 0x50)

        opt = add_rounded_rect(slide, Inches(2), y, Inches(9.3), Inches(0.75), color,
                              border_color=RGBColor(0x2A,0x4A,0x7A))

        # Letter badge
        badge = add_rounded_rect(slide, Inches(10.4), y + Inches(0.1), Inches(0.55), Inches(0.55),
                                TEAL if correct else RGBColor(0x2A,0x4A,0x7A))
        add_textbox(slide, Inches(10.4), y + Inches(0.12), Inches(0.55), Inches(0.5),
                   letter, font_size=16, bold=True, color=WHITE,
                   alignment=PP_ALIGN.CENTER)

        # Option text
        add_textbox(slide, Inches(2.3), y + Inches(0.12), Inches(7.8), Inches(0.5),
                   text, font_size=18, color=WHITE,
                   alignment=PP_ALIGN.RIGHT)

    # Correct answer indicator
    add_textbox(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.4),
               "الاجابة الصحيحة: ب", font_size=12, color=TEAL,
               alignment=PP_ALIGN.LEFT, rtl=False)

    add_slide_number(slide, 7)
    add_notes(slide, "Storyline: سؤال اختيار من متعدد\n- الاجابة الصحيحة: ب (E - Eliminate)\n- عند اختيار الاجابة الصحيحة: تلوين اخضر + رسالة 'احسنت!'\n- عند اختيار اجابة خاطئة: تلوين احمر + رسالة 'حاول مرة اخرى'\n- محاولتان مسموح بهما")

def build_entrepreneurship_slide(prs):
    """Slide 7: What is Entrepreneurship?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SEC_B, "ما هي ريادة الاعمال؟")

    # Definition
    defbox = add_rounded_rect(slide, Inches(2), Inches(1.6), Inches(9.5), Inches(0.9),
                              RGBColor(0xE8,0xF8,0xF5), border_color=TEAL)
    add_textbox(slide, Inches(2.3), Inches(1.7), Inches(9), Inches(0.7),
               "ريادة الاعمال هي تحويل الافكار الى مشاريع ناجحة من خلال تحديد الفرص والمخاطرة",
               font_size=18, bold=True, color=RGBColor(0x00,0x7B,0x83),
               alignment=PP_ALIGN.CENTER)

    # Three requirements
    reqs = [
        ("المرونة", "القدرة على التكيف مع\nالتغيرات والظروف\nالمختلفة بسرعة", RGBColor(0x27,0xAE,0x60)),
        ("الابداع", "التفكير بطرق غير تقليدية\nلايجاد حلول مبتكرة\nوفريدة من نوعها", ORANGE),
        ("حل المشكلات", "القدرة على تحديد المشكلات\nوتحليلها وايجاد\nحلول عملية فعالة", INDIGO),
    ]

    for i, (title, body, color) in enumerate(reqs):
        x = Inches(1.2) + Inches(i * 4.0)
        y = Inches(2.9)

        card = add_rounded_rect(slide, x, y, Inches(3.6), Inches(3.0), WHITE,
                               border_color=LIGHT_LINE)
        add_rect(slide, x + Inches(0.8), y, Inches(2.0), Inches(0.06), color)

        add_textbox(slide, x + Inches(0.2), y + Inches(0.4), Inches(3.2), Inches(0.5),
                   title, font_size=22, bold=True, color=color,
                   alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(0.2), y + Inches(1.1), Inches(3.2), Inches(1.5),
                   body, font_size=14, color=DARK_TEXT,
                   alignment=PP_ALIGN.CENTER, line_spacing=1.6)

    # Key insight
    add_textbox(slide, Inches(1), Inches(6.1), Inches(11.3), Inches(0.5),
               "رواد الاعمال يحددون الفرص ويخاطرون لتطوير منتجات او خدمات جديدة تحل مشكلات حقيقية",
               font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    add_footer_line(slide)
    add_slide_number(slide, 8)
    add_institution_footer(slide)

def build_innovation_mindset_slide(prs):
    """Slide 8: Tech Innovation Mindset"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SEC_C, "عقلية الابتكار التقني", "المحور الثاني")

    # Definition
    add_textbox(slide, Inches(1), Inches(1.6), Inches(11.3), Inches(0.8),
               "عقلية الابتكار التقني هي القدرة على تطوير حلول ابداعية لمشاكل واقعية باستخدام التكنولوجيا، ولا تقتصر على شركات التكنولوجيا بل تنطبق على اي مجال يشهد تحولا رقميا.",
               font_size=18, color=DARK_TEXT, alignment=PP_ALIGN.RIGHT)

    # Three pillars
    pillars = [
        ("التفكير التصميمي", "فهم احتياجات المستخدم\nقبل تطوير الحلول الرقمية", RGBColor(0x6C,0x3A,0xB2)),
        ("المرونة والقدرة على التكيف", "الاستجابة للتغيرات\nالتكنولوجية المتسارعة", RGBColor(0xC0,0x5A,0x20)),
        ("ملاءمة التكنولوجيا للسوق", "ضمان تلبية المنتجات\nلاحتياجات السوق الفعلية", RGBColor(0x1A,0x73,0x5C)),
    ]

    # Central concept with three branches
    # Center label
    center = add_rounded_rect(slide, Inches(5), Inches(2.8), Inches(3.3), Inches(0.8), NAVY)
    add_textbox(slide, Inches(5), Inches(2.85), Inches(3.3), Inches(0.7),
               "المتطلبات الرئيسية", font_size=18, bold=True, color=WHITE,
               alignment=PP_ALIGN.CENTER)

    positions = [
        (Inches(0.8), Inches(4.2)),
        (Inches(4.8), Inches(4.2)),
        (Inches(8.8), Inches(4.2)),
    ]

    for i, (title, body, color) in enumerate(pillars):
        x, y = positions[i]

        card = add_rounded_rect(slide, x, y, Inches(3.6), Inches(2.0), WHITE,
                               border_color=color)
        add_rect(slide, x, y, Inches(3.6), Inches(0.06), color)

        add_number_badge(slide, x + Inches(1.5), y + Inches(0.2), i+1, color)

        add_textbox(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.2), Inches(0.4),
                   title, font_size=16, bold=True, color=color,
                   alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(0.2), y + Inches(1.2), Inches(3.2), Inches(0.7),
                   body, font_size=12, color=DARK_TEXT,
                   alignment=PP_ALIGN.CENTER, line_spacing=1.4)

    add_footer_line(slide)
    add_slide_number(slide, 9)
    add_institution_footer(slide)

def build_quiz2_slide(prs):
    """Slide 9: Knowledge Check 2"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = QUIZ_BG

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), RGBColor(0x0D,0x1B,0x36))
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
               "اختبر معلوماتك (2/5)", font_size=28, bold=True, color=TEAL,
               alignment=PP_ALIGN.RIGHT)
    add_rect(slide, Inches(0), Inches(1.0), SLIDE_W, Inches(0.04), TEAL)

    add_textbox(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(0.8),
               "اي من التالية ليست من المتطلبات الرئيسية لعقلية الابتكار التقني؟",
               font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.RIGHT)

    options = [
        ("أ", "التفكير التصميمي", False),
        ("ب", "المرونة والقدرة على التكيف", False),
        ("ج", "البرمجة بلغة Python", True),
        ("د", "ملاءمة التكنولوجيا للسوق", False),
    ]

    for i, (letter, text, correct) in enumerate(options):
        y = Inches(2.8) + Inches(i * 0.95)
        color = RGBColor(0x1A, 0x2E, 0x50)

        add_rounded_rect(slide, Inches(2), y, Inches(9.3), Inches(0.75), color,
                        border_color=RGBColor(0x2A,0x4A,0x7A))

        badge = add_rounded_rect(slide, Inches(10.4), y + Inches(0.1), Inches(0.55), Inches(0.55),
                                TEAL if correct else RGBColor(0x2A,0x4A,0x7A))
        add_textbox(slide, Inches(10.4), y + Inches(0.12), Inches(0.55), Inches(0.5),
                   letter, font_size=16, bold=True, color=WHITE,
                   alignment=PP_ALIGN.CENTER)

        add_textbox(slide, Inches(2.3), y + Inches(0.12), Inches(7.8), Inches(0.5),
                   text, font_size=18, color=WHITE, alignment=PP_ALIGN.RIGHT)

    add_textbox(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.4),
               "الاجابة الصحيحة: ج", font_size=12, color=TEAL,
               alignment=PP_ALIGN.LEFT, rtl=False)

    add_slide_number(slide, 10)
    add_notes(slide, "Storyline: سؤال اختيار من متعدد\nالاجابة الصحيحة: ج - البرمجة بلغة Python ليست من المتطلبات الثلاثة")

def build_section_transition(prs):
    """Slide 10: Section Transition to Design Thinking"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = SEC_D

    # Decorative elements
    add_rect(slide, Inches(0), Inches(0), Inches(0.1), SLIDE_H, TEAL)
    add_rect(slide, Inches(0), Inches(4.8), SLIDE_W, Inches(0.04),
             RGBColor(0xFF, 0xFF, 0xFF))

    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.5),
               "المحور الثالث", font_size=18, color=RGBColor(0xBB,0xA0,0xE0),
               alignment=PP_ALIGN.RIGHT)

    add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.2),
               "التفكير التصميمي", font_size=52, bold=True, color=WHITE,
               alignment=PP_ALIGN.RIGHT)

    add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
               "Design Thinking", font_size=24, color=RGBColor(0xBB,0xA0,0xE0),
               alignment=PP_ALIGN.RIGHT, font_name=FONT_EN)

    add_textbox(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.8),
               "نهج ابتكاري متمحور حول الانسان يستخدمه شركات مثل جوجل وآبل",
               font_size=18, color=RGBColor(0xD0,0xC0,0xEE),
               alignment=PP_ALIGN.RIGHT)

    add_slide_number(slide, 11)

def build_design_thinking_overview(prs):
    """Slide 11: Design Thinking Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SEC_D, "ما هو التفكير التصميمي؟")

    points = [
        "نهج ابتكاري متمحور حول الانسان لحل المشكلات",
        "يركز على فهم احتياجات المستخدم قبل تطوير الحلول",
        "يستخدم النهج التكراري: بناء الافكار واختبارها وتحسينها",
        "يستخدم من قبل شركات عالمية رائدة مثل جوجل وآبل",
    ]

    for i, point in enumerate(points):
        y = Inches(1.8) + Inches(i * 1.1)

        card = add_rounded_rect(slide, Inches(1.5), y, Inches(10.5), Inches(0.85),
                               WHITE, border_color=LIGHT_LINE)

        add_rect(slide, Inches(11.92), y + Inches(0.1), Inches(0.08), Inches(0.65), SEC_D)
        add_number_badge(slide, Inches(11.1), y + Inches(0.15), i+1, SEC_D)

        add_textbox(slide, Inches(1.8), y + Inches(0.15), Inches(9), Inches(0.55),
                   point, font_size=18, color=DARK_TEXT,
                   alignment=PP_ALIGN.RIGHT)

    # Key insight box
    insight = add_rounded_rect(slide, Inches(2), Inches(6.0), Inches(9.3), Inches(0.7),
                               RGBColor(0xF3,0xE8,0xFF), border_color=SEC_D)
    add_textbox(slide, Inches(2.3), Inches(6.05), Inches(8.7), Inches(0.6),
               "التفكير التصميمي ليس مجرد منهجية بل هو عقلية تضع الانسان في مركز كل قرار تصميمي",
               font_size=14, bold=True, color=SEC_D, alignment=PP_ALIGN.CENTER)

    add_footer_line(slide)
    add_slide_number(slide, 12)
    add_institution_footer(slide)

def build_five_stages_slide(prs):
    """Slide 12: 5 Stages of Design Thinking - Interactive"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SEC_D, "المراحل الخمس للتفكير التصميمي", "تفاعلي: انقر على كل مرحلة")

    stages = [
        ("1", "التعاطف", "فهم احتياجات\nالمستخدم", RGBColor(0xE7,0x4C,0x3C)),
        ("2", "التحديد", "توضيح وتعريف\nالمشكلة", RGBColor(0xF3,0x9C,0x12)),
        ("3", "التفكير", "توليد حلول\nمتعددة", RGBColor(0x27,0xAE,0x60)),
        ("4", "النموذج الاولي", "بناء نماذج\nبسيطة", RGBColor(0x29,0x80,0xB9)),
        ("5", "الاختبار", "تقييم الحل\nوتحسينه", RGBColor(0x8E,0x44,0xAD)),
    ]

    # Draw stages as connected cards
    for i, (num, name, desc, color) in enumerate(stages):
        x = Inches(0.5) + Inches(i * 2.5)
        y = Inches(2.0)

        # Stage card
        card = add_rounded_rect(slide, x, y, Inches(2.2), Inches(3.0), WHITE,
                               border_color=color)
        add_rect(slide, x, y, Inches(2.2), Inches(0.08), color)

        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.75), y + Inches(0.3),
                                        Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Stage name
        add_textbox(slide, x + Inches(0.1), y + Inches(1.2), Inches(2.0), Inches(0.5),
                   name, font_size=20, bold=True, color=color,
                   alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, x + Inches(0.1), y + Inches(1.8), Inches(2.0), Inches(0.8),
                   desc, font_size=13, color=DARK_TEXT,
                   alignment=PP_ALIGN.CENTER, line_spacing=1.4)

        # Arrow between stages
        if i < 4:
            arrow_x = x + Inches(2.2)
            add_textbox(slide, arrow_x, y + Inches(1.2), Inches(0.3), Inches(0.5),
                       "←", font_size=24, color=MID_GRAY,
                       alignment=PP_ALIGN.CENTER, rtl=False, font_name=FONT_EN)

    # Interaction instruction
    add_interaction_box(slide, Inches(0.5), Inches(5.5), Inches(5.5), Inches(0.9),
                       "Sequential Reveal",
                       "عرض المراحل بشكل متتابع. عند النقر تظهر كل مرحلة مع تأثير انزلاق من الاسفل. السهم يظهر بين كل مرحلتين.")

    # Iterative note
    add_textbox(slide, Inches(6.5), Inches(5.6), Inches(6), Inches(0.7),
               "ملاحظة: هذه المراحل ليست خطية - يمكن العودة والتكرار بين المراحل حسب الحاجة",
               font_size=13, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

    add_footer_line(slide)
    add_slide_number(slide, 13)
    add_institution_footer(slide)
    add_notes(slide, "Storylline: عرض متتابع Sequential Reveal\n- المراحل تظهر واحدة تلو الاخرى\n- كل مرحلة تنزلق من الاسفل\n- السهم يظهر بين كل مرحلتين\n- في النهاية يظهر سهم دائري يدل على التكرارية")

def build_empathy_slide(prs):
    """Slide 13: Empathy & Needs Identification"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SEC_D, "التعاطف وتحديد الاحتياجات", "المرحلة الاولى من التفكير التصميمي")

    # Key message
    add_textbox(slide, Inches(1), Inches(1.6), Inches(11.3), Inches(0.6),
               "تحديد احتياجات المستخدم الحقيقية هو اساس الابتكار الهادف",
               font_size=20, bold=True, color=NAVY, alignment=PP_ALIGN.RIGHT)

    add_textbox(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(0.5),
               "التعاطف يساعد المصممين على فهم نقاط الضعف والاحتياجات الحقيقية للمستخدمين",
               font_size=16, color=DARK_TEXT, alignment=PP_ALIGN.RIGHT)

    # Methods
    add_textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(0.5),
               "اساليب تحديد الاحتياجات:", font_size=18, bold=True, color=SEC_D,
               alignment=PP_ALIGN.RIGHT)

    methods = [
        ("الملاحظة", "ملاحظة كيفية استخدام الناس\nللمنتجات في الحياة الواقعية", RGBColor(0xE7,0x4C,0x3C)),
        ("المقابلات", "المشاركة في مقابلات مع\nالمستخدمين للتعرف على\nوجهات نظرهم الشخصية", RGBColor(0x29,0x80,0xB9)),
        ("استطلاعات الرأي", "اجراء استطلاعات رأي\nلجمع آراء المستخدمين\nبشكل واسع", RGBColor(0x27,0xAE,0x60)),
    ]

    for i, (title, body, color) in enumerate(methods):
        x = Inches(1.2) + Inches(i * 4.0)
        y = Inches(3.6)

        card = add_rounded_rect(slide, x, y, Inches(3.6), Inches(2.5), WHITE,
                               border_color=LIGHT_LINE)
        add_rect(slide, x + Inches(0.8), y, Inches(2.0), Inches(0.06), color)
        add_number_badge(slide, x + Inches(1.5), y + Inches(0.25), i+1, color)

        add_textbox(slide, x + Inches(0.2), y + Inches(0.9), Inches(3.2), Inches(0.4),
                   title, font_size=18, bold=True, color=color,
                   alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), y + Inches(1.4), Inches(3.2), Inches(1.0),
                   body, font_size=13, color=DARK_TEXT,
                   alignment=PP_ALIGN.CENTER, line_spacing=1.5)

    add_footer_line(slide)
    add_slide_number(slide, 14)
    add_institution_footer(slide)

def build_define_ideate_prototype(prs):
    """Slide 14: Define, Ideate, Prototype"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SEC_D, "التحديد والتفكير والنمذجة", "المراحل 2 و 3 و 4")

    stages = [
        ("التحديد (عرف)", "توضيح وتعريف المشكلة بدقة بعد فهم\nاحتياجات المستخدم. صياغة المشكلة بشكل\nواضح يساعد في توجيه الحلول.", "2", SEC_D),
        ("التفكير (التصور)", "تبادل الافكار وتوليد حلول متعددة.\nالهدف هو الكم وليس الكيف في هذه\nالمرحلة - كلما زادت الافكار كان افضل.", "3", RGBColor(0x27,0xAE,0x60)),
        ("النموذج الاولي", "بناء نماذج بسيطة واختبارها.\nلا تحتاج لان تكون مثالية - الهدف هو\nاختبار الفكرة بسرعة وبأقل تكلفة.", "4", RGBColor(0x29,0x80,0xB9)),
    ]

    for i, (title, body, num, color) in enumerate(stages):
        y = Inches(1.6) + Inches(i * 1.8)

        card = add_rounded_rect(slide, Inches(1.5), y, Inches(10.5), Inches(1.5),
                               WHITE, border_color=LIGHT_LINE)

        # Left accent bar
        add_rect(slide, Inches(11.92), y + Inches(0.15), Inches(0.08), Inches(1.2), color)

        # Number badge
        add_number_badge(slide, Inches(11.1), y + Inches(0.15), num, color)

        # Title
        add_textbox(slide, Inches(1.8), y + Inches(0.1), Inches(9), Inches(0.45),
                   title, font_size=20, bold=True, color=color,
                   alignment=PP_ALIGN.RIGHT)

        # Body
        add_textbox(slide, Inches(1.8), y + Inches(0.55), Inches(9), Inches(0.9),
                   body, font_size=14, color=DARK_TEXT,
                   alignment=PP_ALIGN.RIGHT, line_spacing=1.5)

    add_footer_line(slide)
    add_slide_number(slide, 15)
    add_institution_footer(slide)

def build_testing_slide(prs):
    """Slide 15: Testing & Iteration"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SEC_D, "الاختبار والتكرار", "المرحلة الخامسة")

    add_textbox(slide, Inches(1), Inches(1.6), Inches(11.3), Inches(0.6),
               "المرحلة الاخيرة: تقييم الحل وتحسينه من خلال التجربة والتكرار المستمر",
               font_size=20, bold=True, color=NAVY, alignment=PP_ALIGN.RIGHT)

    # Cycle: Build -> Test -> Refine -> Repeat
    cycle_items = [
        ("بناء", "بناء النموذج\nالاولي", RGBColor(0x29,0x80,0xB9)),
        ("اختبار", "اختبار النموذج\nمع المستخدمين", RGBColor(0x27,0xAE,0x60)),
        ("صقل", "تحسين النموذج\nبناء على الملاحظات", RGBColor(0xF3,0x9C,0x12)),
        ("تكرار", "تكرار العملية\nحتى الوصول للحل الامثل", RGBColor(0x8E,0x44,0xAD)),
    ]

    for i, (title, desc, color) in enumerate(cycle_items):
        x = Inches(0.8) + Inches(i * 3.1)
        y = Inches(2.8)

        card = add_rounded_rect(slide, x, y, Inches(2.8), Inches(2.2), WHITE,
                               border_color=color)
        add_rect(slide, x + Inches(0.5), y, Inches(1.8), Inches(0.06), color)

        add_textbox(slide, x + Inches(0.1), y + Inches(0.3), Inches(2.6), Inches(0.5),
                   title, font_size=22, bold=True, color=color,
                   alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.9), Inches(2.6), Inches(0.8),
                   desc, font_size=13, color=DARK_TEXT,
                   alignment=PP_ALIGN.CENTER, line_spacing=1.4)

        # Arrow
        if i < 3:
            add_textbox(slide, x + Inches(2.8), y + Inches(0.8), Inches(0.3), Inches(0.5),
                       "←", font_size=20, color=MID_GRAY,
                       alignment=PP_ALIGN.CENTER, rtl=False, font_name=FONT_EN)

    # Key insight
    insight = add_rounded_rect(slide, Inches(2), Inches(5.5), Inches(9.3), Inches(1.0),
                               RGBColor(0xF3,0xE8,0xFF), border_color=SEC_D)
    add_textbox(slide, Inches(2.3), Inches(5.55), Inches(8.7), Inches(0.9),
               "التفكير التصميمي عملية تكرارية وليست خطية - يمكنك العودة لاي مرحلة سابقة في اي وقت بناء على نتائج الاختبار",
               font_size=15, color=SEC_D, alignment=PP_ALIGN.CENTER)

    add_footer_line(slide)
    add_slide_number(slide, 16)
    add_institution_footer(slide)

def build_quiz3_slide(prs):
    """Slide 16: Knowledge Check 3 - Drag and Drop"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = QUIZ_BG

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), RGBColor(0x0D,0x1B,0x36))
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
               "اختبر معلوماتك (3/5)", font_size=28, bold=True, color=TEAL,
               alignment=PP_ALIGN.RIGHT)
    add_rect(slide, Inches(0), Inches(1.0), SLIDE_W, Inches(0.04), TEAL)

    add_textbox(slide, Inches(1), Inches(1.4), Inches(11.3), Inches(0.8),
               "رتب المراحل الخمس للتفكير التصميمي بالترتيب الصحيح:",
               font_size=22, bold=True, color=WHITE, alignment=PP_ALIGN.RIGHT)

    # Scrambled stages for drag & drop
    stages_scrambled = ["النموذج الاولي", "التعاطف", "الاختبار", "التحديد", "التفكير"]

    for i, stage in enumerate(stages_scrambled):
        x = Inches(1.5) + Inches(i * 2.2)
        item = add_rounded_rect(slide, x, Inches(2.5), Inches(2.0), Inches(0.7),
                               RGBColor(0x1A,0x2E,0x50), border_color=RGBColor(0x3A,0x5A,0x8A))
        add_textbox(slide, x, Inches(2.55), Inches(2.0), Inches(0.6),
                   stage, font_size=16, bold=True, color=WHITE,
                   alignment=PP_ALIGN.CENTER)

    # Target slots
    for i in range(5):
        x = Inches(1.5) + Inches(i * 2.2)
        slot = add_rounded_rect(slide, x, Inches(4.0), Inches(2.0), Inches(0.7),
                               RGBColor(0x15,0x25,0x40), border_color=TEAL)
        add_textbox(slide, x, Inches(4.05), Inches(2.0), Inches(0.6),
                   str(i + 1), font_size=20, color=TEAL,
                   alignment=PP_ALIGN.CENTER)

    # Correct order
    add_textbox(slide, Inches(1), Inches(5.2), Inches(11.3), Inches(0.4),
               "الترتيب الصحيح: 1-التعاطف  2-التحديد  3-التفكير  4-النموذج الاولي  5-الاختبار",
               font_size=13, color=TEAL, alignment=PP_ALIGN.CENTER)

    add_interaction_box(slide, Inches(1), Inches(5.8), Inches(5.5), Inches(0.8),
                       "Drag & Drop",
                       "يسحب المتعلم كل مرحلة ويضعها في الخانة المناسبة. تلوين اخضر للصحيح واحمر للخطأ.")

    add_slide_number(slide, 17)
    add_notes(slide, "Storyline: Drag & Drop\n- المراحل تظهر بشكل عشوائي في الاعلى\n- المتعلم يسحب كل مرحلة للخانة المناسبة\n- الترتيب الصحيح: التعاطف > التحديد > التفكير > النموذج الاولي > الاختبار")

def build_flexibility_slide(prs):
    """Slide 17: Flexibility vs Adaptability"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SEC_E, "المرونة مقابل القدرة على التكيف", "المحور الرابع: تفاعلي")

    # Two columns comparison
    # Left card: Adaptability
    card1 = add_rounded_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.2),
                            WHITE, border_color=ORANGE)
    add_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.08), ORANGE)

    add_textbox(slide, Inches(1), Inches(2.1), Inches(5.1), Inches(0.5),
               "القدرة على التكيف", font_size=24, bold=True, color=ORANGE,
               alignment=PP_ALIGN.CENTER)

    adapt_points = [
        "القدرة على التطور والازدهار استجابة للتحولات الجذرية",
        "تتعامل مع التغيرات طويلة الامد في البيئة",
        "تتطلب اعادة تشكيل الاستراتيجيات والنماذج",
        "مثال: تحول نوكيا من الورق الى الاتصالات",
    ]
    for i, point in enumerate(adapt_points):
        add_textbox(slide, Inches(1.2), Inches(2.8) + Inches(i * 0.75), Inches(4.9), Inches(0.6),
                   f"• {point}", font_size=14, color=DARK_TEXT,
                   alignment=PP_ALIGN.RIGHT)

    # VS badge
    vs = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.1), Inches(3.4),
                                Inches(1.1), Inches(1.1))
    vs.fill.solid()
    vs.fill.fore_color.rgb = NAVY
    vs.line.fill.background()
    tf = vs.text_frame
    p = tf.paragraphs[0]
    p.text = "VS"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Right card: Flexibility
    card2 = add_rounded_rect(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(4.2),
                            WHITE, border_color=TEAL)
    add_rect(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(0.08), TEAL)

    add_textbox(slide, Inches(7.2), Inches(2.1), Inches(5.1), Inches(0.5),
               "المرونة", font_size=24, bold=True, color=TEAL,
               alignment=PP_ALIGN.CENTER)

    flex_points = [
        "القدرة على الاستجابة للتحديات او الفرص الفورية",
        "تتعامل مع التغيرات قصيرة المدى بكفاءة",
        "تتطلب سرعة في اتخاذ القرارات والتنفيذ",
        "مثال: تكيف المطاعم مع التوصيل اثناء كورونا",
    ]
    for i, point in enumerate(flex_points):
        add_textbox(slide, Inches(7.4), Inches(2.8) + Inches(i * 0.75), Inches(4.9), Inches(0.6),
                   f"• {point}", font_size=14, color=DARK_TEXT,
                   alignment=PP_ALIGN.RIGHT)

    add_interaction_box(slide, Inches(0.8), Inches(6.2), Inches(5.5), Inches(0.8),
                       "Comparison Reveal",
                       "اظهار كل عمود عند النقر عليه. البطاقة اليمنى تظهر اولا ثم اليسرى مع تأثير انزلاق")

    add_footer_line(slide)
    add_slide_number(slide, 18)
    add_institution_footer(slide)

def build_opportunities_slide(prs):
    """Slide 18: Identifying Business Opportunities"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SEC_E, "تحديد فرص العمل")

    add_textbox(slide, Inches(1), Inches(1.6), Inches(11.3), Inches(0.6),
               "الفكرة يجب ان تحظى بطلب في السوق لتكون مشروعا قابلا للاستمرار",
               font_size=20, bold=True, color=NAVY, alignment=PP_ALIGN.RIGHT)

    add_textbox(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(0.4),
               "فرصة العمل الجيدة تحل مشكلة حقيقية يواجهها الناس",
               font_size=16, color=DARK_TEXT, alignment=PP_ALIGN.RIGHT)

    factors = [
        ("احتياجات العملاء", "فهم ما يحتاجه العملاء\nفعلا وليس ما نظن\nانهم يحتاجونه", TEAL),
        ("دراسة الجدوى", "التأكد من ان المشروع\nقابل للتنفيذ تقنيا\nوماليا وعمليا", ORANGE),
        ("قابلية التوسع", "امكانية نمو المشروع\nوتوسعه للوصول الى\nشريحة اكبر من العملاء", INDIGO),
    ]

    for i, (title, body, color) in enumerate(factors):
        x = Inches(1.2) + Inches(i * 4.0)
        y = Inches(3.2)

        card = add_rounded_rect(slide, x, y, Inches(3.6), Inches(3.0), WHITE,
                               border_color=LIGHT_LINE)
        add_rect(slide, x + Inches(0.8), y, Inches(2.0), Inches(0.06), color)
        add_number_badge(slide, x + Inches(1.5), y + Inches(0.3), i+1, color)

        add_textbox(slide, x + Inches(0.2), y + Inches(1.0), Inches(3.2), Inches(0.5),
                   title, font_size=18, bold=True, color=color,
                   alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), y + Inches(1.5), Inches(3.2), Inches(1.2),
                   body, font_size=13, color=DARK_TEXT,
                   alignment=PP_ALIGN.CENTER, line_spacing=1.5)

    add_footer_line(slide)
    add_slide_number(slide, 19)
    add_institution_footer(slide)

def build_market_fit_slide(prs):
    """Slide 19: Technology & Market Fit"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SEC_E, "التكنولوجيا وملاءمة السوق")

    # Definition
    defbox = add_rounded_rect(slide, Inches(2), Inches(1.6), Inches(9.5), Inches(0.9),
                              RGBColor(0xFE,0xF3,0xE2), border_color=ORANGE)
    add_textbox(slide, Inches(2.3), Inches(1.65), Inches(9), Inches(0.8),
               "ملاءمة السوق = ضمان تلبية المنتجات الرقمية لاحتياجات السوق الفعلية",
               font_size=18, bold=True, color=RGBColor(0xC0,0x5A,0x20),
               alignment=PP_ALIGN.CENTER)

    points = [
        "التكنولوجيا يجب ان تتوافق مع احتياجات المستخدم لتحقيق النجاح",
        "ملاءمة السوق ضرورية للتبني وقابلية التوسع",
        "النهج التكراري يضمن مواءمة التكنولوجيا للسوق واحتياجات المستخدمين",
    ]

    for i, point in enumerate(points):
        y = Inches(2.9) + Inches(i * 0.8)
        add_textbox(slide, Inches(1.5), y, Inches(10.5), Inches(0.6),
                   f"• {point}", font_size=16, color=DARK_TEXT,
                   alignment=PP_ALIGN.RIGHT)

    # Understanding market needs
    add_textbox(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.4),
               "فهم احتياجات السوق:", font_size=18, bold=True, color=SEC_E,
               alignment=PP_ALIGN.RIGHT)

    needs = ["الشركات الناجحة تعالج مشاكل العملاء الفعلية",
             "ابحاث السوق تساعد في تحديد الطلب والمنافسة",
             "ملاحظات العملاء ضرورية لتحسين الحلول"]

    for i, need in enumerate(needs):
        add_textbox(slide, Inches(1.5), Inches(5.5) + Inches(i * 0.45), Inches(10.5), Inches(0.4),
                   f"  {need}", font_size=14, color=DARK_TEXT,
                   alignment=PP_ALIGN.RIGHT)

    add_footer_line(slide)
    add_slide_number(slide, 20)
    add_institution_footer(slide)

def build_prototyping_slide(prs):
    """Slide 20: Prototyping & A/B Testing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SEC_E, "النمذجة والتجريب واختبار A/B")

    # Prototype types
    add_textbox(slide, Inches(1), Inches(1.6), Inches(11.3), Inches(0.5),
               "انواع النماذج الاولية:", font_size=20, bold=True, color=NAVY,
               alignment=PP_ALIGN.RIGHT)

    types = [
        ("نماذج ورقية", "رسومات تخطيطية\nسريعة على الورق", RGBColor(0xE7,0x4C,0x3C)),
        ("نماذج رقمية", "نماذج برمجية تفاعلية\nباستخدام ادوات التصميم", RGBColor(0x29,0x80,0xB9)),
        ("نماذج مادية", "نماذج منتجات اولية\nملموسة للاختبار", RGBColor(0x27,0xAE,0x60)),
    ]

    for i, (title, desc, color) in enumerate(types):
        x = Inches(1.2) + Inches(i * 4.0)
        y = Inches(2.3)

        card = add_rounded_rect(slide, x, y, Inches(3.6), Inches(1.8), WHITE,
                               border_color=LIGHT_LINE)
        add_rect(slide, x + Inches(0.8), y, Inches(2.0), Inches(0.06), color)

        add_textbox(slide, x + Inches(0.2), y + Inches(0.3), Inches(3.2), Inches(0.4),
                   title, font_size=18, bold=True, color=color,
                   alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.2), Inches(0.8),
                   desc, font_size=13, color=DARK_TEXT,
                   alignment=PP_ALIGN.CENTER, line_spacing=1.4)

    # A/B Testing section
    add_textbox(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.5),
               "اختبار A/B:", font_size=20, bold=True, color=NAVY,
               alignment=PP_ALIGN.RIGHT)

    ab_box = add_rounded_rect(slide, Inches(1.5), Inches(5.1), Inches(10.3), Inches(1.3),
                              RGBColor(0xE8,0xF8,0xF5), border_color=TEAL)

    add_textbox(slide, Inches(1.8), Inches(5.2), Inches(9.7), Inches(1.0),
               "اختبار A/B يقارن بين فكرتين لمعرفة الانسب. يتم عرض نسختين مختلفتين من المنتج على مجموعتين من المستخدمين ثم قياس اي النسختين تحقق نتائج افضل. التكرار المستمر يؤدي الى منتجات افضل.",
               font_size=15, color=RGBColor(0x00,0x7B,0x83),
               alignment=PP_ALIGN.RIGHT, line_spacing=1.6)

    add_footer_line(slide)
    add_slide_number(slide, 21)
    add_institution_footer(slide)

def build_quiz4_slide(prs):
    """Slide 21: Knowledge Check 4"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = QUIZ_BG

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), RGBColor(0x0D,0x1B,0x36))
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
               "اختبر معلوماتك (4/5)", font_size=28, bold=True, color=TEAL,
               alignment=PP_ALIGN.RIGHT)
    add_rect(slide, Inches(0), Inches(1.0), SLIDE_W, Inches(0.04), TEAL)

    add_textbox(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(0.8),
               "ما هو الهدف الرئيسي من اختبار A/B؟",
               font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.RIGHT)

    options = [
        ("أ", "اختبار البرمجة فقط", False),
        ("ب", "المقارنة بين فكرتين لمعرفة الانسب", True),
        ("ج", "حذف المنتجات القديمة", False),
        ("د", "تصميم الشعارات", False),
    ]

    for i, (letter, text, correct) in enumerate(options):
        y = Inches(2.8) + Inches(i * 0.95)

        add_rounded_rect(slide, Inches(2), y, Inches(9.3), Inches(0.75),
                        RGBColor(0x1A,0x2E,0x50), border_color=RGBColor(0x2A,0x4A,0x7A))

        add_rounded_rect(slide, Inches(10.4), y + Inches(0.1), Inches(0.55), Inches(0.55),
                        TEAL if correct else RGBColor(0x2A,0x4A,0x7A))
        add_textbox(slide, Inches(10.4), y + Inches(0.12), Inches(0.55), Inches(0.5),
                   letter, font_size=16, bold=True, color=WHITE,
                   alignment=PP_ALIGN.CENTER)

        add_textbox(slide, Inches(2.3), y + Inches(0.12), Inches(7.8), Inches(0.5),
                   text, font_size=18, color=WHITE, alignment=PP_ALIGN.RIGHT)

    add_textbox(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.4),
               "الاجابة الصحيحة: ب", font_size=12, color=TEAL,
               alignment=PP_ALIGN.LEFT, rtl=False)

    add_slide_number(slide, 22)

def build_entrepreneurial_mindset(prs):
    """Slide 22: Building Entrepreneurial Mindset"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SEC_F, "بناء عقلية ريادية", "المحور الخامس")

    add_textbox(slide, Inches(1), Inches(1.6), Inches(11.3), Inches(0.5),
               "كيفية بناء عقلية ريادية؟", font_size=22, bold=True, color=NAVY,
               alignment=PP_ALIGN.RIGHT)

    mindset = [
        ("رؤية الفشل كفرصة للتعلم", "الفشل ليس نهاية الطريق بل هو درس قيم يقربك من النجاح. كل محاولة فاشلة تكشف لك ما لا يعمل وتوجهك نحو ما يعمل.", RGBColor(0xE7,0x4C,0x3C)),
        ("الانفتاح على التغذية الراجعة", "استمع لآراء الآخرين وتقبل النقد البناء. التغذية الراجعة هي مرآتك التي تكشف نقاط القوة والضعف.", RGBColor(0x29,0x80,0xB9)),
        ("البقاء فضوليا والتحسين المستمر", "لا تتوقف عن التعلم والتجربة. الفضول هو محرك الابتكار والتحسين المستمر هو وقوده.", RGBColor(0x27,0xAE,0x60)),
    ]

    for i, (title, body, color) in enumerate(mindset):
        y = Inches(2.3) + Inches(i * 1.5)

        card = add_rounded_rect(slide, Inches(1.5), y, Inches(10.5), Inches(1.3),
                               WHITE, border_color=LIGHT_LINE)
        add_rect(slide, Inches(11.92), y + Inches(0.15), Inches(0.08), Inches(1.0), color)
        add_number_badge(slide, Inches(11.1), y + Inches(0.15), i+1, color)

        add_textbox(slide, Inches(1.8), y + Inches(0.1), Inches(9), Inches(0.4),
                   title, font_size=18, bold=True, color=color,
                   alignment=PP_ALIGN.RIGHT)
        add_textbox(slide, Inches(1.8), y + Inches(0.55), Inches(9), Inches(0.6),
                   body, font_size=13, color=DARK_TEXT,
                   alignment=PP_ALIGN.RIGHT, line_spacing=1.5)

    add_footer_line(slide)
    add_slide_number(slide, 23)
    add_institution_footer(slide)

def build_tools_slide(prs):
    """Slide 23: Tools & Resources"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SEC_F, "ادوات ومصادر لرائد الاعمال")

    tools = [
        ("Lean Startup", "منهجية الشركات الناشئة المرنة", "منهجية للتطوير السريع تركز على\nبناء منتج بالحد الادنى من المميزات\n(MVP) واختباره في السوق بسرعة\nثم التعلم والتحسين بشكل متكرر", RGBColor(0xE7,0x4C,0x3C)),
        ("Business Model Canvas", "نموذج العمل التجاري", "اداة بصرية لتنظيم وهيكلة الافكار\nالتجارية في صفحة واحدة تشمل:\nالعملاء، القيمة المقدمة، القنوات،\nمصادر الايرادات، والتكاليف", RGBColor(0x29,0x80,0xB9)),
        ("التعلم المستمر", "منصات ومصادر التعلم", "منصات عبر الانترنت مثل\nStanford Online و Coursera\nو edX توفر دورات متخصصة\nفي ريادة الاعمال والابتكار", RGBColor(0x27,0xAE,0x60)),
    ]

    for i, (en_title, ar_title, desc, color) in enumerate(tools):
        x = Inches(0.8) + Inches(i * 4.1)
        y = Inches(1.6)

        card = add_rounded_rect(slide, x, y, Inches(3.8), Inches(4.5), WHITE,
                               border_color=LIGHT_LINE)
        add_rect(slide, x, y, Inches(3.8), Inches(0.08), color)

        add_textbox(slide, x + Inches(0.2), y + Inches(0.3), Inches(3.4), Inches(0.4),
                   en_title, font_size=20, bold=True, color=color,
                   alignment=PP_ALIGN.CENTER, font_name=FONT_EN)

        add_textbox(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.4), Inches(0.4),
                   ar_title, font_size=16, bold=True, color=NAVY,
                   alignment=PP_ALIGN.CENTER)

        add_rect(slide, x + Inches(0.8), y + Inches(1.3), Inches(2.2), Inches(0.02), LIGHT_LINE)

        add_textbox(slide, x + Inches(0.2), y + Inches(1.5), Inches(3.4), Inches(2.5),
                   desc, font_size=13, color=DARK_TEXT,
                   alignment=PP_ALIGN.CENTER, line_spacing=1.5)

    # Collaboration section
    collab = add_rounded_rect(slide, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.8),
                              RGBColor(0xE8,0xF5,0xF0), border_color=SEC_F)
    add_textbox(slide, Inches(1.8), Inches(6.35), Inches(9.7), Inches(0.6),
               "التعاون وبناء العلاقات: رواد الاعمال الناجحون يستفيدون من شبكات العلاقات وحاضنات الاعمال لتسريع النمو",
               font_size=14, color=SEC_F, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 24)

def build_challenges_slide(prs):
    """Slide 24: Overcoming Innovation Challenges"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SEC_F, "التغلب على تحديات الابتكار")

    # Challenges
    add_textbox(slide, Inches(1), Inches(1.6), Inches(11.3), Inches(0.4),
               "التحديات:", font_size=20, bold=True, color=RGBColor(0xE7,0x4C,0x3C),
               alignment=PP_ALIGN.RIGHT)

    challenges = [
        ("مقاومة التغيير", "الناس بطبيعتهم يميلون للمألوف ويقاومون التغيير خوفا من المجهول"),
        ("الخوف من الفشل", "الخوف من الفشل يمكن ان يعيق الابداع ويمنع تجربة افكار جديدة"),
    ]

    for i, (title, desc) in enumerate(challenges):
        y = Inches(2.1) + Inches(i * 0.9)
        card = add_rounded_rect(slide, Inches(1.5), y, Inches(10.5), Inches(0.75),
                               RGBColor(0xFD,0xED,0xED), border_color=RGBColor(0xE7,0x4C,0x3C))
        add_textbox(slide, Inches(1.8), y + Inches(0.05), Inches(10), Inches(0.3),
                   title, font_size=16, bold=True, color=RGBColor(0xE7,0x4C,0x3C),
                   alignment=PP_ALIGN.RIGHT)
        add_textbox(slide, Inches(1.8), y + Inches(0.35), Inches(10), Inches(0.35),
                   desc, font_size=13, color=DARK_TEXT, alignment=PP_ALIGN.RIGHT)

    # Solutions
    add_textbox(slide, Inches(1), Inches(4.1), Inches(11.3), Inches(0.4),
               "الحلول:", font_size=20, bold=True, color=TEAL,
               alignment=PP_ALIGN.RIGHT)

    solutions = [
        ("تشجيع التجربة", "خلق بيئة آمنة لتجربة اشياء جديدة بدون خوف من العقوبة على الفشل"),
        ("تعزيز ثقافة الانفتاح", "تشجيع التواصل المفتوح للتعلم من الاخطاء ومشاركة الافكار بحرية"),
        ("الاحتفاء بالتعلم من الفشل", "تحويل كل تجربة فاشلة الى درس مستفاد ومشاركته مع الفريق"),
    ]

    for i, (title, desc) in enumerate(solutions):
        y = Inches(4.6) + Inches(i * 0.75)
        card = add_rounded_rect(slide, Inches(1.5), y, Inches(10.5), Inches(0.65),
                               RGBColor(0xE8,0xF8,0xF5), border_color=TEAL)
        add_textbox(slide, Inches(1.8), y + Inches(0.03), Inches(10), Inches(0.3),
                   title, font_size=15, bold=True, color=TEAL,
                   alignment=PP_ALIGN.RIGHT)
        add_textbox(slide, Inches(1.8), y + Inches(0.3), Inches(10), Inches(0.3),
                   desc, font_size=12, color=DARK_TEXT, alignment=PP_ALIGN.RIGHT)

    add_footer_line(slide)
    add_slide_number(slide, 25)
    add_institution_footer(slide)

def build_quiz5_slide(prs):
    """Slide 25: Knowledge Check 5"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = QUIZ_BG

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), RGBColor(0x0D,0x1B,0x36))
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
               "اختبر معلوماتك (5/5)", font_size=28, bold=True, color=TEAL,
               alignment=PP_ALIGN.RIGHT)
    add_rect(slide, Inches(0), Inches(1.0), SLIDE_W, Inches(0.04), TEAL)

    add_textbox(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(0.8),
               "اي من التالية يعتبر من سمات العقلية الريادية؟",
               font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.RIGHT)

    options = [
        ("أ", "الخوف من المخاطرة وتجنب التجارب الجديدة", False),
        ("ب", "رؤية الفشل كفرصة للتعلم والتحسين المستمر", True),
        ("ج", "رفض التغذية الراجعة من الآخرين", False),
        ("د", "التمسك بالخطة الاصلية دون اي تعديل", False),
    ]

    for i, (letter, text, correct) in enumerate(options):
        y = Inches(2.8) + Inches(i * 0.95)

        add_rounded_rect(slide, Inches(2), y, Inches(9.3), Inches(0.75),
                        RGBColor(0x1A,0x2E,0x50), border_color=RGBColor(0x2A,0x4A,0x7A))

        add_rounded_rect(slide, Inches(10.4), y + Inches(0.1), Inches(0.55), Inches(0.55),
                        TEAL if correct else RGBColor(0x2A,0x4A,0x7A))
        add_textbox(slide, Inches(10.4), y + Inches(0.12), Inches(0.55), Inches(0.5),
                   letter, font_size=16, bold=True, color=WHITE,
                   alignment=PP_ALIGN.CENTER)

        add_textbox(slide, Inches(2.3), y + Inches(0.12), Inches(7.8), Inches(0.5),
                   text, font_size=18, color=WHITE, alignment=PP_ALIGN.RIGHT)

    add_textbox(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.4),
               "الاجابة الصحيحة: ب", font_size=12, color=TEAL,
               alignment=PP_ALIGN.LEFT, rtl=False)

    add_slide_number(slide, 26)

def build_summary_slide(prs):
    """Slide 26: Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, SEC_G, "ملخص الوحدة", "ما تعلمناه اليوم")

    summary_items = [
        ("الابتكار", "ايجاد حلول جديدة باستخدام تقنيات ابداعية مثل SCAMPER", TEAL),
        ("ريادة الاعمال", "تحويل الافكار لمشاريع ناجحة تتطلب حل المشكلات والمرونة", RGBColor(0x29,0x80,0xB9)),
        ("التفكير التصميمي", "5 مراحل: التعاطف، التحديد، التفكير، النمذجة، الاختبار", SEC_D),
        ("ملاءمة السوق", "التكنولوجيا + احتياجات السوق + النمذجة + اختبار A/B", ORANGE),
        ("العقلية الريادية", "الفشل فرصة، التغذية الراجعة ضرورية، التعلم مستمر", SEC_F),
    ]

    for i, (title, desc, color) in enumerate(summary_items):
        y = Inches(1.6) + Inches(i * 1.0)

        card = add_rounded_rect(slide, Inches(1.5), y, Inches(10.5), Inches(0.85),
                               WHITE, border_color=LIGHT_LINE)
        add_rect(slide, Inches(11.92), y + Inches(0.1), Inches(0.08), Inches(0.65), color)

        # Badge
        badge = add_rounded_rect(slide, Inches(10.3), y + Inches(0.12), Inches(1.5), Inches(0.6), color)
        add_textbox(slide, Inches(10.3), y + Inches(0.15), Inches(1.5), Inches(0.55),
                   title, font_size=13, bold=True, color=WHITE,
                   alignment=PP_ALIGN.CENTER)

        add_textbox(slide, Inches(1.8), y + Inches(0.15), Inches(8.3), Inches(0.55),
                   desc, font_size=16, color=DARK_TEXT,
                   alignment=PP_ALIGN.RIGHT)

    add_footer_line(slide)
    add_slide_number(slide, 27)
    add_institution_footer(slide)

def build_closing_slide(prs):
    """Slide 27: Next Steps / Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY

    # Decorative elements
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, TEAL)
    add_rect(slide, Inches(0), Inches(5.8), SLIDE_W, Inches(0.04), TEAL)
    add_rect(slide, Inches(11.5), Inches(0), Inches(1.8), Inches(0.04), ORANGE)
    add_rect(slide, Inches(13.15), Inches(0), Inches(0.04), Inches(1.2), ORANGE)

    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.5),
               "الخطوات القادمة", font_size=20, color=TEAL,
               alignment=PP_ALIGN.RIGHT)

    add_textbox(slide, Inches(1), Inches(2.3), Inches(11), Inches(1.0),
               "اكملت المحاضرة بنجاح!", font_size=36, bold=True, color=WHITE,
               alignment=PP_ALIGN.RIGHT)

    next_steps = [
        "اكمال الانشطة التفاعلية المرتبطة بالوحدة",
        "المشاركة في نشاط النقاش حول العقلية الرقمية",
        "تطبيق التفكير التصميمي في حل مشكلة رقمية (الواجب)",
        "الاستعداد للاختبار البعدي",
    ]

    for i, step in enumerate(next_steps):
        y = Inches(3.5) + Inches(i * 0.55)
        add_textbox(slide, Inches(2), y, Inches(9.5), Inches(0.5),
                   f"  {step}", font_size=18, color=RGBColor(0xC0,0xD0,0xE0),
                   alignment=PP_ALIGN.RIGHT)

    # Thank you
    add_textbox(slide, Inches(1), Inches(5.9), Inches(11), Inches(0.6),
               "جامعة نجران - كلية علوم الحاسب ونظم المعلومات",
               font_size=14, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

    add_slide_number(slide, 28)

# ─── MAIN ────────────────────────────────────────────────────

def main():
    print("Building interactive lecture from scratch...")
    print("Design: Dark navy + teal accent + warm orange")
    print()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Build all 28 slides
    builders = [
        ("0: Title", build_title_slide),
        ("1: Objectives", build_objectives_slide),
        ("2: Agenda", build_agenda_slide),
        ("3: What is Innovation", build_innovation_slide),
        ("4: Idea Generation", build_idea_generation_slide),
        ("5: SCAMPER", build_scamper_slide),
        ("6: Quiz 1", build_quiz1_slide),
        ("7: Entrepreneurship", build_entrepreneurship_slide),
        ("8: Innovation Mindset", build_innovation_mindset_slide),
        ("9: Quiz 2", build_quiz2_slide),
        ("10: Section Transition", build_section_transition),
        ("11: Design Thinking Overview", build_design_thinking_overview),
        ("12: 5 Stages", build_five_stages_slide),
        ("13: Empathy", build_empathy_slide),
        ("14: Define/Ideate/Prototype", build_define_ideate_prototype),
        ("15: Testing", build_testing_slide),
        ("16: Quiz 3", build_quiz3_slide),
        ("17: Flexibility vs Adaptability", build_flexibility_slide),
        ("18: Business Opportunities", build_opportunities_slide),
        ("19: Market Fit", build_market_fit_slide),
        ("20: Prototyping & A/B", build_prototyping_slide),
        ("21: Quiz 4", build_quiz4_slide),
        ("22: Entrepreneurial Mindset", build_entrepreneurial_mindset),
        ("23: Tools & Resources", build_tools_slide),
        ("24: Challenges", build_challenges_slide),
        ("25: Quiz 5", build_quiz5_slide),
        ("26: Summary", build_summary_slide),
        ("27: Closing", build_closing_slide),
    ]

    for name, builder in builders:
        print(f"  Building slide {name}...")
        builder(prs)

    output_path = "/Users/qusaiabushanap/dev/storyboard/output/NJR01/U02/NJR01_U02_Interactive_Lecture.pptx"
    prs.save(output_path)
    print(f"\nSUCCESS! Saved {len(prs.slides)} slides to:")
    print(f"  {output_path}")

if __name__ == "__main__":
    main()
